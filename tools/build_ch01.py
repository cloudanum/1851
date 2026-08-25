#!/usr/bin/env python3
"""build_ch01.py — assemble the A.2 'Introduction to AI in Software Development' chapter.

Carrier = _bak0/1851-Ch01.pptx (restored fresh every run -> idempotent).
Keeps the definitions / AI-vs-rules / Narrow-General-Generative / Turing-Test /
SDLC / security core (25 slides), cuts the stale 2025 market survey, MLOps,
hardware/frameworks, dev-readiness, regulation, and old case-study blocks (23
slides), authors 13 new GenAI-first slides, rewrites objectives / agenda /
summary in place, adds Chapter-8 pointer lines to the security/risks pair, then
arranges everything into the target order. Activity 1.1 is kept untouched.

Run:  /Users/iahmad/Creator/Courses_and_conferences/LT/.venv-courseware/bin/python tools/build_ch01.py
"""
from __future__ import annotations
import copy
import shutil
import sys
import zipfile
from collections import Counter
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import slidekit as sk  # noqa: F401  (import activates the partname-collision fix)
import pptx_tools as pt
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
BAK = ROOT / "decks" / "_bak0"
OUT = ROOT / "decks" / "1851-Ch01.pptx"

# --- carrier slides kept (0-based, against _bak0/1851-Ch01.pptx) --------------
# 0-2 openers; 3-6 best four of definitions/AI-vs-rules (s8 venn dropped);
# 8-10 Narrow/General/Generative; 11-15 Turing block + Activity 1.1 (untouched);
# 23-24 ML lifecycle best two (rest is Ch03); 40-43 SDLC block;
# 44-45 security/risks; 46-47 summary + recap.
KEEP = [0, 1, 2, 3, 4, 5, 6, 8, 9, 10, 11, 12, 13, 14, 15,
        23, 24, 40, 41, 42, 43, 44, 45, 46, 47]

# --- rewritten text for kept carrier slides ----------------------------------
OBJECTIVES = [
    "Explain how predictive, generative, and agentic AI differ—and where each fits in software work",
    "Differentiate AI systems from traditional rule-based systems",
    "Classify AI by capability: Narrow, General, and Generative",
    "Map AI assistance across the SDLC—from requirements to deployment",
    "Build a working mental model of tokens, context windows, and cost",
    "Apply trust-but-verify practices: review, verification, and accountability for AI output",
]
AGENDA = [
    "Defining AI—rules vs. learned behavior",
    "Narrow, General, and Generative AI",
    "AI for software engineers in 2026—assistants, agents, tokens, and cost",
    "The Turing Test (Activity 1.1)",
    "The GenAI-augmented SDLC",
    "Security, risks, and team readiness",
    "Case study, course map, and Activity 1.2",
]
SUMMARY = [
    "AI has shifted from hand-written rules to learned models—and now to systems that generate and act",
    "Predictive and generative AI are commodity tooling in 2026; agents are the emerging wave",
    "Tokens, context windows, and cost shape every GenAI design decision",
    "AI assists every SDLC phase, but engineer accountability, review, and verification do not change",
    "Classic ML lifecycle skills transfer directly to LLM evals later in this course",
    "Security, privacy, and governance are first-class concerns—Chapters 2 and 8 go deep",
]
SEC_POINTER = ("GenAI-specific threats—prompt injection, data leakage, model supply chain—"
               "are covered in Chapter 8")
RISK_POINTER = ("GenAI-specific risks—hallucinated output, data leakage, IP and licensing—"
                "are covered in Chapter 8")

# --- 13 new slides -------------------------------------------------------------
# kind 'b' = bullets on 'Content with Header. Full Page'; 't' = table;
# kind 'a' = activity on 'Do Now with Typing Hands'. Every slide gets notes.
NEW_SLIDES = [
    dict(key="history", kind="b", title="From Rules to Learning to Generation",
         bullets=[
             "1950s–1980s: hand-written rules and expert systems—powerful but brittle",
             "1990s–2000s: machine learning replaces rules with models trained on data",
             "2012: deep learning wins on perception—vision and speech jump in quality",
             "2017–2022: transformers and large language models master language—and code",
             "2023–2026: coding assistants become default tooling; agents begin executing multi-step work",
             "Each wave absorbed the last: rules, ML, and GenAI now coexist in the same systems",
         ],
         notes="A 70-year arc in one slide. The pattern that matters: each wave did not erase the "
               "previous one—rules still guard safety checks, classic ML still scores transactions, "
               "and GenAI drafts and reasons on top. In 2026 the engineer's job is picking the right "
               "wave for each problem; the 'three waves' slide later in this chapter makes that concrete."),
    dict(key="a", kind="b", title="AI for Software Engineers in 2026",
         bullets=[
             "Predictive AI: classifies and forecasts—fraud scores, failure prediction, prioritization (commodity)",
             "Generative AI: drafts text, code, tests, docs, and designs from natural-language prompts (commodity)",
             "Agentic AI: plans and executes multi-step work with tools—runs tests, opens PRs, triages tickets (emerging)",
             "GPT-5-class models and mature coding assistants are standard issue in most engineering orgs",
             "Commodity means: the differentiator is no longer access to AI, but how well you direct and verify it",
             "This course treats AI as engineering material—capabilities, limits, and controls",
         ],
         notes="Set the year's baseline: predictive and generative AI are commodities—built into every "
               "IDE and CI pipeline—while agents are crossing into production. The audience has already "
               "used coding assistants, so acknowledge that and move fast. The key reframe: advantage "
               "moved from 'having AI' to directing and verifying it well."),
    dict(key="b", kind="b", title="The Three Waves in One Team",
         bullets=[
             "Classic ML: predictions on structured data—scoring, forecasting, classification (Chapter 3)",
             "GenAI assistants: a pair-programmer and drafter inside your tools (Chapters 4–5)",
             "Agents: goal-driven automation that uses tools and runs workflows (Chapter 7)",
             "Same organization, different jobs: fraud scoring, drafting a service, triaging incidents",
             "Choosing the wave is an architecture decision—cost, latency, risk, and data boundaries differ",
         ],
         notes="One team can run all three waves at once: an ML model scores risk, an assistant helps "
               "developers ship, an agent triages overnight tickets. They differ in autonomy and risk—ML "
               "answers, GenAI drafts, agents act. Point out the chapter references; this course covers "
               "each wave in turn."),
    dict(key="myths", kind="b", title="GenAI Myths vs. Reality for Engineers",
         bullets=[
             "Myth: “AI will replace software engineers” — Reality: it replaces tasks; accountability shifts up-stack",
             "Myth: “The model understands the code” — Reality: it predicts plausible text; it can be confidently wrong",
             "Myth: “Bigger model is always better” — Reality: latency, cost, and data boundaries often favor smaller tools",
             "Myth: “Generated code is production-ready” — Reality: review and tests still decide what ships",
             "Myth: “Prompting is a passing fad” — Reality: directing AI is becoming core engineering literacy",
         ],
         notes="Calibrate expectations early—every class arrives with strong opinions from headlines. "
               "The through-line: fluency beats hype. Each myth maps to a skill this course teaches—"
               "verification (Ch05), model limits (Ch04), evals (Ch06)—so keep it light and moving."),
    dict(key="d", kind="b", title="Where AI Coding Help Shines — and Where It Fails",
         bullets=[
             "Shines: boilerplate, scaffolding, migrations, and common patterns",
             "Shines: unit tests, documentation, commit messages, and code explanation",
             "Shines: unfamiliar languages and APIs—a fast first draft to react to",
             "Fails: novel algorithms and domain-specific business rules",
             "Fails: subtle logic—concurrency, boundary conditions, security-sensitive paths",
             "Fails silently: output is plausible but wrong—always verify against tests and intent",
         ],
         notes="Give permission to trust assistants for the boring 80%—that is where the measured gains "
               "are. Then be blunt about the failure mode: errors are plausible and silent, not loud. "
               "Novel algorithms, concurrency, and security paths need human-led design with AI as "
               "drafter. This slide motivates the verification habits taught in Chapter 5."),
    dict(key="landscape", kind="b", title="The Coding Assistant Landscape in 2026",
         bullets=[
             "Inline autocomplete: next-line and whole-function suggestions in the editor",
             "Chat pair-programmers: explain, refactor, and draft from natural language",
             "Agentic CLI/IDE tools: read the repo, edit files, run tests, iterate autonomously",
             "CI-integrated review: automated PR feedback, test generation, and security scanning",
             "Selection is a policy decision: data boundaries, licensing, and approved-tool lists come first",
         ],
         notes="Categories, not vendors—the market moves too fast for a slide of logos to survive a "
               "quarter. What matters is the shape of the interaction: suggest, converse, or act "
               "autonomously. Stress the last bullet: which tools may see your code is an organizational "
               "policy decision, not a personal preference."),
    dict(key="g", kind="b", title="Tokens, Context, and Cost in 60 Seconds",
         bullets=[
             "Models read and write tokens: ~4 characters each; a page of text is roughly 500–700 tokens",
             "The context window is the model's working memory—everything it can see at once",
             "Prompt + history + retrieved docs + output all compete for that window",
             "Cost scales with tokens in and out—long chats and big files are not free",
             "Larger windows do not mean better answers: relevant context beats more context",
             "Chapter 4 turns this model into hands-on prompting technique",
         ],
         notes="The minimum viable mental model before the hands-on work in Chapter 4. Two intuitions "
               "to land: the window is finite working memory, so curate what goes in; and the meter runs "
               "on tokens in both directions. If learners remember 'relevant beats more' and 'output "
               "tokens cost too,' they are ready."),
    dict(key="c", kind="t", title="The GenAI-Augmented SDLC",
         header=["SDLC Phase", "AI Assist in 2026", "Human Stays Accountable For"],
         rows=[
             ["Requirements", "Drafting stories, acceptance criteria, gap analysis of specs",
              "Business value, priorities, sign-off"],
             ["Design", "Architecture options, trade-off summaries, diagram drafts",
              "Boundaries, constraints, final decisions"],
             ["Coding", "Scaffolding, implementations, refactors, explanations",
              "Intent, code review, what merges"],
             ["Testing", "Test-case generation, edge-case ideas, TDD red-green loops (Ch05)",
              "Coverage strategy, meaningful assertions"],
             ["Review", "PR summaries, style and security pre-checks",
              "Judgment, architecture fit, approval"],
             ["Deploy", "Runbooks, release notes, anomaly detection in ops",
              "Rollback decisions, production risk"],
         ],
         notes="Read a row at a time: every phase now has a credible AI assist, and every phase keeps a "
               "human-owned core. The right-hand column is the course's trust-but-verify theme in table "
               "form. Chapter 5 works through each phase hands-on, including AI-assisted TDD."),
    dict(key="f", kind="b", title="The Engineer’s Responsibilities Don’t Change",
         bullets=[
             "You are accountable for what merges—no matter who or what drafted it",
             "Review AI output like a talented intern’s: fast, confident, and sometimes wrong",
             "Trust but verify: run the tests, check the edge cases, read the diff",
             "Own the intent: AI generates solutions; engineers define the problem and the trade-offs",
             "Escalate surprises: security, licensing, and data-privacy questions go to the org, not the chatbot",
         ],
         notes="This is the culture slide. The intern analogy works because it captures both sides: real "
               "productivity and real supervision. Verification is not bureaucracy—it is the job. When in "
               "doubt about data or licensing, the answer is the organization's policy channel, never the "
               "tool itself."),
    dict(key="h", kind="b", title="AI Readiness for a Software Team",
         bullets=[
             "Data access: can AI tools reach the code, docs, and tickets they need—safely?",
             "Tooling policy: approved tools, data boundaries, and licensing rules everyone knows",
             "Skills: prompting, review discipline, and knowing when not to use AI",
             "Evals culture: measure AI output quality like you measure code quality (Chapter 6)",
             "Workflow integration: AI steps inside CI/CD and code review, not beside them",
             "Start small: one low-risk workflow, measured, then expand",
         ],
         notes="Readiness is mostly organizational, not technical. Data access and policy are the usual "
               "blockers—models are the easy part. The evals point previews Chapter 6: teams that measure "
               "AI output adopt faster and safer. Recommend pilot-and-measure over a big-bang rollout."),
    dict(key="i", kind="b", title="Case Study: One Sprint with GenAI",
         bullets=[
             "Ticket: add export-to-CSV to a reporting service—estimated at 3 days",
             "Requirements: assistant drafts acceptance criteria; engineer trims scope—30 minutes saved",
             "Coding: scaffolding and tests generated; engineer designs the streaming logic—day 1 done",
             "Review: PR bot flags a license issue; reviewer catches a date-format bug AI introduced—day 2",
             "Deploy: release notes drafted, canary watched by anomaly detection—shipped on day 2.5",
             "Result: ~40% faster, and both defects were caught by human review and policy checks",
         ],
         notes="A realistic composite, not a miracle story: the win is real (~40% on a routine ticket) "
               "and so are the two catches—an AI-introduced date bug and a licensing flag that only "
               "humans and policy gates caught. Ask the class where their own review would have caught "
               "these. The honest lesson: speed comes from drafting, safety comes from review."),
    dict(key="e", kind="b", title="How This Course Is Organized",
         bullets=[
             "Chapter 2: ethics, quality, and governance—recurring themes throughout",
             "Chapter 3: machine learning essentials—the classic wave, condensed",
             "Chapter 4: prompting and the OpenAI API—hands-on GenAI",
             "Chapter 5: GenAI across the SDLC, including AI-assisted TDD",
             "Chapter 6: evals—measuring whether AI output is good enough",
             "Chapters 7–8: agents that act, then security and guardrails for all of it",
         ],
         notes="Orient the week: we climb the autonomy ladder—predict, generate, act—with quality and "
               "governance running alongside every chapter. Chapter 4 is the hands-on core; Chapter 6's "
               "eval mindset is what separates demos from production. Point back to the three-waves "
               "slide if learners ask why the order."),
    dict(key="j", kind="a", title="Activity 1.2: Map Your Own SDLC",
         bullets=[
             "Work in pairs: list the phases of YOUR team’s SDLC as it really runs today",
             "Mark each phase: where could AI help TODAY? Where could it help in a YEAR?",
             "Note one blocker per phase: data access, policy, skills, or trust",
             "10 minutes, then debrief: one “today” and one “in a year” from each pair",
             "Capture your map—you will revisit it after Chapter 5",
         ],
         notes="Make it concrete: this is their pipeline, not the textbook's. The 'today vs. in a year' "
               "split separates commodity assists from emerging agentic ones. During debrief, group the "
               "blockers—data access and policy usually dominate, which reinforces the readiness slide. "
               "Tell them to keep the map; Chapter 5 revisits it phase by phase."),
]


# --- formatting-preserving text rewrite helpers (same pattern as build_ch03) --
def _set_para_text(p_el, text):
    """Replace a paragraph's text, keeping the first run's formatting."""
    runs = p_el.findall(qn("a:r"))
    if not runs:
        r = p_el.makeelement(qn("a:r"), {})
        r.append(p_el.makeelement(qn("a:t"), {}))
        p_el.append(r)
        runs = [r]
    first = runs[0]
    for extra in runs[1:]:
        p_el.remove(extra)
    for br in p_el.findall(qn("a:br")):
        p_el.remove(br)
    t = first.find(qn("a:t"))
    if t is None:
        t = first.makeelement(qn("a:t"), {})
        first.append(t)
    t.text = text


def rewrite_bullets(shape, lines):
    """Rewrite a body placeholder's bullets, reusing paragraph formatting."""
    txBody = shape.text_frame._txBody
    paras = txBody.findall(qn("a:p"))
    for i, line in enumerate(lines):
        if i < len(paras):
            _set_para_text(paras[i], line)
        else:
            newp = copy.deepcopy(paras[-1])
            txBody.append(newp)
            paras.append(newp)
            _set_para_text(newp, line)
    for p in paras[len(lines):]:
        txBody.remove(p)


def append_pointer(slide, text):
    """Append a top-level pointer bullet, preserving the existing bullets."""
    body = pt._body_placeholder(slide)
    lines = [p.text for p in body.text_frame.paragraphs if p.text.strip()]
    rewrite_bullets(body, lines + [text])
    body.text_frame.paragraphs[-1].level = 0


def layout_named(prs, name):
    for lay in prs.slide_layouts:
        if (lay.name or "") == name:
            return lay
    raise KeyError(f"layout {name!r} not found")


def add_new_slide(prs, spec, lay_content, lay_activity):
    if spec["kind"] == "t":
        slide = sk.add_table_slide(prs, spec["title"], spec["header"], spec["rows"],
                                   notes=spec["notes"], layout=lay_content)
        for sh in slide.shapes:                      # compact table text
            if sh.has_table:
                for ri, row in enumerate(sh.table.rows):
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            for r in p.runs:
                                r.font.size = Pt(13)
                                if ri == 0:
                                    r.font.bold = True
        return slide
    lay = lay_activity if spec["kind"] == "a" else lay_content
    slide = prs.slides.add_slide(lay)
    sk.set_title(slide, spec["title"])
    body = pt._body_placeholder(slide)               # largest text area (idx 1)
    tf = body.text_frame
    tf.word_wrap = True
    for k, b in enumerate(spec["bullets"]):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.text = b
    sk.set_notes(slide, spec["notes"])
    return slide


def main():
    # 0. fresh carrier -> idempotent
    shutil.copy(BAK / "1851-Ch01.pptx", OUT)
    prs = sk.open_prs(OUT)
    assert len(prs.slides) == 48, "carrier deck should have 48 slides"

    lay_content = layout_named(prs, "Content with Header. Full Page")
    lay_activity = layout_named(prs, "Do Now with Typing Hands")

    # 1. author the new slides (appended at end), record indices by key
    new_idx = {}
    for spec in NEW_SLIDES:
        add_new_slide(prs, spec, lay_content, lay_activity)
        new_idx[spec["key"]] = len(prs.slides) - 1

    # 2. edit kept carrier slides in place (original 0-based indices)
    rewrite_bullets(pt._body_placeholder(prs.slides[1]), OBJECTIVES)    # Objectives
    rewrite_bullets(pt._body_placeholder(prs.slides[2]), AGENDA)        # Agenda
    append_pointer(prs.slides[44], SEC_POINTER)                         # IT Security and AI
    append_pointer(prs.slides[45], RISK_POINTER)                        # Challenges and Risks
    rewrite_bullets(pt._body_placeholder(prs.slides[46]), SUMMARY)      # Summary
    rewrite_bullets(pt._body_placeholder(prs.slides[47]), OBJECTIVES)   # recap

    # 3. arrange into the target arc (drops everything not referenced)
    order = (
        [0, 1, 2]                                   # title, objectives, agenda
        + [3, 4, 5, 6]                              # definitions & AI-vs-rules (best 4)
        + [new_idx["history"]]                      # rules -> learning -> generation
        + [8, 9, 10]                                # Narrow / General / Generative
        + [new_idx[k] for k in ("a", "b", "myths", "d", "landscape", "g")]
        + [11, 12, 13, 14, 15]                      # Turing Test block + Activity 1.1
        + [23, 24]                                  # ML lifecycle + data-driven principles
        + [new_idx["c"], new_idx["f"]]              # SDLC table + responsibilities
        + [40, 41, 42, 43]                          # SDLC integration block
        + [44, 45]                                  # security & risks (Ch08 pointers)
        + [new_idx["h"], new_idx["i"]]              # readiness + case study
        + [new_idx["e"], new_idx["j"]]              # course map + Activity 1.2
        + [46, 47]                                  # summary + recap
    )
    pt.arrange(prs, order)
    pt.save(prs, OUT)

    # 4. verify: reopen, titles, blanks, duplicate partnames, zip integrity
    prs2 = Presentation(OUT)
    n = len(prs2.slides)
    print(f"\n{OUT.name}: {n} slides")
    print(f"kept carrier: {len(KEEP)} | cut: {48 - len(KEEP)} | authored new: {len(NEW_SLIDES)}")
    empty = []
    for i, s in enumerate(prs2.slides, 1):
        t = sk.slide_title(s)
        if not t.strip():
            empty.append(i)
        print(f"{i:>3} | {t}")
    assert not empty, f"empty titles at {empty}"
    assert 38 <= n <= 44, f"slide count {n} outside 38-44"
    with zipfile.ZipFile(OUT) as z:
        dupes = [p for p, c in Counter(z.namelist()).items() if c > 1]
        assert not dupes, f"duplicate zip partnames: {dupes}"
        assert z.testzip() is None, "zip integrity check failed"
    print("\nOK: no empty titles, no duplicate partnames, zip integrity verified")


if __name__ == "__main__":
    main()
