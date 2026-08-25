#!/usr/bin/env python3
"""build_engagement_ch04.py — ENGAGEMENT pass on 1851-Ch04 (A.2, 48 -> 51 slides).

Restores the A.2 deck fresh from decks/_bak1 every run (idempotent), then:
  1. inserts 'Do Now: Name That Temperature' right after 'Temperature and
     Sampling' (layout 'Do Now Group Work'),
  2. inserts 'Do Now: Fix This Prompt' right after the prompting best-practices
     block's 'Iterative Prompting Techniques' (layout 'Do Now with Typing Hands'),
  3. inserts 'Activity 4.2: Prompt Pattern Clinic' right after the existing
     'Lab 4.1: First OpenAI API Calls' exercise slide (layout 'Exercise
     Reference Slide with Typing Hands'),
  4. adds Activity 4.2 to the Agenda bullet that mentions Lab 4.1.

Run:  /Users/iahmad/Creator/Courses_and_conferences/LT/.venv-courseware/bin/python tools/build_engagement_ch04.py
"""
from __future__ import annotations
import shutil
import sys
import zipfile
from collections import Counter
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import slidekit as sk  # noqa: F401  (import activates the partname-collision fix)
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent
BAK1 = ROOT / "decks" / "_bak1"
OUT = ROOT / "decks" / "1851-Ch04.pptx"

# --- new slide content -------------------------------------------------------
DONOW_TEMP = ("Do Now: Name That Temperature", "Do Now Group Work",
    ["Same prompt, three runs — only the temperature changed. Which is which?",
     "Prompt: \"Finish this commit message: 'Fix null pointer in login by…'\"",
     "Output A: \"…adding a null check before reading the session object.\"",
     "Output B: \"…validating credentials before the session lookup, so expired tokens no longer crash it.\"",
     "Output C: \"…teaching the login gremlin to waltz past the null abyss. Confetti included.\"",
     "Match A, B, C to temperature 0, 0.7, and 1.3. Time box: 5 minutes, groups of three."],
    "Answers: A = temperature 0 (safe, deterministic, the obvious completion); B = 0.7 (coherent, but "
    "invents plausible detail beyond the prompt); C = 1.3 (creative to the point of being unusable in a "
    "commit message). Debrief: low temperature for exact tasks — parsing, schema output, tests; higher "
    "only when divergence is the goal, like brainstorming names. Ask what they would pick for generating "
    "unit tests and why (0 — you want boring).")

DONOW_FIX = ("Do Now: Fix This Prompt", "Do Now with Typing Hands",
    ["The broken prompt: \"Write some code to parse dates.\"",
     "In pairs, rewrite it so it specifies:",
     "A role for the model — who is it?",
     "The language, the inputs, and the exact output format",
     "Constraints: edge cases, allowed libraries, what NOT to do",
     "Time box: 7 minutes — one rewritten prompt per pair, then we compare."],
    "Give pairs 7 minutes to rewrite the prompt on the slide. Model answer: 'You are a senior Python "
    "engineer. Write a function parse_date(text: str) -> datetime.date that accepts ISO 8601 and common "
    "US formats (MM/DD/YYYY), raises ValueError on anything else, uses only the standard library, and "
    "includes three unit tests.' Debrief: have two pairs read theirs aloud; point out how role, format, "
    "and constraints each removed one ambiguity that made the original unanswerable — which language, "
    "which formats, what happens on bad input.")

ACTIVITY_42 = ("Activity 4.2: Prompt Pattern Clinic",
    ["Follow the detailed instructions in the Activity 4.2 notebook on your VM",
     "Practice six prompt patterns from this chapter on real code artifacts",
     "Score each pattern's output with the rubric in the notebook",
     "Swap one prompt with another pair and peer-review it against the rubric"],
    "In-class paired activity, about 40 minutes: 30 for the notebook steps, 10 for the swap and debrief. "
    "Pairs apply six patterns — role, few-shot, constraints, structured output, grounding, and "
    "chain-of-thought — to real code artifacts, score each output with the notebook's rubric, then trade "
    "one prompt with a neighboring pair for peer review. Debrief: ask which pattern moved the rubric "
    "score the most; it is usually constraints and grounding. The notebook reuses the Lab 4.1 API setup, "
    "so no new credentials are needed.")


# --- helpers ------------------------------------------------------------------
def pick_layout_exact(prs, name):
    """Exact layout-name match (pick_layout's substring match would return the
    'Numbered …' variant first)."""
    for master in prs.slide_masters:
        for lay in master.slide_layouts:
            if (lay.name or "") == name:
                return lay
    raise AssertionError(f"layout not found: {name!r}")


def add_donow_slide(prs, title, layout_name, bullets, notes):
    """Add a Do Now / Discussion slide on the named template layout (graphic on
    the left, text column on the right). Returns the new slide."""
    lay = pick_layout_exact(prs, layout_name)
    slide = prs.slides.add_slide(lay)
    sk.set_title(slide, title)
    body = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1 and ph.has_text_frame:
            body = ph
            break
    assert body is not None, f"content placeholder not found on {layout_name!r}"
    # the layout's content placeholder is stubby; give it room for 5+ lines
    body.height = Inches(4.5)
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
    sk.set_notes(slide, notes)
    return slide


def add_exercise_slide(prs, title, bullets, notes):
    """Same exercise-layout builder as build_ch04.add_exercise_slide."""
    lay = pick_layout_exact(prs, "Exercise Reference Slide with Typing Hands")
    slide = prs.slides.add_slide(lay)
    sk.set_title(slide, title)
    body = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1 and ph.has_text_frame:
            body = ph
            break
    assert body is not None, "content placeholder not found on exercise layout"
    body.height = Inches(4.5)
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
    sk.set_notes(slide, notes)
    return slide


def find_slide(prs, title_sub):
    """0-based index of the unique slide whose title contains title_sub."""
    hits = [i for i, s in enumerate(prs.slides) if title_sub in sk.slide_title(s)]
    assert len(hits) == 1, f"anchor {title_sub!r}: expected 1 hit, got {hits}"
    return hits[0]


def _set_para_text(p_el, text):
    runs = p_el.findall(qn("a:r"))
    if not runs:
        r = p_el.makeelement(qn("a:r"), {})
        r.append(p_el.makeelement(qn("a:t"), {}))
        p_el.append(r)
        runs = [r]
    first = runs[0]
    for extra in runs[1:]:
        p_el.remove(extra)
    t = first.find(qn("a:t"))
    if t is None:
        t = first.makeelement(qn("a:t"), {})
        first.append(t)
    t.text = text


def replace_in_slide(slide, old, new):
    """Run-level first (preserves formatting); paragraph-level fallback."""
    n = 0
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if old in r.text:
                    r.text = r.text.replace(old, new)
                    n += 1
    if n:
        return n
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p_el in sh.text_frame._txBody.findall(qn("a:p")):
            full = "".join(t.text or "" for t in p_el.iter(qn("a:t")))
            if old in full:
                _set_para_text(p_el, full.replace(old, new))
                n += 1
    return n


def main():
    # 0. fresh restore from _bak1 -> idempotent
    shutil.copy(BAK1 / OUT.name, OUT)
    prs = sk.open_prs(OUT)
    assert len(prs.slides) == 48, f"expected 48-slide A.2 deck, got {len(prs.slides)}"

    # 1. agenda: add Activity 4.2 alongside the Lab 4.1 mention (slide 3)
    n = replace_in_slide(prs.slides[2], "Lab 4.1, and what's next",
                         "Lab 4.1 and Activity 4.2, and what's next")
    assert n > 0, "agenda bullet for Lab 4.1 not found"

    # 2+3. append each new slide and immediately move it into place — descending
    # target index, so each move only shifts already-placed slides and every
    # anchor lookup runs against the current order (no stale indices).
    add_exercise_slide(prs, ACTIVITY_42[0], ACTIVITY_42[1], ACTIVITY_42[2])
    at_act = find_slide(prs, "Lab 4.1: First OpenAI API Calls") + 1   # right after Lab 4.1
    sk.move_slide(prs, len(prs.slides) - 1, at_act)

    add_donow_slide(prs, DONOW_FIX[0], DONOW_FIX[1], DONOW_FIX[2], DONOW_FIX[3])
    at_fix = find_slide(prs, "Iterative Prompting Techniques") + 1    # after pitfalls/iteration
    sk.move_slide(prs, len(prs.slides) - 1, at_fix)

    add_donow_slide(prs, DONOW_TEMP[0], DONOW_TEMP[1], DONOW_TEMP[2], DONOW_TEMP[3])
    at_temp = find_slide(prs, "Temperature and Sampling") + 1         # right after temperature slide
    sk.move_slide(prs, len(prs.slides) - 1, at_temp)

    print(f"inserted at 1-based: {at_temp+1} (Name That Temperature), "
          f"{at_fix+1} (Fix This Prompt), {at_act+1} (Activity 4.2)")

    sk.save(prs, OUT)

    # 4. verify: reopen, count, titles, blanks, duplicate partnames, zip integrity
    prs2 = Presentation(OUT)
    n = len(prs2.slides)
    print(f"\n{OUT.name}: {n} slides")
    empty = []
    for i, s in enumerate(prs2.slides, 1):
        t = sk.slide_title(s)
        if not t.strip():
            empty.append(i)
        print(f"{i:>3} | {t}")
    assert not empty, f"empty titles at {empty}"
    assert n == 51, f"slide count {n} != 51"
    titles = [sk.slide_title(s) for s in prs2.slides]
    # positions are relative to neighbors, so verify the adjacency contract
    # directly (insertion order makes captured indices stale)
    assert titles.index(DONOW_TEMP[0]) == titles.index("Temperature and Sampling") + 1
    assert titles.index(DONOW_FIX[0]) == titles.index("Iterative Prompting Techniques") + 1
    assert titles.index(ACTIVITY_42[0]) == titles.index("Lab 4.1: First OpenAI API Calls") + 1
    agenda_txt = "\n".join(sh.text_frame.text for sh in prs2.slides[2].shapes if sh.has_text_frame)
    assert "Lab 4.1 and Activity 4.2" in agenda_txt, "agenda update missing"
    with zipfile.ZipFile(OUT) as z:
        dupes = [p for p, c in Counter(z.namelist()).items() if c > 1]
        assert not dupes, f"duplicate zip partnames: {dupes}"
        assert z.testzip() is None, "zip integrity check failed"
    print("\nOK: 51 slides, no empty titles, no duplicate partnames, zip integrity verified")


if __name__ == "__main__":
    main()
