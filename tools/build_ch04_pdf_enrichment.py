#!/usr/bin/env python3
"""Enrich 1851 Ch04 with 'The Art of Agent Prompting' (book ch.3) material.

- Rewrites the two Do Now slides (s12, s33) around two-layer prompts and PTCF
- Updates Agenda (s3), Agents-Preview bridge (s45), Lab 4.1 / Activity 4.2
  slides (s48, s49), Summary (s50)
- Inserts 5 new slides with PDF-derived content (3 with SVG diagrams)

Content source: 30_Agents_Ch_03_The Art of Agent Prompting.pdf
Idempotent: every edit is guarded by a text/identity check.

Usage: build_ch04_pdf_enrichment.py <1851-Ch04.pptx>
"""
import sys
from pathlib import Path

from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).parent))
import slidekit as sk  # noqa: E402

DIAGRAMS = Path(__file__).parent / 'diagrams'

# ------------------------------------------------------------- new slides
# (title, at_index into the ORIGINAL 51-slide deck, bullets, levels, notes,
#  diagram png or None). Applied in given order = descending at_index.
NEW_SLIDES = [
    ('From Instructions to Constitutions', 44, [
        'A chatbot prompt is a one-off command; an agent prompt is a constitution',
        'It persists across turns, tools, and handoffs — shaping every decision',
        'Written and edited in natural language — no code redeploy needed',
        'Carries identity (persona), mission (task), law (context), voice (format)',
        'Multi-agent handoffs must re-establish the constitution in each message',
        'Chapter 7 builds agents on exactly this foundation',
    ], None,
     'Prompt engineering becomes cognitive programming: the prompt is the '
     'constitutional bedrock of agent behavior, not a transient instruction.',
     None),
    ('PTCF Walk-Through: An Enterprise Billing Agent', 34, [
        'Persona: empathetic senior support specialist, 5 yrs enterprise SaaS',
        'Professional, approachable, solution-oriented',
        'Task: resolve billing inquiries — diagnose, guide, escalate within 24 h',
        'Context: Fortune 500 SLA; never request passwords; privacy regulations',
        'Format: acknowledge \u2192 diagnose \u2192 resolve, numbered, with case references',
        'Each component is authored, audited, and iterated independently',
        'Anti-pattern check: does any component fight the others?',
    ], [0, 1, 0, 0, 0, 0, 0],
     'Walk-through of the book\'s PTCF-enhanced system prompt for an '
     'enterprise customer support agent.',
     None),
    ('The PTCF Blueprint', 34, [
        'Persona — who the agent is (never just "helpful assistant")',
        'Task — the core mission, with explicit "must not" boundaries',
        'Context — operational law: SLAs, regulations, conflict rules',
        'Format — the output contract: structure, limits, fallback',
        'Ambiguity is the enemy of alignment — PTCF decomposes the system prompt',
    ], None,
     'PTCF = Persona, Task, Context, Format. Adopted in LangChain, CrewAI, '
     'AutoGen communities; CRISPE is a more granular alternative.',
     'ptcf-blueprint.png'),
    ('Designing Thinking Agents: The Prompting Spectrum', 27, [
        'L1 reactive agents: direct, unambiguous commands',
        'L2 tool-using agents: prompt what, how, and when to use tools',
        'L3 planning agents: "think step by step" — structured decomposition',
        'L4 learning agents: metacognitive prompts — reason about reasoning',
        'Rule: prompt sophistication must scale with agent capability',
    ], None,
     'Agent capability spectrum (book fig 3.2). LangChain targets L2-L3; '
     'CrewAI orchestration reaches L3-L4.',
     'thinking-spectrum.png'),
    ('The Two-Layer Prompt Architecture', 21, [
        'System prompt = the constitution: how the agent behaves',
        'Identity, boundaries, style — loaded once, persistent',
        'User prompt = the stimulus: what the agent should do',
        'This turn\'s task, data, question — changes every turn',
        'Decouples personality from task: same agent, any request',
        'Design constraint: system tokens are spent on every API call',
    ], [0, 1, 0, 1, 0, 0],
     'The diplomat analogy: system prompt is national policy, user prompt is '
     'the negotiation at hand.',
     'two-layer-architecture.png'),
]

FULL_BODY = (0.29, 1.04, 9.51, 5.60)
LEFT_BODY = (0.29, 1.04, 4.60, 5.60)
PIC_POS = (5.05, 1.35, 4.70)          # left, top, width


def body_shape(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            return ph
    return None


def set_geom(shape, geom):
    l, t, w, h = geom
    shape.left, shape.top, shape.width, shape.height = (Inches(l), Inches(t),
                                                        Inches(w), Inches(h))


def rebody(slide, marker, lines):
    """Replace the body text of the shape containing `marker`. Returns True
    when a replacement happened (False = already done / marker absent)."""
    for sh in slide.shapes:
        if sh.has_text_frame and marker in sh.text_frame.text:
            tf = sh.text_frame
            tf.clear()
            for i, b in enumerate(lines):
                (tf.paragraphs[0] if i == 0 else tf.add_paragraph()).text = b
            return True
    return False


def find_slide(prs, title_start):
    for i, s in enumerate(prs.slides):
        if sk.slide_title(s).startswith(title_start):
            return i, s
    return None, None


def insert_after(slide, anchor_text, new_line):
    """Insert new_line as a paragraph right after the paragraph containing
    anchor_text, inside the same text frame. Idempotent."""
    for sh in slide.shapes:
        if not sh.has_text_frame or anchor_text not in sh.text_frame.text:
            continue
        tf = sh.text_frame
        if new_line in tf.text:
            return False
        anchor = next((p for p in tf.paragraphs if anchor_text in p.text), None)
        if anchor is None:
            return False
        p = tf.add_paragraph()
        p.text = new_line
        p.level = anchor.level
        anchor._p.addnext(p._p)
        return True
    return False


def main(path):
    prs = sk.open_prs(path)
    changed = 0

    # ---- 1. Do Now s12 -> two-layer classification drill
    _, s12 = find_slide(prs, 'Do Now: Name That Temperature')
    if s12 is not None:
        sk.set_title(s12, 'Do Now: Constitution or Command?')
        rebody(s12, 'Same prompt, three runs', [
            'System layer (constitution) or user layer (stimulus)?',
            '1) "You never ask for passwords."  2) "Summarize this invoice."',
            '3) "You are a methodical billing specialist."  '
            '4) "Book the cheapest flight to Toronto."',
            '5) "When instructions conflict, escalate to a human reviewer."',
            'Time box: 5 minutes, groups of three — then we vote.',
        ])
        sk.set_notes(s12, 'Answer key — SYSTEM (constitution, persists): 1, 3, 5. '
                          'USER (stimulus, per-turn task): 2, 4.')
        changed += 1

    # ---- 2. Do Now s33 -> PTCF misalignment clinic
    _, s33 = find_slide(prs, 'Do Now: Fix This Prompt')
    if s33 is not None and rebody(s33, 'parse dates', [
            'An agent\'s system prompt: [PERSONA] "You are a creative, '
            'experimental assistant."',
            '[TASK] "Help users troubleshoot enterprise billing issues."',
            '[FORMAT] "Always respond with a numbered list."',
            'Why does this agent oscillate between whimsical and procedural?',
            'In pairs: rewrite all three PTCF components so they reinforce '
            'each other. Time box: 7 minutes.',
    ]):
        sk.set_notes(s33, 'Corrected: [PERSONA] methodical enterprise billing '
                          'specialist, 5 yrs SaaS financial operations, professional and '
                          'solution-oriented. [TASK] Resolve billing inquiries: diagnose '
                          'discrepancies, explain charges, escalate within 24 h. '
                          '[FORMAT] Numbered list: (1) acknowledge (2) diagnose '
                          '(3) resolution or escalation. Persona, task, and format now '
                          'reinforce each other.')
        changed += 1

    # ---- 3. Agenda
    _, s3 = find_slide(prs, 'Agenda')
    if s3 is not None and insert_after(
            s3, 'zero-shot to reasoning',
            'Agent-grade prompting: two-layer architecture and the PTCF blueprint'):
        changed += 1

    # ---- 4. Lab 4.1 slide: retitle + rebody
    _, s48 = find_slide(prs, 'Lab 4.1: First OpenAI API Calls')
    if s48 is not None:
        sk.set_title(s48, 'Lab 4.1: Building a PTCF Agent with the OpenAI API')
        rebody(s48, 'Sweep temperature', [
            'Follow the detailed instructions in the Lab 4.1 notebook on your VM',
            'Assemble a system prompt element by element: persona, task, context, format',
            'Make Responses API calls with your PTCF constitution and sweep temperature',
            'Ablation: remove one PTCF element and watch the agent\'s behavior drift',
        ])
        changed += 1

    # ---- 5. Activity 4.2 slide: rebody (title unchanged)
    _, s49 = find_slide(prs, 'Activity 4.2: Prompt Pattern Clinic')
    if s49 is not None and rebody(s49, 'six prompt patterns', [
            'Follow the detailed instructions in the Activity 4.2 notebook on your VM',
            'Diagnose broken agent system prompts with the PTCF checklist',
            'Rewrite misaligned persona/task/context/format so they reinforce each other',
            'Peer-review one repaired constitution against the rubric in the notebook',
    ]):
        changed += 1

    # ---- 6. Bridge + Summary lines
    _, s45 = find_slide(prs, 'Agents: Preview')
    if s45 is not None and insert_after(
            s45, 'one prompt, one response',
            'An agent\'s prompt is its constitution — two layers, PTCF-designed'):
        changed += 1
    _, s50 = find_slide(prs, 'Summary')
    if s50 is not None and insert_after(
            s50, 'RAG grounds prompts',
            'Agents run on constitutions: the two-layer architecture and PTCF (Ch. 7 builds them)'):
        changed += 1

    # ---- 7. Insert the 5 new slides (descending at_index, original numbering)
    for title, at_index, bullets, levels, notes, png in NEW_SLIDES:
        if find_slide(prs, title)[1] is not None:
            continue
        slide = sk.add_bullets_slide(prs, title, bullets, levels, notes,
                                     ref_index=16, at_index=at_index)
        ph = body_shape(slide)
        if png:
            set_geom(ph, LEFT_BODY)
            slide.shapes.add_picture(str(DIAGRAMS / png),
                                     Inches(PIC_POS[0]), Inches(PIC_POS[1]),
                                     width=Inches(PIC_POS[2]))
        else:
            set_geom(ph, FULL_BODY)
        changed += 1

    sk.save(prs, path)
    print(f'{path}: {changed} edits applied, {len(prs.slides._sldIdLst)} slides')


if __name__ == '__main__':
    main(sys.argv[1])
