#!/usr/bin/env python3
"""build_ch07.py — assemble the A.2 'AI Agents and Agentic Workflows' chapter.

Carrier = _bak0/1851-Ch07.pptx (restored fresh every run -> idempotent).
The A.1 neural-network content slides are all dropped (that content moved to
Ch03); only the openers (title/objectives/agenda) and closers (summary/recap)
are kept and rewritten. The body is cloned cross-deck (images re-embedded via
pptx_tools.clone_slide) from:

  16_1  = 2016-agentic-security/2016-slides/enriched/Ch01-Pass5.pptx
  16_5  = 2016-agentic-security/2016-slides/enriched/Ch05-Final.pptx
  1258  = 1258-applied-ai.../1258a4-author-input/decks/1258-Ch07-Building-with-LLMs.pptx

Four slides are authored new (From Prompts to Agents / Agents in the SDLC /
Evaluating Agents / Lab 7.1). Cloned slides get light branding patches
(wrong chapter/lab refs, 'SOC takeaway', government phrasing) via PATCHES.

Run:  /Users/iahmad/Creator/Courses_and_conferences/LT/.venv-courseware/bin/python tools/build_ch07.py
"""
from __future__ import annotations
import copy
import shutil
import sys
import zipfile
from collections import Counter
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import slidekit as sk  # noqa: F401  (import activates the partname-collision fix)
import pptx_tools as pt
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent
BAK = ROOT / "decks" / "_bak0"
OUT = ROOT / "decks" / "1851-Ch07.pptx"
SRC = {
    "16_1": Path("/Users/iahmad/Creator/Courses_and_conferences/LT/courses/"
                 "2016-agentic-security/2016-slides/enriched/Ch01-Pass5.pptx"),
    "16_5": Path("/Users/iahmad/Creator/Courses_and_conferences/LT/courses/"
                 "2016-agentic-security/2016-slides/enriched/Ch05-Final.pptx"),
    "1258": Path("/Users/iahmad/Creator/Courses_and_conferences/LT/courses/"
                 "1258-applied-ai-for-government-it-professionals/"
                 "1258a4-author-input/decks/1258-Ch07-Building-with-LLMs.pptx"),
}

# --- clone lists (1-based, against source decks) ------------------------------
CLONES = {
    "16_1": [21, 22, 23, 24, 25,          # what is an agent
             30, 32,                      # tool use / MCP
             26, 27,                      # PEV pattern, single vs multi
             37, 38, 39, 40, 41,          # vendor landscape + OpenAI stack
             35],                         # limitations
    "16_5": [22, 23,                      # PEV in production + human gate, SOC case study
             25, 26, 27,                  # Kafka / Redis+Celery / event backbone
             28, 29,                      # LangChain-CrewAI-AutoGen + compared
             41, 42,                      # full agentic stack
             43, 44],                     # human-on-the-loop rules
    "1258": [36, 37,                      # answering->acting, LLM+tools+loop
             12, 38,                      # JSON schemas, function-calling mechanics
             20, 43, 22,                  # MCP trio
             46, 47,                      # allowed_tools, loop caps
             40, 41, 42,                  # single / workflow / multi-agent
             44, 45,                      # HITL approvals, least privilege
             49, 50],                     # failure modes, injection via tool results
}

# --- branding patches on cloned slides: (old, new) substring per paragraph ----
PATCHES = {
    ("16_1", 21): [("Security lens", "Risk lens")],
    ("16_1", 23): [("Memory poisoning and context manipulation are exploited in Ch4",
                    "Memory poisoning and context manipulation are real-world attack vectors"),
                   ("LT 1376 Ch3 (memory systems)", "LT 1376 (agent memory systems)")],
    ("16_1", 24): [("Why a security course cares — context is an attack surface:",
                    "Why engineers care — context is an attack surface:"),
                   ("Indirect injection (Ch3) smuggles hostile instructions",
                    "Indirect injection smuggles hostile instructions"),
                   ("Payload splitting (Ch4) hides an attack", "Payload splitting hides an attack")],
    ("16_1", 26): [("Why security cares: the Verifier is a control point — policy checks, "
                    "output validation, and guardrails live here. Ch5 builds SOC pipelines on "
                    "exactly this topology.",
                    "Why it matters: the Verifier is a control point — policy checks, output "
                    "validation, and guardrails live here. Production pipelines are built on "
                    "exactly this topology.")],
    ("16_1", 32): [('FastMCP("secops-tools")', 'FastMCP("agent-tools")')],
    ("16_1", 35): [("(Ch4 slopsquatting)", "(slopsquatting)"),
                   ("Ch3 weaponizes injection; Ch4 weaponizes hallucination and resource "
                    "limits (see next slide)",
                    "Injection, hallucination, and resource limits are weaponized against "
                    "unguarded agents")],
    ("16_1", 39): [("SOC takeaway:", "Takeaway:")],
    ("16_1", 40): [("SOC takeaway:", "Takeaway:")],
    ("16_1", 41): [("SOC takeaway:", "Takeaway:")],
    ("16_5", 22): [("Lab 5.1: build this with CrewAI or LangChain over a real alert stream",
                    "Lab 7.1: build this pattern with a human gate on irreversible actions")],
    ("16_5", 23): [("(maps to Lab 5.1)", "(maps to Lab 7.1)")],
    ("16_5", 25): [("Lab 5.1: publish alert events, subscribe with the triage agent, "
                    "verify delivery",
                    "In the lab: publish events, subscribe with an agent, verify delivery")],
    ("16_5", 28): [("(mind Ch4 poisoning)", "(see the memory slides)"),
                   ("Lab 5.1 uses CrewAI or LangChain to build the pipeline",
                    "Lab 7.1 uses an agent framework to build the pipeline")],
    ("16_5", 41): [("(→ Ch6 governance)", "(governance)"),
                   ("The complete defensive stack", "The complete agentic stack")],
    ("16_5", 42): [("(→ Ch6)", "(governance)")],
    ("16_5", 43): [("This is the Lab 5.1 control and the bridge to Ch6 governance",
                    "This is the Lab 7.1 control and the bridge to governance")],
    ("16_5", 44): [("(Lab 5.1)", "(Lab 7.1)")],
    ("1258", 12): [("The mechanism the agent lab (7.3) is built on",
                    "The mechanism the chapter lab is built on")],
    ("1258", 20): [("worth knowing as the standard way to expose agency tools",
                    "worth knowing as the standard way to expose tools")],
    ("1258", 36): [("You will build a real one in Lab 7.3",
                    "You will build a real one in Lab 7.1")],
    ("1258", 38): [("Lab 7.3 Stage A traces exactly one such call end to end",
                    "Lab 7.1 traces exactly one such call end to end")],
    ("1258", 40): [("Covers a large share of real government use cases",
                    "Covers a large share of real-world use cases")],
    ("1258", 41): [("For government, prefer workflows where the process is known and "
                    "must be auditable",
                    "Prefer workflows where the process is known and must be auditable")],
    ("1258", 43): [("Lab 7.3 shows an optional MCP demo after you hand-wire tools first",
                    "The lab shows an optional MCP demo after you hand-wire tools first")],
    ("1258", 44): [("Lab 7.3 Stage C gates a notification behind a typed approval",
                    "Lab 7.1 gates a consequential action behind a typed approval")],
    ("1258", 46): [("Lab 7.3 Stage C turns these controls from policy into running code",
                    "Lab 7.1 turns these controls from policy into running code")],
    ("1258", 50): [("Lab 7.3 Stage C runs this exact attack and shows the guardrails catch it",
                    "Lab 7.1 runs this exact attack and shows the guardrails catch it")],
}

RETITLES = {
    ("16_5", 22): "Planner-Executor-Verifier in Production — with a Human Gate",
    ("16_5", 41): "The Agentic Stack",
    ("16_5", 42): "The Agentic Stack — Layer by Layer",
}

# --- rewritten text for kept carrier slides -----------------------------------
OBJECTIVES = [
    "Explain what makes a system agentic: the cognitive loop, the LLM + tools + loop anatomy, and agent memory and context",
    "Implement tool use with function calling, JSON schemas, and MCP — scoped by tool allow-lists and loop caps",
    "Compare orchestration patterns — single-agent, workflow, and multi-agent — including Planner-Executor-Verifier",
    "Design agentic workflows with human-in-the-loop approvals, least-privilege tools, and event-backbone plumbing",
    "Position agents across the software development lifecycle and navigate the vendor landscape, including the OpenAI Responses/Agents stack",
    "Recognize agent failure modes and apply agent-specific evaluation techniques",
]
AGENDA = [
    "What Is an Agent? Cognitive Loops, Memory, and ReAct",
    "Tools and Protocols: Function Calling, JSON Schemas, and MCP",
    "Orchestration Patterns: Single-Agent, Workflows, and Multi-Agent",
    "Agentic Workflows and Automation: Guardrails, Plumbing, and the SDLC",
    "The Vendor Landscape and the OpenAI Agents Stack",
    "Agent Failure Modes and Evaluation",
    "Lab 7.1: Build a Guarded Agent",
]
SUMMARY = [
    "An agent = LLM + tools + loop + oversight: the model reasons, tools act, the loop iterates, and guardrails bound it",
    "Function calling and MCP standardize how agents reach tools; allow-lists, approval gates, and loop caps bound what they may do",
    "Patterns scale from single agent to workflow to multi-agent; Planner-Executor-Verifier builds a control point into the topology",
    "Agentic workflows automate SDLC toil — with human gates on irreversible actions and a tamper-evident audit trail",
    "Failure modes are real — endless loops, hallucinated calls, injection via tool output — so evals (Chapter 6) and guardrails are mandatory",
]

# --- 4 new slides --------------------------------------------------------------
NEW_FROM_PROMPTS = (
    "From Prompts to Agents",
    ["Chapter 4 gave you prompting: one call, one response — the model only answers",
     "Add tools (function calling): the model can request actions on the world",
     "Add a loop: plan → call tool → observe result → re-plan until the goal is met",
     "Add memory and guardrails: the agent becomes production-capable",
     "This chapter climbs that ladder — then builds workflows on top of agents"],
    "This is the bridge from Chapter 4 to the rest of the course. Each rung is a small, "
    "mechanical step: tools, then a loop, then memory and guardrails — no magic. By the end "
    "students should see an agent as an engineering composition of things they already know.",
)
NEW_SDLC = (
    "Agents in the Software Development Lifecycle",
    ["Requirements & design — agents draft user stories, critique architectures, surface edge cases",
     "Implementation — coding copilots and agentic IDEs generate, refactor, and navigate code",
     "Testing & review — agents write tests, review pull requests, and triage failures",
     "Operations — incident triage, runbook execution, and remediation behind human gates",
     "Chapter 5 goes deep on GenAI across the SDLC; this chapter supplies the agent mechanics"],
    "Position agents as SDLC automation, not a toy: every phase has toil an agent can absorb. "
    "Chapter 5 owns the full GenAI-across-the-SDLC treatment, so keep this as the map and "
    "forward-reference the details. Stress that irreversible operations work stays behind "
    "human gates — the rule from the previous slides.",
)
NEW_EVAL = (
    "Evaluating Agents",
    ["Chapter 6 covers evaluation in depth — this slide is the agent-specific preview",
     "Trajectory evals: did the agent pick the right tools, in the right order, with valid arguments?",
     "Outcome evals: task success rate on a golden set of end-to-end tasks",
     "Cost & safety evals: steps and tokens per task, guardrail violations, human-gate triggers",
     "Evals are the regression suite for agents — rerun them on every prompt, tool, or model change"],
    "Point firmly at Chapter 6 for the full evals treatment; here, give only the agent-specific "
    "twist. Agents fail in the trajectory as well as the answer, so evals must check the steps, "
    "not just the output. Frame evals as the regression suite that makes agent changes safe.",
)
LAB_TITLE = "Lab 7.1: Build a Guarded Agent"
LAB_BULLETS = [
    "Follow the detailed instructions in the Lab 7.1 notebook on your VM",
    "Wire tools behind JSON schemas and an allow-list — least privilege by default",
    "Run the ReAct loop with step and spend caps",
    "Gate irreversible actions behind human approval",
    "Watch the guardrails catch a prompt-injection attack",
]
LAB_NOTES = (
    "The capstone lab for the chapter: students assemble a small agent with the exact "
    "guardrails from the lecture — allow-listed tools, loop caps, and a human approval gate — "
    "then fire a prompt-injection attack at it and watch the controls hold. Circulate during "
    "the injection stage; that is where the lesson lands."
)


# --- formatting-preserving text rewrite helpers (same as build_ch03) -----------
def _set_para_text(p_el, text):
    runs = p_el.findall(qn("a:r"))
    if not runs:
        r = p_el.makeelement(qn("a:r"), {})
        r.append(p_el.makeelement(qn("a:t"), {}))
        p_el.append(r)
        runs = [r]
    first = runs[0]
    for extra in runs[1:]:
        p_el.remove(extra)
    for br in p_el.findall(qn("a:br")):
        p_el.remove(br)
    t = first.find(qn("a:t"))
    if t is None:
        t = first.makeelement(qn("a:t"), {})
        first.append(t)
    t.text = text


def rewrite_bullets(shape, lines):
    txBody = shape.text_frame._txBody
    paras = txBody.findall(qn("a:p"))
    for i, line in enumerate(lines):
        if i < len(paras):
            _set_para_text(paras[i], line)
        else:
            newp = copy.deepcopy(paras[-1])
            txBody.append(newp)
            paras.append(newp)
            _set_para_text(newp, line)
    for p in paras[len(lines):]:
        txBody.remove(p)


def _iter_paras(shape):
    """Yield every a:p element in a shape tree (recurses into group shapes)."""
    if shape.shape_type == 6:  # GROUP
        for sub in shape.shapes:
            yield from _iter_paras(sub)
    elif shape.has_text_frame:
        yield from shape.text_frame._txBody.findall(qn("a:p"))


def patch_slide(slide, pairs, tag):
    """Substring-replace inside paragraphs, preserving first-run formatting."""
    missed = []
    paras = [p for sh in slide.shapes for p in _iter_paras(sh)]
    for old, new in pairs:
        hit = False
        for p in paras:
            joined = "".join(t.text or "" for t in p.iter(qn("a:t")))
            if old in joined:
                _set_para_text(p, joined.replace(old, new))
                hit = True
                break
        if not hit:
            missed.append(old[:60])
    if missed:
        print(f"  !! patch MISS on {tag}: {missed}")


def main():
    # 0. fresh carrier -> idempotent
    shutil.copy(BAK / "1851-Ch07.pptx", OUT)
    prs = sk.open_prs(OUT)

    # 1. clone from the three sources (appended at end), patch + retitle
    cidx = {}
    decks = {tag: pt.open_deck(path) for tag, path in SRC.items()}
    for tag, keeps in CLONES.items():
        for k in keeps:
            dst = pt.clone_slide(prs, decks[tag].slides[k - 1])
            cidx[(tag, k)] = len(prs.slides) - 1
            if (tag, k) in PATCHES:
                patch_slide(dst, PATCHES[(tag, k)], f"{tag} s{k}")
            if (tag, k) in RETITLES:
                sk.set_title(dst, RETITLES[(tag, k)])

    # 2. author the 4 new slides (appended at end)
    n_fp = len(prs.slides)
    pt.append_content(prs, NEW_FROM_PROMPTS[0], NEW_FROM_PROMPTS[1], NEW_FROM_PROMPTS[2])
    n_sd = len(prs.slides)
    pt.append_content(prs, NEW_SDLC[0], NEW_SDLC[1], NEW_SDLC[2])
    n_ev = len(prs.slides)
    pt.append_content(prs, NEW_EVAL[0], NEW_EVAL[1], NEW_EVAL[2])
    # lab slide on the exercise layout
    lab_layout = sk.pick_layout(prs, "Exercise Reference Slide with Typing Hands")
    lab = prs.slides.add_slide(lab_layout)
    sk.set_title(lab, LAB_TITLE)
    body = pt._body_placeholder(lab)
    if body is None:
        body = lab.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(11.5), Inches(5.2))
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(LAB_BULLETS):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
    sk.set_notes(lab, LAB_NOTES)
    n_lab = len(prs.slides) - 1

    # 3. rewrite kept carrier slides in place (0-based: 0,1,2,26,27)
    s1 = prs.slides[0]
    _set_para_text(s1.shapes.title.text_frame._txBody.findall(qn("a:p"))[0],
                   "AI Agents and Agentic Workflows")
    rewrite_bullets(pt._body_placeholder(prs.slides[1]), OBJECTIVES)
    rewrite_bullets(pt._body_placeholder(prs.slides[2]), AGENDA)
    rewrite_bullets(pt._body_placeholder(prs.slides[26]), SUMMARY)
    rewrite_bullets(pt._body_placeholder(prs.slides[27]), OBJECTIVES)

    # 4. arrange into the final teaching order
    order = (
        [0, 1, 2, n_fp]                                                     # openers + bridge
        + [cidx[("1258", k)] for k in [36, 37]]                             # agent intro
        + [cidx[("16_1", k)] for k in [21, 22, 23, 24, 25]]                 # loop/anatomy/memory/ReAct
        + [cidx[("16_1", 30)], cidx[("1258", 12)], cidx[("1258", 38)],
           cidx[("16_1", 32)], cidx[("1258", 20)], cidx[("1258", 43)],
           cidx[("1258", 22)], cidx[("1258", 46)], cidx[("1258", 47)]]      # tools & protocols
        + [cidx[("1258", k)] for k in [40, 41, 42]]                         # patterns
        + [cidx[("16_1", 26)], cidx[("16_1", 27)],
           cidx[("16_5", 22)], cidx[("16_5", 23)],
           cidx[("16_5", 28)], cidx[("16_5", 29)]]
        + [cidx[("1258", 44)], cidx[("1258", 45)]]                          # HITL, least privilege
        + [cidx[("16_5", k)] for k in [25, 26, 27, 41, 42]]                 # plumbing + stack
        + [n_sd]                                                            # agents in the SDLC
        + [cidx[("16_5", 43)], cidx[("16_5", 44)]]                          # human-on-the-loop rules
        + [cidx[("16_1", k)] for k in [37, 38, 39, 40, 41]]                 # vendor/OpenAI stack
        + [cidx[("16_1", 35)], cidx[("1258", 49)], cidx[("1258", 50)]]      # failure modes
        + [n_ev, n_lab]                                                     # eval pointer + lab
        + [26, 27]                                                          # summary + recap
    )
    pt.arrange(prs, order)
    pt.save(prs, OUT)

    # 5. verify: reopen, titles, blanks, duplicate partnames, zip integrity
    prs2 = Presentation(OUT)
    n = len(prs2.slides)
    print(f"\n{OUT.name}: {n} slides")
    print(f"clones: 2016-Ch01={len(CLONES['16_1'])} 2016-Ch05={len(CLONES['16_5'])} "
          f"1258-Ch07={len(CLONES['1258'])}  new=4  carrier-kept=5")
    empty = []
    for i, s in enumerate(prs2.slides, 1):
        t = sk.slide_title(s)
        if not t.strip():
            empty.append(i)
        print(f"{i:>3} | {t}")
    assert not empty, f"empty titles at {empty}"
    assert 46 <= n <= 54, f"slide count {n} outside 46-54"
    with zipfile.ZipFile(OUT) as z:
        dupes = [p for p, c in Counter(z.namelist()).items() if c > 1]
        assert not dupes, f"duplicate zip partnames: {dupes}"
        assert z.testzip() is None, "zip integrity check failed"
    print("\nOK: no empty titles, no duplicate partnames, zip integrity verified")


if __name__ == "__main__":
    main()
