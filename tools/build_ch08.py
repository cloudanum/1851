#!/usr/bin/env python3
"""build_ch08.py — refresh Ch08 'AI Security and Vulnerability Testing' for 1851 A.2.

Carrier = _bak0/1851-Ch08.pptx (restored fresh every run -> idempotent).
Keeps the proven A.1 core (incl. all three activities: ATLAS Do Now s20,
Activity 8.1 s31, Activity 8.2 s47). Cuts 5 dated/redundant slides, clones 4
agentic-era slides from the 2016 'Agentic Security' enriched decks, authors 2
new slides, and rewrites the Summary with agentic-era points.

Cuts (1-based carrier):  9 NIST predictive-taxonomy figure (redundant w/ s10-11),
16 FGSM (dated technique), 17 Risks of Adversarial Examples (dup of s15),
30 Broader Implications of Data Minimization (recap dup), 45 Threat Modeling
(dup of s37, truncated text).

Clones (2016 enriched):  Ch03 s15 Indirect Prompt Injection (Greshake 2023),
Ch03 s22 Agent Hijacking — and the Guard Against It, Ch06 s44 Determinism
Sandwich, Ch06 s50 Zero Trust Reference Architecture for an Agent.
Clone notes are ported with minimal edits that remove 2016-only lab/section
references (Lab 3.1, capstone, PQ section) — everything else kept verbatim.

Run:  /Users/iahmad/Creator/Courses_and_conferences/LT/.venv-courseware/bin/python tools/build_ch08.py
"""
from __future__ import annotations
import copy
import re
import shutil
import sys
import zipfile
from collections import Counter
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import slidekit as sk  # noqa: F401  (import activates the partname-collision fix)
import pptx_tools as pt
from pptx import Presentation
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
BAK = ROOT / "decks" / "_bak0"
OUT = ROOT / "decks" / "1851-Ch08.pptx"
SRC16 = Path("/Users/iahmad/Creator/Courses_and_conferences/LT/courses/"
             "2016-agentic-security/2016-slides/enriched")

# --- clone plan: (source deck, 1-based slide, on-slide text edits, notes edits)
CLONES = [
    {
        "src": SRC16 / "Ch03-Pass5.pptx", "idx": 15, "key": "indirect",
        "text_edits": [],
        "notes_edits": [
            (" -- and it has existing notes worth incorporating, so let me "
             "build on them. The core idea:", ". The core idea:"),
        ],
    },
    {
        "src": SRC16 / "Ch03-Pass5.pptx", "idx": 22, "key": "hijack",
        "text_edits": [
            ("Lab 3.1 defense: an input-validation guard",
             "Defense in practice: an input-validation guard"),
        ],
        "notes_edits": [
            ("the guard you build against it in Lab 3.1.",
             "the guard you build against it."),
            ("because it is the structure students must own -- it is on the "
             "chapter review.", "because it is the structure students must own."),
            ("this is why those two OWASP entries are the backbone of the "
             "chapter.", "this is why injection and excessive agency are the "
             "backbone of agentic risk."),
            ("which is exactly the lesson from the multi-language and "
             "refusal-suppression slides.",
             "which is exactly what a plain keyword filter misses."),
            ("and this is the constructive heart of the lab.",
             "and this is the constructive heart of the defense."),
            ("Lab facilitation: in Part 6 you add this guard to the sandbox "
             "agent and re-test all five injection variants from earlier. The "
             "'done' criterion is that the guard blocks the obvious injections "
             "and the out-of-allowlist tool calls while still permitting "
             "benign requests -- and that you can articulate which variant "
             "slips past and why. Expect the multi-language and "
             "cleverly-reworded ones to challenge the similarity check; that "
             "is the intended lesson.", ""),
        ],
    },
    {
        "src": SRC16 / "Ch06-Final.pptx", "idx": 44, "key": "sandwich",
        "text_edits": [],
        "notes_edits": [
            ("This single abstraction connects everything - the previous "
             "slide's micro-segmentation is the top bread, the upcoming "
             "Verifier and output checks are the bottom bread, and Chapter 5's "
             "Planner-Executor-Verifier loop was an instance of exactly this "
             "sandwich. Each slice of bread is a Zero Trust boundary.",
             "This single abstraction connects everything in this block - "
             "input filtering is the top bread, output validation is the "
             "bottom bread. Each slice of bread is a Zero Trust boundary."),
            ("Concrete example: the EDD judge from earlier is probabilistic "
             "filling - so you wrap it in deterministic schema and allow-list "
             "checks; you would never let the judge alone gate a release, just "
             "as you would never let the model alone authorize an irreversible "
             "action.",
             "Concrete example: an LLM-as-judge evaluation (Chapter 6) is "
             "probabilistic filling - so you wrap it in deterministic schema "
             "and allow-list checks; you would never let the judge alone gate "
             "a release, just as you would never let the model alone authorize "
             "an irreversible action."),
            ("Transition: the most important slice of deterministic bread for "
             "governance is the audit trail.",
             "Transition: let's turn these layers into practical "
             "prompt-injection defenses."),
        ],
    },
    {
        "src": SRC16 / "Ch06-Final.pptx", "idx": 50, "key": "zerotrust",
        "text_edits": [],
        "notes_edits": [
            (", and note that the capstone lab checks the most critical of "
             "these layers directly.", "."),
            ("Crypto: TLS today, with a planned PQC migration to FIPS "
             "203/204/205 - which is the on-ramp to the final section.",
             "Crypto: TLS today, with a planned PQC migration to FIPS "
             "203/204/205."),
            ("Stress that this is not aspirational - every layer here was "
             "built in this section, and together they are a coherent stack "
             "you can deploy.",
             "Stress that this is not aspirational - each layer maps to a "
             "control from this chapter, and together they are a coherent "
             "stack you can deploy."),
            ("Make the lab connection explicit: the capstone's "
             "zero_trust_audit checks that the audit plane is present on every "
             "component - and in the sample, the executor-agent's missing "
             "audit plane is what BLOCKS go-live. So this reference "
             "architecture is not just a diagram; its audit-plane row is the "
             "literal pass/fail line in the lab.",
             "Stress the audit plane: an agent without a signed, append-only "
             "audit trail on every component is not ready for production - "
             "treat the audit row of this architecture as a go/no-go line in "
             "any deployment review."),
            ("Concrete example: drop the incident-response agent into this "
             "architecture and every earlier control finds its home - the "
             "quarantine request flows through PDP, input boundary, action "
             "gate, and audit plane in order, exactly as the worked example "
             "showed.",
             "Concrete example: drop any tool-calling agent into this "
             "architecture and every control finds its home - each action "
             "flows through PDP, input boundary, action gate, and audit plane "
             "in order."),
            ("no - harvest-now-decrypt-later means today's crypto choices have "
             "a multi-year exposure, which the next section explains.",
             "no - harvest-now-decrypt-later means today's crypto choices "
             "have a multi-year exposure."),
            ("Transition: that crypto row opens the final section - "
             "post-quantum and advanced threats.",
             "Transition: with the architecture in place, let's secure the "
             "workflows that run on it."),
        ],
    },
]

# --- 2 new authored slides (title, bullets, speaker notes) -------------------
NEW_SLIDES = [
    ("Prompt Injection Defenses in Practice",
     ["Filter inputs — pattern matching and semantic similarity against known injection signatures",
      "Validate outputs — schema-check responses and tool-call arguments before anything executes",
      "Separate privileges — the model proposes actions; a deterministic layer approves and executes them",
      "Treat retrieved content as data — isolate and quote it; never blend it into the instruction channel",
      "Recap: the determinism sandwich — deterministic guardrails above and below the probabilistic core"],
     "Move from naming defenses to wiring them into a real application: no single layer stops "
     "injection, so engineer for the failure of each one. Point back to the Agent Hijacking guard "
     "as a concrete instance of input filtering plus a tool allowlist. Close with the determinism "
     "sandwich as the organizing pattern — trust lives in the deterministic bread, never in the "
     "model's judgment."),
    ("Securing Agentic Workflows",
     ["Scope every tool to least privilege — allowlisted actions, narrow credentials, no standing access",
      "Gate high-impact actions on humans — approvals for payments, deletions, and external messages",
      "Log everything — prompts, tool calls, and outputs as signed, tamper-evident audit events",
      "Segment trust between agents — validate inter-agent messages like any other untrusted input",
      "Bound the blast radius — rate limits, spend caps, and sandboxed execution for code tools"],
     "This slide complements Chapter 7's agentic workflows with the controls that make them safe "
     "to operate. The organizing idea is least privilege applied to tools rather than users: an "
     "agent should hold exactly the permissions its current task requires. Human gates plus signed "
     "audit logs turn raw autonomy into accountable autonomy."),
]

# --- rewritten Summary (carrier s50) with agentic-era points -----------------
SUMMARY = [
    "Testing confirms AI reliability and supports safe deployment, especially in public and regulatory environments",
    "Adversarial threats like evasion and data poisoning require proactive detection and defensive strategies",
    "Prompt injection—direct and indirect—is the top GenAI threat; against agents it escalates from bad output to hijacked actions",
    "Defend in depth: filter inputs, validate outputs, enforce least-privilege tool scopes, and wrap the probabilistic core in deterministic guardrails",
    "Zero Trust for agents: verify every action, gate high-impact operations with human approval, and audit-log every step",
    "Explainability tools such as LIME and SHAP provide insights into model decisions, enhancing trust and auditability",
    "Compliance frameworks (e.g., AIACT, EU AI Act) demand traceability, oversight, and documentation",
    "Ethical AI demands transparency, robustness, and clear communication with both technical and non-technical audiences",
]


# --- formatting-preserving text rewrite helpers (same pattern as build_ch03) --
def _set_para_text(p_el, text):
    """Replace a paragraph's text, keeping the first run's formatting."""
    runs = p_el.findall(qn("a:r"))
    if not runs:
        r = p_el.makeelement(qn("a:r"), {})
        r.append(p_el.makeelement(qn("a:t"), {}))
        p_el.append(r)
        runs = [r]
    first = runs[0]
    for extra in runs[1:]:
        p_el.remove(extra)
    for br in p_el.findall(qn("a:br")):
        p_el.remove(br)
    t = first.find(qn("a:t"))
    if t is None:
        t = first.makeelement(qn("a:t"), {})
        first.append(t)
    t.text = text


def rewrite_bullets(shape, lines):
    """Rewrite a body placeholder's bullets, reusing paragraph formatting."""
    txBody = shape.text_frame._txBody
    paras = txBody.findall(qn("a:p"))
    for i, line in enumerate(lines):
        if i < len(paras):
            _set_para_text(paras[i], line)
        else:
            newp = copy.deepcopy(paras[-1])
            txBody.append(newp)
            paras.append(newp)
            _set_para_text(newp, line)
    for p in paras[len(lines):]:
        txBody.remove(p)


def replace_on_slide(slide, old, new):
    """Run-level text replace across all shapes; falls back to paragraph level."""
    hits = 0
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if old in r.text:
                    r.text = r.text.replace(old, new)
                    hits += 1
            if old in p.text and not any(old in r.text for r in p.runs):
                _set_para_text(p._p, p.text.replace(old, new))
                hits += 1
    return hits


def edit_notes(slide, pairs, label):
    nt = slide.notes_slide.notes_text_frame
    txt = nt.text
    for old, new in pairs:
        if old in txt:
            txt = txt.replace(old, new)
        else:
            print(f"  WARN [{label}]: notes string not found: {old[:70]}...")
    txt = re.sub(r"\n{3,}", "\n\n", txt).strip()
    nt.text = txt


def main():
    # 0. fresh carrier -> idempotent
    shutil.copy(BAK / "1851-Ch08.pptx", OUT)
    prs = sk.open_prs(OUT)

    # 1. clone the 4 agentic-era slides from 2016 (appended at end)
    clone_idx = {}
    for spec in CLONES:
        src = pt.open_deck(spec["src"])
        new_slide = pt.clone_slide(prs, src.slides[spec["idx"] - 1])
        for old, new in spec["text_edits"]:
            n = replace_on_slide(new_slide, old, new)
            if n == 0:
                print(f"  WARN [{spec['key']}]: slide text not found: {old[:60]}...")
        if spec["notes_edits"]:
            edit_notes(new_slide, spec["notes_edits"], spec["key"])
        clone_idx[spec["key"]] = len(prs.slides) - 1

    # 2. author the 2 new slides (appended at end)
    new_idx = {}
    for title, bullets, notes in NEW_SLIDES:
        pt.append_content(prs, title, bullets, notes)
        new_idx[title] = len(prs.slides) - 1

    # 3. rewrite the Summary (carrier s50, 0-based 49) with agentic-era points
    rewrite_bullets(pt._body_placeholder(prs.slides[49]), SUMMARY)

    # 4. arrange into final order (0-based carrier indices; cuts simply absent)
    #    cuts: 8 (s9), 15 (s16 FGSM), 16 (s17), 29 (s30), 44 (s45)
    order = (
        list(range(0, 8))           # title, objectives, agenda, 4 dimensions, call to action
        + [9, 10]                   # attacks on predictive AI: poisoning, privacy
        + [11, 12]                  # attacks on generative AI
        + [13, 14]                  # ATLAS intro, what are adversarial examples
        + [17]                      # adversarial robustness strategies
        + [18, 19]                  # ATLAS landscape, ATLAS Matrix (Do Now activity)
        + list(range(20, 29))       # CIA triad, privacy, securing AI, data minimization
        + [30]                      # Activity 8.1: Data Minimization Checklist
        + list(range(31, 39))       # datasets, controls, SecOps, policies, threat modeling,
                                    # prompt engineering, hacking prompts/injection risks
        + [clone_idx["indirect"], clone_idx["hijack"]]
        + [39]                      # defenses against prompt-based attacks
        + [clone_idx["sandwich"], new_idx["Prompt Injection Defenses in Practice"],
           clone_idx["zerotrust"], new_idx["Securing Agentic Workflows"]]
        + list(range(40, 44))       # PII block
        + [45]                      # insider threats
        + [46]                      # Activity 8.2: Detect and Mask PII
        + [47, 48]                  # suspicious behavior, insider-threat mitigation
        + [49, 50]                  # summary (rewritten), objectives recap
    )
    pt.arrange(prs, order)
    pt.save(prs, OUT)

    # 5. verify: reopen, titles, blanks, activities, duplicate partnames, zip
    prs2 = Presentation(OUT)
    n = len(prs2.slides)
    print(f"\n{OUT.name}: {n} slides "
          f"(kept 46 carrier + {len(CLONES)} clones + {len(NEW_SLIDES)} new)")
    empty = []
    titles = []
    for i, s in enumerate(prs2.slides, 1):
        t = sk.slide_title(s)
        titles.append(t)
        if not t.strip():
            empty.append(i)
        print(f"{i:>3} | {t}")
    assert not empty, f"empty titles at {empty}"
    assert 48 <= n <= 52, f"slide count {n} outside 48-52"
    for req in ("Data Minimization Checklist", "Activity 8.2", "ATLAS Matrix"):
        assert any(req in t for t in titles), f"missing activity: {req}"
    with zipfile.ZipFile(OUT) as z:
        dupes = [p for p, c in Counter(z.namelist()).items() if c > 1]
        assert not dupes, f"duplicate zip partnames: {dupes}"
        assert z.testzip() is None, "zip integrity check failed"
    print("\nOK: no empty titles, all three activities present, "
          "no duplicate partnames, zip integrity verified")


if __name__ == "__main__":
    main()
