#!/usr/bin/env python3
"""Fix Ch03 slides with invisible body text (missing xfrm extent -> 0x0 at -1,-1)
and enrich them with SVG-derived diagrams.

Targets slides 38, 45, 50, 51, 52 (1-based) of 1851-Ch03.pptx.
Idempotent: geometry is always set; diagram pictures are inserted only once
(recognized by shape name 'diagram-sNN').

Usage: fix_ch03_empty.py <path-to-1851-Ch03.pptx>
"""
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

DIAGRAM_DIR = Path(__file__).parent / 'diagrams'

# slide index -> (content geometry L,T,W,H in inches, diagram png, diagram pos L,T,W)
RIGHT_PIC = (0.29, 1.04, 5.05, 2.10)   # bullets beside right-hand stock photo
FIXES = {
    38: (RIGHT_PIC, 's38-eval-loop.png', (0.44, 3.30, 4.75)),
    45: (RIGHT_PIC, 's45-metric-selection.png', (0.44, 3.30, 4.75)),
    50: (RIGHT_PIC, 's50-forward-prop.png', (0.44, 3.30, 4.75)),
    52: (RIGHT_PIC, 's52-training-loop.png', (0.44, 3.30, 4.75)),
    # s51 has a full-width photo across the top -> bullets below-left, diagram below-right
    51: ((0.29, 4.02, 4.75, 2.80), 's51-backprop.png', (5.15, 3.95, 4.50)),
}


def main(path: str) -> None:
    prs = Presentation(path)
    changed = 0
    for idx, ((cl, ct, cw, ch), png, (pl, pt, pw)) in FIXES.items():
        slide = prs.slides[idx - 1]
        for sh in slide.shapes:
            if sh.name.startswith('Content') and sh.has_text_frame:
                sh.left, sh.top = Inches(cl), Inches(ct)
                sh.width, sh.height = Inches(cw), Inches(ch)
                changed += 1
        tag = f'diagram-s{idx}'
        img = DIAGRAM_DIR / png
        with Image.open(img) as im:
            aspect = im.height / im.width
        existing = next((sh for sh in slide.shapes if sh.name == tag), None)
        if existing is None:
            existing = slide.shapes.add_picture(
                str(img), Inches(pl), Inches(pt), width=Inches(pw))
            existing.name = tag
        else:
            existing.left, existing.top = Inches(pl), Inches(pt)
            existing.width, existing.height = Inches(pw), Inches(pw * aspect)
        changed += 1
    prs.save(path)
    print(f'{path}: {changed} fixes applied')


if __name__ == '__main__':
    main(sys.argv[1])
