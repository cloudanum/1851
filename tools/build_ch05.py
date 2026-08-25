#!/usr/bin/env python3
"""build_ch05.py — author the NEW A.2 chapter 'Generative AI Across the
Software Development Lifecycle' (TDD as its centerpiece).

Carrier = _bak0/1851-Ch04.pptx (A.1 'Data Preparation and Handling', restored
fresh every run -> idempotent). The A.1 content moved to the new Ch03, so the
carrier is used purely as an LT-branded shell: only the openers (title /
objectives / agenda) and closers (summary / objectives recap) are kept and
rewritten in place; every old content slide is dropped via pt.arrange.

All 45 body slides are authored fresh on the LT layouts:
  'Content with Header. Full Page'  (bullets, 4-6 per slide)
  'Two Column Full Page'            (good/bad and IDE-vs-API compares)
  'Exercise Reference Slide with Typing Hands' (Lab 5.1)
  + one add_table_slide for the phase map. Every authored slide carries 2-4
  sentence speaker notes. Forward references stay consistent: Ch04 prompting /
  OpenAI API / RAG, Ch06 evals, Ch07 agents, Ch08 security.

Run:  /Users/iahmad/Creator/Courses_and_conferences/LT/.venv-courseware/bin/python tools/build_ch05.py
"""
from __future__ import annotations
import copy
import shutil
import sys
import zipfile
from collections import Counter
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import slidekit as sk  # noqa: F401  (import activates the partname-collision fix)
import pptx_tools as pt
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent
BAK = ROOT / "decks" / "_bak0"
CARRIER = BAK / "1851-Ch04.pptx"
OUT = ROOT / "decks" / "1851-Ch05.pptx"

# --- rewritten text for kept carrier slides (idx 0,1,2 and 61,62) ------------
TITLE = "Generative AI Across the Software Development Lifecycle"
SUBTITLE = "Chapter 5"
TITLE_NOTES = (
    "Welcome to Chapter 5, the heart of the A.2 revision. Earlier chapters covered what modern "
    "AI is (Chapters 1-3) and how to drive it through prompting and the OpenAI API (Chapter 4); "
    "this chapter applies all of it to the work software engineers do every day, phase by phase. "
    "The throughline: tests and evals are what make AI-generated work trustworthy.")
OBJECTIVES = [
    "Map each SDLC phase to concrete generative AI uses, from requirements to operations",
    "Apply LLMs to requirements work: user stories, acceptance criteria, and ambiguity detection",
    "Use AI for design: architecture options, trade-off analysis, ADRs, and review-ready diagrams",
    "Get reliable code from assistants with small units, clear contracts, and examples in the prompt",
    "Practice TDD with an AI pair — tests as the spec that keeps generated code honest",
    "Extend GenAI to review, docs, CI/CD, and operations — and adopt it behind quality gates",
]
OBJECTIVES_NOTES = (
    "Walk the objectives as the chapter's arc: we map the SDLC, then work it end to end — "
    "requirements, design, coding, testing, review, and operations. Objective five is the "
    "centerpiece: TDD with an AI pair is how we keep generated code honest, and it connects "
    "directly to evals in Chapter 6. Set the expectation that engineers stay accountable "
    "throughout — the AI drafts, the engineer decides.")
AGENDA = [
    "The GenAI-augmented SDLC and the engineer's new role",
    "Requirements and planning with LLM support",
    "AI-assisted design: options, trade-offs, and ADRs",
    "Coding with assistants and the OpenAI API",
    "Testing and TDD with an AI pair — the heart of this chapter",
    "Review, documentation, CI/CD, and operations",
    "Adoption: metrics, quality gates, and Lab 5.1",
]
AGENDA_NOTES = (
    "Preview the seven stops. Flag that the testing and TDD section is the longest for a "
    "reason: it is where GenAI either becomes trustworthy engineering or stays a demo. "
    "The chapter closes with adoption — how a team measures impact and sets norms — and "
    "Lab 5.1 puts the TDD loop in students' hands.")
SUMMARY = [
    "Generative AI augments every SDLC phase — but the engineer specifies, verifies, and owns the result",
    "Requirements and design: AI drafts stories, criteria, options, and ADRs; humans decide and sign off",
    "Coding: small units, clear contracts, examples in the prompt — then read and run everything",
    "Tests are the spec: TDD with an AI pair turns generated code into trustworthy code",
    "Evals (Chapter 6) are TDD for AI behavior; agents (Chapter 7) automate the workflow toil",
    "Adopt deliberately: measure impact, keep quality gates, and know when not to use GenAI",
]
SUMMARY_NOTES = (
    "Land the three messages. First, GenAI helps in every phase — this chapter walked the whole "
    "lifecycle. Second, tests and evals are the quality backbone: they are what let you trust "
    "code and behavior you did not write. Third, accountability never moved — the engineer who "
    "ships it owns it. Point forward: Chapter 6 turns test thinking into evals, Chapter 7 into "
    "agents, Chapter 8 into security.")
RECAP_NOTES = (
    "Close by mapping back to the objectives. Ask the room which objective they will apply first "
    "on their own codebase. Then bridge to Lab 5.1 if not already done, and to Chapter 6, where "
    "'tests are the spec' becomes 'evals are the spec' for AI behavior.")

# --- 45 authored body slides, in final teaching order --------------------------
# kinds: b = bullets | tc = two-column | tbl = phase-map table | lab = exercise
SLIDES = [

# ===== FRAMING (6) =====
("b", "The SDLC at a Glance",
 ["Requirements and planning — what to build, for whom, and why",
  "Design — architecture, interfaces, data models, and trade-offs",
  "Implementation — writing, reading, and refactoring code",
  "Testing and review — verifying behavior and holding the quality bar",
  "Deployment and operations — ship, monitor, respond, maintain",
  "Every phase produces text and structured artifacts — and GenAI works on exactly those"],
 "A quick refresher so the room shares one vocabulary for the lifecycle — keep it to a minute. "
 "The key point is the last bullet: every phase is language-plus-structure work, which is precisely "
 "what generative models are good at. That observation is the chapter's entire premise."),

("tbl", "The GenAI-Augmented SDLC",
 (["Phase", "GenAI uses"],
  [["Requirements", "Draft user stories; detect ambiguity and gaps; generate acceptance criteria"],
   ["Design", "Draft architecture options; analyze trade-offs; write ADRs and review diagrams"],
   ["Coding", "Generate, explain, refactor, and translate code with assistants"],
   ["Testing", "Generate unit tests, edge cases, and test data; TDD with an AI pair"],
   ["Review", "AI-assisted PR review; standards checks; hallucination hunting"],
   ["Deploy & Ops", "CI/CD assistance, release notes, incident triage, log summarization"]]),
 "This table is the map for the rest of the chapter — each row becomes a section. Do not read "
 "every cell; give one example per row and flag that the testing row gets the deepest treatment. "
 "Tell students to mark the rows where their own team loses the most time."),

("b", "Levels of AI Assistance",
 ["Autocomplete — next-token suggestions in the editor: fast, low-context, always on",
  "Chat — ask, paste code, iterate: the engineer drives, the model advises",
  "Task agents — hand over a goal ('fix this issue'); the model plans, edits, and runs tests in a loop",
  "Each level up gets more context and more autonomy — and demands stronger verification",
  "This chapter works all three levels; Chapter 7 goes deep on agents"],
 "Draw the ladder: autocomplete, chat, agentic. The practical point is the fourth bullet — "
 "assistance and verification scale together, so the more you let the AI do, the more you need "
 "the testing discipline this chapter teaches. Defer agent mechanics to Chapter 7."),

("b", "What Changed by 2026",
 ["Context windows in the hundreds of thousands of tokens — whole modules fit in one prompt",
  "Reasoning models (GPT-5.x era) plan multi-step changes instead of guessing the next token",
  "Tool use is native: assistants run tests, read files, and call APIs — not just chat",
  "The OpenAI Responses API (Chapter 4) makes text, reasoning, and tools one call shape",
  "Assistants are now default tooling — Copilot, Codex, Cursor, and API-driven scripts"],
 "Why is a chapter like this possible now when it was not a few years ago? Three capability "
 "shifts: long context, reasoning, and native tool use — plus the protocols (MCP and friends) "
 "that let assistants reach real dev tools. Avoid model-version trivia; the point is that these "
 "shifts are what the rest of the chapter exploits."),

("tc", "What GenAI Is Good — and Bad — At",
 (("Strong at",
   ["Boilerplate, glue code, and well-trodden patterns",
    "Explaining unfamiliar code and error messages",
    "Drafting tests, docs, and migration scripts",
    "Exploring options — many fast drafts to react to"]),
  ("Weak at",
   ["Inventing APIs and libraries that do not exist",
    "Subtle logic errors that look plausible",
    "Business context it was never given",
    "Knowing when it is wrong — it will not tell you"])),
 "The calibration slide — engineers who skip this slide over-trust or under-use the tools for "
 "the rest of the course. The right column drives everything that follows: because the model "
 "fails confidently and silently, verification is not optional. That is why tests are the "
 "centerpiece of this chapter."),

("b", "The Engineer's New Role: Director of AI Output",
 ["Specify — precise intent, contracts, and constraints beat clever prompts",
  "Verify — read, run, and test everything; unreviewed AI code is still your code",
  "Own — accountability for shipped behavior stays with the engineer and the team",
  "The skill shift: less typing, more reviewing, designing, and deciding",
  "Tests and evals are the verification backbone — the theme of this chapter"],
 "Reframe the job, not eliminate it: the engineer moves from author to director of AI output. "
 "Specify, verify, own — those three verbs repeat all chapter. Land the last bullet hard: the "
 "verification backbone is tests for code and evals for AI behavior, and both are coming."),

# ===== REQUIREMENTS & PLANNING (5) =====
("b", "Drafting User Stories with LLMs",
 ["Feed context: persona, problem, and constraints — get a first-draft story set in seconds",
  "Ask for the standard shape: 'As a <role>, I want <capability>, so that <benefit>'",
  "Iterate conversationally: split epics, merge duplicates, sharpen scope",
  "Paste real artifacts (support tickets, call notes) for grounded drafts — mind privacy (Chapter 8)",
  "The draft is a starting point for refinement, not a finished backlog"],
 "Show the workflow: context in, story draft out, then refine by conversation. The win is speed "
 "to a first draft that humans then argue with — refinement is still a human meeting. Remind "
 "students that pasting real customer artifacts into a prompt has privacy implications; "
 "Chapter 8 covers that properly."),

("b", "Generating Acceptance Criteria",
 ["Acceptance criteria turn a story into testable statements — Given/When/Then works well",
  "Prompt with the story plus domain rules; ask for happy-path, edge, and failure cases",
  "Example: Given a locked account, when the user logs in, then show unlock guidance",
  "Check every criterion against real business rules — the model invents plausible ones",
  "These criteria feed directly into the TDD workflow later in this chapter"],
 "Acceptance criteria are where requirements become executable thinking. Walk one Given/When/Then "
 "example end to end, then stress the review bullet: the model will cheerfully invent rules that "
 "sound right. Plant the forward pointer — these criteria become the tests in the TDD section."),

("b", "Ambiguity and Gap Detection in Requirements",
 ["LLMs are strong critics: 'List every ambiguity, contradiction, and gap in this requirement'",
  "Ask for the questions a new hire would ask — it surfaces implicit assumptions",
  "Probe negative cases: what must the system NOT do, and what happens on failure?",
  "Cross-check stories against each other for duplicates, conflicts, and missing flows",
  "Cheap at requirements time — the same defects cost 10-100x more in production"],
 "This is the highest-value requirements use: the model as tireless critic rather than author. "
 "Give the new-hire prompt verbatim — it reliably exposes assumptions the team forgot they were "
 "making. Close on the economics: an ambiguity fixed here costs minutes; found in production it "
 "costs an incident."),

("b", "Estimation and Risk Support",
 ["Ask for a work breakdown: components, integration points, and unknowns",
  "Surface risk classes — third-party dependencies, data migration, concurrency, scale",
  "Have it argue against your estimate: 'What am I underestimating, and why?'",
  "Never outsource the number — estimates commit the team and need human judgment",
  "Best use: a structured checklist of what to think about, not a magic estimate"],
 "Be careful here — LLMs sound confident about estimates and are often wrong. Position the model "
 "as a structured-thinking partner: breakdowns, risk lists, and adversarial challenge, while the "
 "team owns the actual number. The 'argue against my estimate' prompt is the one to demo live."),

("b", "Human Accountability for Requirements",
 ["The model does not know your users, regulators, or business strategy — context is your job",
  "Requirement sign-off stays with product owners and engineers, AI draft or not",
  "Wrong requirements fail quietly: everything downstream gets built confidently wrong",
  "Keep a decision log of what was accepted, rejected, and why — ADRs are next",
  "Rule of thumb: AI proposes, humans dispose"],
 "End the section on accountability, because requirements errors are the most expensive kind and "
 "the quietest. 'AI proposes, humans dispose' is the line to remember — it generalizes to the "
 "whole chapter. The decision-log bullet tees up ADRs in the design section."),

# ===== DESIGN (4) =====
("b", "Architecture Drafting and Trade-Off Analysis",
 ["Describe the system, scale, and constraints; ask for two or three candidate architectures",
  "Demand trade-offs, not just diagrams: consistency vs. availability, cost vs. latency",
  "Play adversary: 'Attack this design — where does it break at 10x load?'",
  "Strong for breadth: it will name patterns you forgot to consider",
  "Validate every claim against your real constraints — model benchmarks are generic"],
 "The model is a breadth machine: it recalls more patterns than any one engineer, which makes it "
 "excellent for option generation and devil's-advocate review. But its performance and cost claims "
 "are generic, so every trade-off must be re-derived against your constraints. Demo idea: have it "
 "attack the team's current architecture."),

("b", "API and Schema Design Assistance",
 ["Draft REST endpoints, resource names, and error shapes from a capability description",
  "Generate JSON Schema or OpenAPI specs — and iterate on them conversationally",
  "Ask for a consumer's-eye critique: confusing names, missing pagination, breaking changes",
  "Generate example requests and responses to sanity-check the contract early",
  "The OpenAI API itself speaks JSON Schema (Chapter 4) — a design skill that compounds"],
 "Contract-first design pairs beautifully with LLMs because contracts are structured text. "
 "Emphasize the consumer's-eye critique — naming and pagination mistakes are cheap to catch now "
 "and painful forever after. Note the compounding effect: the JSON Schema skills from Chapter 4 "
 "make students better API designers, not just better API callers."),

("b", "Architecture Decision Records with AI",
 ["ADRs capture context, decision, alternatives, and consequences — AI drafts the skeleton",
  "Feed the debate: paste the options discussed; ask for the strongest case for each",
  "Keep ADRs short and honest — one decision, stated reasons, known downsides",
  "Review ADRs like code: the reasoning, not the prose, is what matters",
  "Future maintainers — and future AI assistants — read ADRs to understand why"],
 "ADRs are the perfect AI-assisted document: highly structured, short, and mostly argument. The "
 "model's job is to steelman every option and format the record; the team's job is the actual "
 "decision. Nice payoff line: well-kept ADRs become grounding material for future AI assistants "
 "working in the repo."),

("b", "Diagrams and Docs for Design Reviews",
 ["Generate diagrams as text: Mermaid, PlantUML, or C4 from a design description",
  "Draft the review narrative: problem, options considered, decision, rollout plan",
  "Ask the model to write the FAQ a skeptical reviewer would ask",
  "Regenerate as the design evolves — docs that cost nothing stay current",
  "Chapter 4's prompting patterns — structure, constraints, grounding — apply verbatim"],
 "Text-renderable diagrams are the unlock: because Mermaid and PlantUML are code, the model can "
 "diff and regenerate them as the design changes, so diagrams stop going stale. The skeptical-FAQ "
 "trick is the sleeper hit — it prepares the presenter for the actual review meeting."),

# ===== CODING (7) =====
("tc", "The 2026 Assistant Landscape",
 (("IDE assistants",
   ["Autocomplete and chat in the editor (Copilot, Cursor, and peers)",
    "Agentic modes edit files and run tests from a goal",
    "Always-on, repo-aware context",
    "Best for: in-flow individual work"]),
  ("API-driven tooling",
   ["Your own scripts and services on the OpenAI API",
    "Batch jobs: migrations, doc sweeps, test generation",
    "Custom gates and steps in CI (later this chapter)",
    "Best for: team-level, repeatable automation"])),
 "Two complementary families, not competitors: IDE assistants help the individual in flow, while "
 "API-driven tooling encodes the team's repeatable work. Students already know the IDE side; the "
 "next slide shows the API side they can build themselves with Chapter 4 skills. Both families "
 "end up in the same pipelines by the ops section."),

("b", "The OpenAI API in Your Own Dev Tools",
 ["Chapter 4's Responses API patterns — messages, structured output, function calling — are the toolkit",
  "Wrap the API in scripts: explain-this-error, generate-migration, summarize-this-log",
  "Use Structured Outputs (strict JSON Schema) when the result feeds another tool",
  "Pin model and parameters inside your tooling for reproducible behavior",
  "Keep secrets in environment variables or a vault — never in prompts or code (Chapter 8)"],
 "This is where Chapter 4 pays rent: the same Responses API calls become small team tools. Walk "
 "one example — a thirty-line script that summarizes a failing test log. Two engineering rules: "
 "structured output when machines consume the result, and pinned models so the tool behaves the "
 "same tomorrow."),

("b", "Code Generation That Works",
 ["Generate small: one function, one class, one migration — never 'write the app'",
  "State the contract: signature, types, inputs/outputs, error behavior",
  "Show, don't tell: a usage example in the prompt beats a paragraph of prose",
  "Ask for the test with the code — if it cannot be tested, the spec is unclear",
  "Read every line before you commit: generated code is a draft, not a delivery"],
 "The five habits that separate productive use from slop generation. Small units plus stated "
 "contracts plus a concrete example — that is most of prompt quality for code. The fourth bullet "
 "is the pivot to the testing section: demanding a test with the code forces spec clarity and "
 "hands you the verification harness for free."),

("b", "Explaining Legacy Code",
 ["Paste a function; ask for a plain-English summary, then a line-by-line walkthrough",
  "Ask for the implied contract: inputs, outputs, side effects, invariants",
  "Interrogate: 'What breaks if this returns null? Who calls this with an empty string?'",
  "Build a dependency map by walking call sites conversationally",
  "Verify explanations against behavior — write a test before you trust the explanation"],
 "The most immediately useful coding use case for most engineers: inheriting code nobody "
 "explains. Teach the interrogation style — the model is better when you cross-examine it than "
 "when you accept the first summary. The last bullet is crucial: an explanation of code you "
 "cannot run is a hypothesis until a test confirms it."),

("b", "Refactoring with LLMs: A Safe Workflow",
 ["Step 1 — Characterize: generate tests that pin down current behavior first",
  "Step 2 — Refactor in small steps: one transformation per prompt, tests green after each",
  "Step 3 — Verify: run the suite, then review the diff line by line",
  "Never refactor and change behavior in the same step",
  "Legacy rule stands: without tests you are not refactoring — you are editing"],
 "Give the three-step workflow as a ritual: characterize, refactor, verify. The characterization "
 "tests are the safety net that makes AI-assisted refactoring safe instead of reckless — and they "
 "stay in the repo as a permanent asset. Quote Feathers' rule on the last bullet; the AI does not "
 "repeal it."),

("b", "Language and Framework Translation",
 ["A strong use case: port idiomatic code across languages or framework versions",
  "Provide the target's conventions: 'use dataclasses', 'prefer error returns over exceptions'",
  "Translate in slices — module by module — with the tests traveling first",
  "Expect semantic drift: concurrency, null handling, and numerics do not map one-to-one",
  "The characterize-refactor-verify workflow from the previous slide applies unchanged"],
 "Translation demos well, which is why it gets over-trusted: syntax ports easily, semantics do "
 "not. The discipline is the same as refactoring — tests travel first, work in slices, verify "
 "each slice. Call out the drift hotspots (concurrency models, nullability, numeric behavior) "
 "because those are where ported code rots."),

("b", "Prompt Patterns for Code",
 ["Role and context: 'You are a senior Go engineer; this is a payment service'",
  "Constraints up front: versions, banned libraries, style rules, performance budget",
  "Few-shot: two input/output examples lock in the format (Chapter 4)",
  "Grounding: paste the real code, schema, or error — never describe it vaguely",
  "Iterate on failure: feed the compiler error or failing test back verbatim"],
 "This slide compresses Chapter 4's prompting discipline into the five patterns that matter most "
 "for code work. Grounding is the big one — most bad code generation comes from describing a "
 "problem instead of pasting it. The iterate-on-failure loop (error message back in verbatim) is "
 "also the exact loop Lab 5.1 automates."),

# ===== TESTING & TDD WITH GENAI (10 — the core) =====
("b", "Why Tests Matter More with AI-Generated Code",
 ["AI writes plausible code, not correct code — plausibility is exactly what tests check",
  "Generated code fails differently: invented APIs and subtly wrong logic, not typos",
  "Your review throughput is now the bottleneck — tests scale review",
  "Tests are executable documentation for code you did not write",
  "Teams that generate more code must generate more tests — same AI, both directions"],
 "This slide justifies why the next nine exist. The core argument: generation moved the "
 "bottleneck from writing to verifying, and tests are the only verification that scales. The "
 "last bullet reframes AI from a test-skipping shortcut into a test-writing accelerator — the "
 "mature teams run it in both directions."),

("b", "Generating Unit Tests from a Function Contract",
 ["Prompt with the function plus its contract: signature, preconditions, error cases",
  "Ask for the matrix: happy path, boundaries, invalid input, error handling",
  "Demand runnable output: imports, fixtures, and assertions that actually assert",
  "Run immediately — treat generated tests as untrusted until they execute",
  "Delete vacuous tests: an assert on a tautology is worse than no test"],
 "The mechanics of test generation: contract in, test matrix out, run it before you trust it. "
 "Spend the most time on the failure modes — tests that do not compile are obvious, but vacuous "
 "asserts pass silently and rot the suite. A suite full of tautologies is a false sense of "
 "security, not coverage."),

("b", "Edge Cases and Property-Based Test Ideas",
 ["LLMs brainstorm well: 'List 15 edge cases for this parser' beats one engineer's memory",
  "Classic lists: empty/null, boundaries, unicode, huge inputs, concurrency, time zones",
  "Ask for adversarial inputs: what would make this function lie?",
  "Property-based testing (Hypothesis, fast-check): specify invariants, not examples",
  "AI proposes the properties; the framework hammers them with thousands of cases"],
 "Example-based tests check what you thought of; property-based tests check what you did not. The "
 "division of labor is lovely: the model is a strong brainstormer of edge cases and invariants, "
 "and the property framework does the exhausting. Even teams that never adopt property-based "
 "testing should steal the adversarial-input prompt."),

("b", "Test Data Generation",
 ["Generate realistic fixtures: users, orders, events — varied shapes, not one golden record",
  "Specify distributions and constraints: '1% nulls, names in five locales, dates spanning DST'",
  "Ask for the nasty cases: malformed rows, encoding traps, duplicate keys",
  "Synthetic only — never paste production data into a prompt (Chapter 8, privacy)",
  "Version the generated fixtures with the tests they serve"],
 "Realistic, varied test data is drudge work that models do happily — and the constraint spec "
 "(nulls, locales, DST boundaries) is where the value hides. Two rules: synthetic data only, "
 "because production data in a prompt is a privacy incident; and fixtures get versioned like "
 "code, because tests depend on them."),

("b", "Tests Are the Spec",
 ["A precise test is the clearest requirement statement an engineer can write",
  "Acceptance criteria (Given/When/Then) translate directly into test cases",
  "With tests in place, 'did the AI do what I meant?' becomes a runnable question",
  "The spec outlives the code: implementations change, the spec keeps them honest",
  "This is the bridge to TDD — tests first, implementation second"],
 "The conceptual hinge of the whole chapter. If a test is an executable spec, then directing an "
 "AI stops being prompt-crafting and becomes spec-writing — a skill engineers already have. "
 "Trace the path backward: acceptance criteria from the requirements section become these tests. "
 "Everything after this slide is a consequence of it."),

("b", "TDD Refresher: Red-Green-Refactor",
 ["Red — write one failing test that states the next small behavior",
  "Green — write the simplest code that passes it",
  "Refactor — clean up with tests green, in small safe steps",
  "The cycle forces small units, clear contracts, and continuous verification",
  "It is also the perfect discipline for directing an AI pair — next slides"],
 "Thirty-second refresher for the rusty: red, green, refactor, with the emphasis on SMALL steps. "
 "Then make the turn: notice that red-green-refactor's demands (small units, clear contracts, "
 "constant verification) are exactly the properties that make AI-generated code safe. TDD was "
 "always about controlling complexity; now it controls the AI too."),

("b", "TDD with an AI Pair: You Write the Tests",
 ["The strongest pattern: human writes failing tests from the spec; AI implements until green",
  "You keep authority over what correct means — the AI only fills in how",
  "Prompt: 'Here are the failing tests — make them pass, change nothing else'",
  "The red bar proves the AI did not game the tests; the green bar proves the implementation",
  "Lab 5.1 drills exactly this loop"],
 "Pattern one, the recommended default: the human owns the spec as tests, the AI owns the "
 "implementation, and the harness arbitrates. Walk the loop once on a whiteboard — red, API call, "
 "green, refactor — because Lab 5.1 is exactly this. The red-bar-first detail matters: a test "
 "that never failed may not be testing anything."),

("b", "TDD with an AI Pair: AI Drafts the Tests",
 ["Faster variant: AI drafts tests from your spec — you review, correct, and extend them",
  "Review test intent first: is this the behavior the business asked for?",
  "Beware tests that mirror the implementation — AI tests can rubber-stamp AI code",
  "Keep at least one human-authored test per feature as an anchor",
  "Then the same loop: AI implements until the corrected tests pass"],
 "Pattern two trades control for speed, so the review discipline tightens: you are auditing the "
 "spec itself now. Name the failure mode plainly — a model grading its own homework writes tests "
 "that mirror its implementation — and give the countermeasure: at least one human-authored "
 "anchor test per feature. Both patterns end at the same red-green loop."),

("b", "Evals Are TDD for AI Behavior",
 ["The same idea one level up: define expected behavior, run it, watch it fail, then fix",
  "Evals pin down prompt and model behavior exactly as unit tests pin down functions",
  "Red-green maps directly: failing eval, improve prompt or model, passing eval",
  "Chapter 6 builds the full practice: golden sets, rubrics, LLM-as-judge, regression gates",
  "If your product ships AI behavior, evals are its regression suite"],
 "The bridge to Chapter 6 — do not teach evals here, just name the isomorphism. Students who "
 "internalized 'tests are the spec' should feel 'evals are TDD for AI behavior' click "
 "immediately. One sentence on why it matters: any feature whose behavior comes from a model "
 "needs evals the same way any function needs tests."),

("b", "Mutation Testing and Coverage with AI Help",
 ["Coverage tells you what ran; mutation testing tells you what is actually checked",
  "AI closes the gap: 'Line 42 survived three mutants — write the test that kills them'",
  "Ask the model to play saboteur: introduce a subtle bug, see if your suite catches it",
  "Use AI to triage mutation reports — kill the mutants that matter, skip the equivalents",
  "The goal: a suite that fails when behavior changes, not one that is merely green"],
 "Coverage is the most gamed metric in engineering — mutation testing is the audit. Two AI uses: "
 "generating the killer tests for surviving mutants, and triaging the report (mutation tools are "
 "noisy, and models are decent at spotting equivalent mutants). The saboteur game is a fun live "
 "demo and doubles as a suite health check."),

# ===== CODE REVIEW (3) =====
("b", "AI-Assisted Code Review Workflows",
 ["First-pass review by AI: style, obvious bugs, missing tests — before human eyes",
  "AI as review partner: paste the diff, ask 'what would you flag, and why?'",
  "Summarize large PRs for human reviewers: intent, risk areas, suggested focus",
  "Keep humans on the merge decision — accountability does not delegate",
  "Watch review load: generated code makes PRs bigger, so set size norms (adoption section)"],
 "The pragmatic workflow: machine first pass, human judgment pass. The PR-summary use is "
 "underrated — it raises review quality on big diffs because reviewers know where to look. Warn "
 "about the new bottleneck: when code is cheap to generate, review load explodes, so PR-size "
 "norms become a team-level control."),

("b", "Review Prompts and Checklists",
 ["Make review prompts team assets: version them, tune them, share them (Chapter 4 discipline)",
  "A solid checklist: correctness, error handling, security, tests, readability",
  "Point the model at your standards: paste the style guide or ADR it must honor",
  "Ask for severity-ranked findings with line references — not a wall of prose",
  "Chapter 8 adds the security pass: injection, secrets, unsafe dependencies"],
 "Review quality from an LLM is mostly prompt quality, so the prompts get engineered like "
 "everything else — versioned, reviewed, improved. Two concrete upgrades: ground the review in "
 "your actual standards documents, and demand severity-ranked, line-referenced output so findings "
 "are actionable. Security review gets its own chapter later."),

("b", "Hallucination Traps in Generated Code",
 ["Invented APIs: methods and parameters that do not exist in your library version",
  "Subtly wrong logic: off-by-one, inverted comparisons, swallowed exceptions",
  "Confident nonsense: comments describing code that is not there",
  "Outdated patterns: deprecated APIs and old idioms from training data",
  "Defense in depth: run it, test it, review it — and pin dependency versions"],
 "Give the rogues' gallery with a real example of each if you can — especially the invented-API "
 "trap, which every student will hit. The through-line is that these failures look like fine "
 "code, which is why the testing section preceded this one: traps are what the harness catches. "
 "Pinning dependency versions removes a whole class of version-mismatch hallucinations."),

# ===== DOCS, DEPLOYMENT & AUTOMATION (6) =====
("b", "Documentation Generation and Upkeep",
 ["Generate docstrings and README sections from the code — then edit for truth",
  "Draft tutorials from working examples, keeping the code real and runnable",
  "Docs-as-code: regenerate on change in CI so docs track the implementation",
  "Ask for the reader's view: 'What is confusing about this README?'",
  "Stale docs are worse than none — automate the upkeep or delete them"],
 "Docs are the first place teams see ROI and the first place drift appears. The honest framing: "
 "generation is cheap, truth is expensive — so the win is not writing docs once but regenerating "
 "them continuously in CI. If a doc cannot be regenerated or verified, question whether it should "
 "exist."),

("b", "Repo Q&A with RAG",
 ["Question-answering over your own codebase, docs, tickets, and ADRs",
  "RAG retrieves the relevant files into the prompt — the grounding pattern from Chapter 4",
  "Killer use case: onboarding — 'where is retry logic configured?' answered with citations",
  "Keep the index fresh: re-embed on merge, or answers drift from the code",
  "Citations are mandatory — an answer without a source is a rumor"],
 "One slide only because Chapter 4 owns RAG mechanics — here it is applied to the repo. Sell the "
 "onboarding use case; every engineer remembers their first week of unanswerable questions. The "
 "two operational rules: fresh index (re-embed on merge) and mandatory citations, otherwise the "
 "answers quietly become folklore."),

("b", "LLM Steps in CI/CD Pipelines",
 ["LLMs as pipeline steps: PR summaries, test-failure triage, changelog drafts",
  "Call the OpenAI API from CI scripts — the Chapter 4 patterns, run headless",
  "Structure outputs with JSON Schema so downstream steps can branch on the result",
  "Gate, don't block: AI findings land as comments and warnings; humans own merge and release",
  "Control cost and flakiness: cache results, pin models, set token budgets"],
 "Everything from Chapter 4, now running unattended — which is exactly why the engineering bar "
 "goes up: structured outputs for machine consumers, pinned models for repeatability, budgets for "
 "cost. The gate-don't-block principle keeps a flaky model step from holding the pipeline "
 "hostage while still putting its findings in front of humans."),

("b", "Release Notes and Changelog Automation",
 ["Draft release notes from merged PRs and commit ranges, grouped by user impact",
  "Ask for audiences: user-facing notes vs. operator-facing upgrade notes",
  "Flag breaking changes and migrations explicitly — the model must hunt for them",
  "Human review before publish: tone is cheap, accuracy is not",
  "Close the loop: link notes back to tickets and PRs for traceability"],
 "A toil task everyone hates and models do well — with one sharp edge: breaking-change detection "
 "must be demanded explicitly or the model will smooth it over. Human review stays, but it shifts "
 "from writing to fact-checking. The traceability bullet (notes linked to PRs and tickets) makes "
 "audits and rollbacks far easier."),

("b", "Incident Triage and Log Summarization",
 ["Paste the alert and logs; get a timeline, hypotheses, and next diagnostic steps",
  "Strong at pattern-matching across noisy stack traces and metric descriptions",
  "Draft the comms: status updates and a postmortem skeleton from the timeline",
  "Keep humans on the commands — AI suggests, on-call executes (agents, Chapter 7)",
  "Scrub secrets and customer data before anything enters a prompt (Chapter 8)"],
 "At 3 a.m. the model's value is triage speed: timelines, hypotheses, and drafted comms while the "
 "human thinks. Two hard boundaries: the model never runs the remediation commands (that is "
 "Chapter 7's gated-agent territory), and logs get scrubbed first because they are full of "
 "secrets and personal data. The postmortem skeleton alone usually pays for the workflow."),

("b", "Automating Dev-Workflow Toil with Agents",
 ["Beyond one-shot calls: agents chain steps — triage ticket, find code, propose patch",
  "Built on function calling, tool allow-lists, and loop caps — Chapter 7's core mechanics",
  "Good first targets: dependency updates, flaky-test quarantine, issue labeling",
  "Irreversible actions stay behind human approval gates",
  "Chapter 7 builds these hands-on; this slide is the map"],
 "A forward pointer, not a lesson: show that the chapter's one-shot API patterns compose into "
 "agents, and name safe first targets where a bad outcome is reversible. The gate rule carries "
 "over from incident response — agents propose, humans approve anything irreversible. Chapter 7 "
 "teaches the loop, the tools, and the guardrails properly."),

# ===== ADOPTION (3) =====
("b", "Measuring Impact",
 ["Baseline first: cycle time, review load, defect escape rate — before and after",
  "DORA-style metrics adapt well: lead time, change failure rate, deployment frequency",
  "Measure rework, not just speed — AI code that bounces in review is not faster",
  "Collect the qualitative signal too: developer surveys on friction and flow",
  "Beware vanity metrics: lines of code generated is not productivity"],
 "Adoption without measurement becomes folklore in both directions — hype or backlash. Push for "
 "before/after baselines and outcome metrics, and call out rework explicitly: generated code that "
 "bounces in review or escapes as defects is negative productivity. The vanity-metric warning is "
 "the line to remember."),

("b", "Quality Gates and Team Norms",
 ["Non-negotiables: tests pass, human review, no unowned code — AI-written or not",
  "Norm the workflow: when to use autocomplete vs. chat vs. agent; disclose AI use in PRs",
  "Prompts and review checklists are team assets — versioned and improved together",
  "The definition of done includes verified AI output, not just generated AI output",
  "Start on low-risk paths, prove the gates, then expand"],
 "The chapter's practices only work as team norms, not individual habits. The gates did not "
 "change when AI arrived — tests, review, ownership — but they must now be stated to include "
 "AI-generated work explicitly. Recommend the crawl-walk-run rollout: low-risk paths first, "
 "expand once the gates prove they hold."),

("b", "When Not to Use Generative AI",
 ["Safety-critical and regulated paths where every line needs traced provenance",
  "Secrets, credentials, and customer data — never into a prompt (Chapter 8)",
  "Licensed-code risk: generated code may echo training data, so run license scanning",
  "Novel algorithms whose correctness is subtle and no test oracle exists yet",
  "When the cost of being wrong exceeds the cost of writing it yourself"],
 "Credibility slide — a chapter that never says 'don't' is marketing, not engineering. Work down "
 "the list with examples; the last bullet is the general rule that generates the others. Note "
 "that 'no test oracle exists' is the interesting case: verification is the whole game, so where "
 "you cannot verify, you should not generate."),

# ===== LAB (1) =====
("lab", "Lab 5.1: TDD with an AI Pair",
 ["Follow the detailed instructions in the Lab 5.1 notebook on your VM",
  "Write failing tests from a spec — you define what correct means",
  "Use the OpenAI API to implement until the tests go green",
  "Refactor with AI assistance while keeping the suite green"],
 "The chapter's capstone lab, and the one to protect time for. Students run the recommended "
 "pattern end to end: human-written failing tests, an OpenAI API loop that implements until "
 "green, then an AI-assisted refactor with the suite as referee. The VM has the API key "
 "pre-injected as in Chapter 4; circulate during the first red-green cycle, because students who "
 "skip the red bar miss the point."),
]


# --- formatting-preserving text rewrite helpers (same as build_ch04/ch07) ------
def _set_para_text(p_el, text):
    runs = p_el.findall(qn("a:r"))
    if not runs:
        r = p_el.makeelement(qn("a:r"), {})
        r.append(p_el.makeelement(qn("a:t"), {}))
        p_el.append(r)
        runs = [r]
    first = runs[0]
    for extra in runs[1:]:
        p_el.remove(extra)
    for br in p_el.findall(qn("a:br")):
        p_el.remove(br)
    t = first.find(qn("a:t"))
    if t is None:
        t = first.makeelement(qn("a:t"), {})
        first.append(t)
    t.text = text


def rewrite_bullets(shape, lines):
    txBody = shape.text_frame._txBody
    paras = txBody.findall(qn("a:p"))
    for i, line in enumerate(lines):
        if i < len(paras):
            _set_para_text(paras[i], line)
        else:
            newp = copy.deepcopy(paras[-1])
            txBody.append(newp)
            paras.append(newp)
            _set_para_text(newp, line)
    for p in paras[len(lines):]:
        txBody.remove(p)


# --- authored-slide builders ---------------------------------------------------
def author_bullets(prs, title, bullets, notes):
    lay = sk.pick_layout(prs, "Content with Header. Full Page")
    slide = prs.slides.add_slide(lay)
    sk.set_title(slide, title)
    body = pt._body_placeholder(slide)
    body.height = Inches(4.4)          # layout box is stubby for 4-6 bullets
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
    sk.set_notes(slide, notes)
    return slide


def _fill_column(ph, heading, bullets):
    ph.height = Inches(4.4)            # layout boxes are stubby (1.7")
    tf = ph.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = heading
    if p0.runs:
        p0.runs[0].font.bold = True
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b


def author_two_col(prs, title, payload, notes):
    (l_head, l_bullets), (r_head, r_bullets) = payload
    lay = sk.pick_layout(prs, "Two Column Full Page")
    slide = prs.slides.add_slide(lay)
    sk.set_title(slide, title)
    cols = {}
    for ph in slide.placeholders:
        if ph.placeholder_format.idx in (1, 13) and ph.has_text_frame:
            cols[ph.placeholder_format.idx] = ph
    assert 1 in cols and 13 in cols, "two-column placeholders not found"
    _fill_column(cols[1], l_head, l_bullets)
    _fill_column(cols[13], r_head, r_bullets)
    sk.set_notes(slide, notes)
    return slide


def author_table(prs, title, payload, notes):
    header, rows = payload
    lay = sk.pick_layout(prs, "Content with Header. Full Page")
    slide = sk.add_table_slide(prs, title, header, rows, notes=notes, layout=lay)
    # drop the empty OBJECT placeholder so its prompt text doesn't sit under the table
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx != 0 and ph.has_text_frame \
                and not ph.text_frame.text.strip():
            ph._element.getparent().remove(ph._element)
    return slide


def author_lab(prs, title, bullets, notes):
    lay = sk.pick_layout(prs, "Exercise Reference Slide with Typing Hands")
    assert lay is not None, "exercise layout not found in carrier template"
    slide = prs.slides.add_slide(lay)
    sk.set_title(slide, title)
    body = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1 and ph.has_text_frame:
            body = ph
            break
    if body is None:
        body = pt._body_placeholder(slide)
    body.height = Inches(4.5)          # the layout's content placeholder is stubby
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
    sk.set_notes(slide, notes)
    return slide


_BUILDERS = {"b": author_bullets, "tc": author_two_col,
             "tbl": author_table, "lab": author_lab}


def main():
    # 0. fresh carrier -> idempotent
    shutil.copy(CARRIER, OUT)
    prs = sk.open_prs(OUT)
    n_carrier = len(prs.slides._sldIdLst)
    assert n_carrier == 63, f"carrier expected 63 slides, got {n_carrier}"

    # 1. author all 45 body slides (appended at end, in final teaching order)
    for kind, title, payload, notes in SLIDES:
        _BUILDERS[kind](prs, title, payload, notes)
    n_new = len(SLIDES)

    # 2. rewrite the 5 kept carrier slides in place (0-based: 0,1,2,61,62)
    s1 = prs.slides[0]
    _set_para_text(s1.shapes.title.text_frame._txBody.findall(qn("a:p"))[0], TITLE)
    for sh in s1.shapes:                              # subtitle: Chapter 4 -> Chapter 5
        if sh.has_text_frame and "Chapter 4" in sh.text_frame.text:
            _set_para_text(sh.text_frame._txBody.findall(qn("a:p"))[0], SUBTITLE)
    sk.set_notes(s1, TITLE_NOTES)
    rewrite_bullets(pt._body_placeholder(prs.slides[1]), OBJECTIVES)   # Objectives
    sk.set_notes(prs.slides[1], OBJECTIVES_NOTES)
    rewrite_bullets(pt._body_placeholder(prs.slides[2]), AGENDA)       # Agenda
    sk.set_notes(prs.slides[2], AGENDA_NOTES)
    rewrite_bullets(pt._body_placeholder(prs.slides[61]), SUMMARY)     # Summary
    sk.set_notes(prs.slides[61], SUMMARY_NOTES)
    rewrite_bullets(pt._body_placeholder(prs.slides[62]), OBJECTIVES)  # objectives recap
    sk.set_notes(prs.slides[62], RECAP_NOTES)

    # 3. arrange: openers -> authored body (already in order) -> closers;
    #    the 58 old content slides are dropped
    order = [0, 1, 2] + list(range(n_carrier, n_carrier + n_new)) + [61, 62]
    pt.arrange(prs, order)
    pt.save(prs, OUT)

    # 4. verify: reopen, count + titles, blanks, notes coverage, partnames, zip
    prs2 = Presentation(OUT)
    n = len(prs2.slides)
    print(f"\n{OUT.name}: {n} slides (carrier-kept=5, authored={n_new})")
    empty, no_notes = [], []
    for i, s in enumerate(prs2.slides, 1):
        t = sk.slide_title(s)
        if not t.strip():
            empty.append(i)
        if not s.has_notes_slide or not s.notes_slide.notes_text_frame.text.strip():
            no_notes.append(i)
        print(f"{i:>3} | {t}")
    assert not empty, f"empty titles at {empty}"
    assert not no_notes, f"missing speaker notes at {no_notes}"
    assert 46 <= n <= 52, f"slide count {n} outside 46-52"
    with zipfile.ZipFile(OUT) as z:
        dupes = [p for p, c in Counter(z.namelist()).items() if c > 1]
        assert not dupes, f"duplicate zip partnames: {dupes}"
        assert z.testzip() is None, "zip integrity check failed"
    print("\nOK: no empty titles, notes on every slide, "
          "no duplicate partnames, zip integrity verified")


if __name__ == "__main__":
    main()
