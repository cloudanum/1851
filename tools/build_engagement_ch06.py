#!/usr/bin/env python3
"""build_engagement_ch06.py — ENGAGEMENT pass on 1851-Ch06 (A.2, 44 -> 46 slides).

Restores the A.2 deck fresh from decks/_bak1 every run (idempotent), then:
  1. inserts 'Do Now: Grade This Answer' at the start of the Evaluation Methods
     section, right before 'Exact Match and Programmatic Assertions' (layout
     'Do Now with Typing Hands'),
  2. inserts 'Do Now: Flaky or Broken?' in the evals-in-CI block, right after
     'Eval Suites in CI' (layout 'Discussion').

Lab 6.1 already has its exercise slide — no lab-slide work here.

Run:  /Users/iahmad/Creator/Courses_and_conferences/LT/.venv-courseware/bin/python tools/build_engagement_ch06.py
"""
from __future__ import annotations
import shutil
import sys
import zipfile
from collections import Counter
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import slidekit as sk  # noqa: F401  (import activates the partname-collision fix)
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent
BAK1 = ROOT / "decks" / "_bak1"
OUT = ROOT / "decks" / "1851-Ch06.pptx"

# --- new slide content -------------------------------------------------------
DONOW_GRADE = ("Do Now: Grade This Answer", "Do Now with Typing Hands",
    ["User question: \"A junior dev suggests hashing our passwords with MD5. What do you recommend?\"",
     "LLM answer: \"MD5 is a good choice: it is fast and produces fixed-length hashes, which keeps "
     "lookups efficient. Add a salt — a random value per password — so identical passwords hash "
     "differently. For extra security, apply MD5 twice.\"",
     "Mini-rubric:  5 = correct, safe, complete   3 = partly right, missing something key   "
     "1 = wrong or dangerous advice",
     "Score it 1–5 and list every flaw you find. Time box: 7 minutes."],
    "Intended score: 1 (a generous 2). Two subtle flaws: (1) it endorses MD5 for password hashing — its "
    "speed is exactly why it is wrong here; the right advice is a slow password hash such as bcrypt, "
    "scrypt, or Argon2; (2) 'apply MD5 twice' is cargo-cult security, not a real defense. The salt "
    "sentence is correct, which is what makes the answer dangerous: fluent prose smuggling in harmful "
    "advice. Debrief: this is why evals need domain rubrics and human spot-checks — a fluency check "
    "would have passed it.")

DONOW_FLAKY = ("Do Now: Flaky or Broken?", "Discussion",
    ["Three failures from last week's CI eval runs. Real regression or flaky eval?",
     "A. Faithfulness fell 0.91 → 0.62 right after a prompt edit — and stays there on every re-run.",
     "B. The gate failed with 'judge call timed out after 30 s'; the same commit passed on re-run, "
     "code unchanged.",
     "C. The judge scored one answer 2/5, then 5/5 on two re-runs — same input, same judge model; "
     "the judge's temperature was never pinned.",
     "Decide each, and say what you would check next. Time box: 6 minutes in pairs."],
    "Answers: A is a real regression — a reproducible drop tied to a specific change; bisect the prompt "
    "diff. B is flaky infrastructure — a judge timeout says nothing about the system; add retries and "
    "alert on judge health, not on the score. C is a flaky eval, not a flaky system — an unpinned judge "
    "makes the measurement itself nondeterministic; pin judge model and temperature and version judge "
    "configs like datasets. Debrief: trust the signal only when the harness is deterministic — that is "
    "why this chapter insists on pinning, versioning, and canarying.")


# --- helpers ------------------------------------------------------------------
def pick_layout_exact(prs, name):
    """Exact layout-name match (pick_layout's substring match would return the
    'Numbered …' variant first)."""
    for master in prs.slide_masters:
        for lay in master.slide_layouts:
            if (lay.name or "") == name:
                return lay
    raise AssertionError(f"layout not found: {name!r}")


def add_donow_slide(prs, title, layout_name, bullets, notes):
    """Add a Do Now / Discussion slide on the named template layout (graphic on
    the left, text column on the right). Returns the new slide."""
    lay = pick_layout_exact(prs, layout_name)
    slide = prs.slides.add_slide(lay)
    sk.set_title(slide, title)
    body = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1 and ph.has_text_frame:
            body = ph
            break
    assert body is not None, f"content placeholder not found on {layout_name!r}"
    # the layout's content placeholder is stubby; give it room for 5+ lines
    body.height = Inches(4.5)
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
    sk.set_notes(slide, notes)
    return slide


def find_slide(prs, title_sub):
    """0-based index of the unique slide whose title contains title_sub."""
    hits = [i for i, s in enumerate(prs.slides) if title_sub in sk.slide_title(s)]
    assert len(hits) == 1, f"anchor {title_sub!r}: expected 1 hit, got {hits}"
    return hits[0]


def main():
    # 0. fresh restore from _bak1 -> idempotent
    shutil.copy(BAK1 / OUT.name, OUT)
    prs = sk.open_prs(OUT)
    assert len(prs.slides) == 44, f"expected 44-slide A.2 deck, got {len(prs.slides)}"

    # 1+2. append each new slide and immediately move it into place — descending
    # target index, so the second move cannot disturb the first placement.
    add_donow_slide(prs, DONOW_FLAKY[0], DONOW_FLAKY[1], DONOW_FLAKY[2], DONOW_FLAKY[3])
    at_flaky = find_slide(prs, "Eval Suites in CI") + 1                          # after it
    sk.move_slide(prs, len(prs.slides) - 1, at_flaky)

    add_donow_slide(prs, DONOW_GRADE[0], DONOW_GRADE[1], DONOW_GRADE[2], DONOW_GRADE[3])
    at_grade = find_slide(prs, "Exact Match and Programmatic Assertions")        # before it
    sk.move_slide(prs, len(prs.slides) - 1, at_grade)

    print(f"inserted at 1-based: {at_grade+1} (Grade This Answer), "
          f"{at_flaky+1} (Flaky or Broken?)")

    sk.save(prs, OUT)

    # 3. verify: reopen, count, titles, blanks, duplicate partnames, zip integrity
    prs2 = Presentation(OUT)
    n = len(prs2.slides)
    print(f"\n{OUT.name}: {n} slides")
    empty = []
    for i, s in enumerate(prs2.slides, 1):
        t = sk.slide_title(s)
        if not t.strip():
            empty.append(i)
        print(f"{i:>3} | {t}")
    assert not empty, f"empty titles at {empty}"
    assert n == 46, f"slide count {n} != 46"
    titles = [sk.slide_title(s) for s in prs2.slides]
    # positions are relative to neighbors, so verify the adjacency contract
    assert titles.index(DONOW_GRADE[0]) == \
        titles.index("Exact Match and Programmatic Assertions") - 1
    assert titles.index(DONOW_FLAKY[0]) == titles.index("Eval Suites in CI") + 1
    with zipfile.ZipFile(OUT) as z:
        dupes = [p for p, c in Counter(z.namelist()).items() if c > 1]
        assert not dupes, f"duplicate zip partnames: {dupes}"
        assert z.testzip() is None, "zip integrity check failed"
    print("\nOK: 46 slides, no empty titles, no duplicate partnames, zip integrity verified")


if __name__ == "__main__":
    main()
