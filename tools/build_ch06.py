#!/usr/bin/env python3
"""build_ch06.py — assemble the A.2 'Evaluating Generative AI Systems' chapter.

Carrier = _bak0/1851-Ch05.pptx (A.1 'Model Evaluation Metrics', restored fresh
every run -> idempotent). The A.1 metric content slides all moved into Ch03, so
only the openers (title/objectives/agenda) and closers (summary/recap) are kept
and rewritten; the body is authored fresh, plus six high-value clones:

  16    = 2016-agentic-security/2016-slides/enriched/Ch06-Final.pptx
          s23 TEVV, s24 Eval-Driven Development, s25 NIST ARIA, s26 KPIs
  58_4  = 1258.../decks/1258-Ch04-ModernGenAI-PromptEng.pptx  s54 rubric scoring
  58_7  = 1258.../decks/1258-Ch07-Building-with-LLMs.pptx     s33 Ragas

Clones get branding/audience patches (RMF/government/security framing softened
to organizational quality framing, 2016/1258 chapter and lab refs fixed) via
PATCHES / NOTES_PATCHES / RETITLES. Authored slides sit on 'Content with
Header. Full Page' (or 'Two Column Full Page' for the tests-vs-evals contrast),
each with speaker notes. Lab 6.1 uses the typing-hands exercise layout.

Run:  /Users/iahmad/Creator/Courses_and_conferences/LT/.venv-courseware/bin/python tools/build_ch06.py
"""
from __future__ import annotations
import copy
import shutil
import sys
import zipfile
from collections import Counter
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import slidekit as sk  # noqa: F401  (import activates the partname-collision fix)
import pptx_tools as pt
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent
BAK = ROOT / "decks" / "_bak0"
OUT = ROOT / "decks" / "1851-Ch06.pptx"
CARRIER = BAK / "1851-Ch05.pptx"
SRC = {
    "16": Path("/Users/iahmad/Creator/Courses_and_conferences/LT/courses/"
               "2016-agentic-security/2016-slides/enriched/Ch06-Final.pptx"),
    "58_4": Path("/Users/iahmad/Creator/Courses_and_conferences/LT/courses/"
                 "1258-applied-ai-for-government-it-professionals/"
                 "1258a4-author-input/decks/1258-Ch04-ModernGenAI-PromptEng.pptx"),
    "58_7": Path("/Users/iahmad/Creator/Courses_and_conferences/LT/courses/"
                 "1258-applied-ai-for-government-it-professionals/"
                 "1258a4-author-input/decks/1258-Ch07-Building-with-LLMs.pptx"),
}

# --- clone lists (1-based, against source decks) ------------------------------
CLONES = {
    "16": [23, 24, 25, 26],     # TEVV / EDD / NIST ARIA / program KPIs
    "58_4": [54],               # rubric scoring 1-5 + LLM-as-judge
    "58_7": [33],               # Ragas RAG evaluation
}

# --- retitles on cloned slides ------------------------------------------------
RETITLES = {
    ("16", 24): "Eval-Driven Development: Regression Tests for GenAI Quality",
    ("16", 26): "KPIs for an AI Quality Program",
}

# --- body/table patches on cloned slides (exact substrings, asserted) ---------
PATCHES = {
    ("16", 23): [
        ("TEVV is how the RMF Measure function gets evidence",
         "TEVV is how a serious AI program gets its evaluation evidence"),
        ("Test — exercise the agent against benign and adversarial inputs",
         "Test — exercise the system against benign and adversarial inputs"),
    ],
    ("16", 24): [
        ("Build a golden dataset of adversarial cases — every jailbreak, injection, and "
         "unsafe output from Ch3–Ch4 becomes a permanent test case",
         "Build a golden dataset of real cases — every bug, failure, and edge case you "
         "have ever seen becomes a permanent test case"),
        ("LLM-as-judge scores each output against a rubric (Did it refuse? Leak data? "
         "Follow policy?)",
         "LLM-as-judge scores each output against a rubric (Correct? Grounded? Followed "
         "instructions?)"),
        ("gate releases on the robustness score (refusal rate, injection-block rate)",
         "gate releases on the suite score (quality, groundedness, safety)"),
        ("Drift detection: an update that lowers the block rate fails the build",
         "Drift detection: an update that lowers a score fails the build"),
        ("Security caveat: the judge is itself a probabilistic, injectable component — "
         "calibrate it, and never let it be the only gate (pair it with the deterministic "
         "shell)",
         "Caveat: the judge is itself a probabilistic model — calibrate it against human "
         "review, and never let it be the only gate (pair it with deterministic checks)"),
    ],
    ("16", 25): [
        ("Operationalizes the RMF 'Measure' function at scale",
         "Operationalizes continuous, independent evaluation at scale"),
        ("Findings feed back into Manage — fixes, guardrails, and gating",
         "Findings feed back into development — fixes, guardrails, and gating"),
    ],
    ("16", 26): [
        ("% agents with risk tier + named owner",
         "% AI features with eval suite + named owner"),
        ("% actions with signed audit event",
         "% releases with recorded eval results"),
        ("Guardrail efficacy", "Regression catch rate"),
        ("% red-team attempts blocked", "% known regressions caught by the suite"),
        ("MTTR (AI incidents)", "MTTR (AI regressions)"),
        ("Time to contain an agent incident", "Time to detect + revert a bad change"),
        ("% high-risk actions human-approved", "% high-impact changes human-reviewed"),
    ],
}

# --- speaker-notes patches on cloned slides (substring replace, warn on miss) -
NOTES_PATCHES = {
    ("16", 23): [
        ("TEVV is how the Measure function actually gets evidence",
         "TEVV is how a serious evaluation program actually gets its evidence"),
        ("Test means exercising the agent against both benign and adversarial inputs",
         "Test means exercising the system against both benign and adversarial inputs"),
    ],
    ("16", 24): [
        ("test-driven development for adversarial robustness",
         "test-driven development for GenAI quality"),
        ("every jailbreak, injection, and unsafe output the class produced back in "
         "Chapters 3 and 4 becomes a permanent regression test",
         "every failure you have ever seen in review or production becomes a permanent "
         "regression test"),
        ("did it refuse, did it leak data, did it follow policy?",
         "was it correct, was it grounded, did it follow instructions?"),
        ("gate releases on the robustness score - refusal rate, injection-block rate",
         "gate releases on the suite score - quality, groundedness, safety"),
        ("an update that lowers the block rate fails the build, just like a broken unit "
         "test",
         "an update that lowers a score fails the build, just like a broken unit test"),
        ("Now stress the security caveat hard, because it is the chapter's signature "
         "insight: the judge is itself a probabilistic, injectable component.",
         "Now stress this caveat hard, because it is the chapter's signature insight: "
         "the judge is itself a probabilistic model."),
        ("This is the Determinism Sandwich we formalize in 6.4: probabilistic judgment "
         "wrapped in deterministic checks.",
         "Think of it as a determinism sandwich: probabilistic judgment wrapped in "
         "deterministic checks."),
        ("Concrete example: a model upgrade improves helpfulness but quietly weakens "
         "injection resistance; the EDD suite catches the dropped block rate and fails "
         "the build before it ships.",
         "Concrete example: a model upgrade improves helpfulness but quietly weakens "
         "instruction-following; the EDD suite catches the dropped score and fails the "
         "build before it ships."),
    ],
    ("16", 25): [
        ("Frame it as how a serious organization operationalizes the Measure function "
         "beyond a benchmark number, and as the institutional big brother of the EDD "
         "loop on the previous slide.",
         "Frame it as how a serious organization operationalizes evaluation beyond a "
         "benchmark number, and as the institutional big brother of the EDD loop from "
         "this section."),
        ("generate the credible evidence the Measure function demands",
         "generate the credible evidence a serious evaluation program demands"),
        ("they flow straight into the Manage backlog as new guardrails and gates",
         "they flow straight into the engineering backlog as new guardrails and gates"),
        ("Measure feeds Manage; that closed loop is the point.",
         "Evaluation feeds fixes; that closed loop is the point."),
        ("Position red-teaming for security learners as the AI-native penetration test",
         "Position red-teaming as the AI-native penetration test"),
        ("that finding becomes a Manage-side input-validation control",
         "that finding becomes an input-validation control"),
    ],
    ("16", 26): [
        ("Coverage: the percentage of agents that have both a risk tier and a named "
         "owner - target 100 percent, because an ungoverned agent is the one that hurts "
         "you, and you cannot govern what you have not inventoried.",
         "Coverage: the percentage of AI features that have both an eval suite and a "
         "named owner - target 100 percent, because an unmeasured feature is the one "
         "that hurts you, and you cannot improve what you have not instrumented."),
        ("Audit completeness: the percentage of actions that emit a signed audit event - "
         "target 100 percent, and flag that this is the exit criterion the lab enforces.",
         "Audit completeness: the percentage of releases with recorded eval results - "
         "target 100 percent, because a release without recorded results is an "
         "unmeasured release."),
        ("Guardrail efficacy: the percentage of red-team attempts blocked - target "
         "increasing, fed directly by your ARIA-style red-teaming.",
         "Regression catch rate: the percentage of known regressions the suite catches "
         "before release - target increasing, fed by every failure you add to the "
         "golden set."),
        ("MTTR for AI incidents: time to contain an agent incident - target decreasing.",
         "MTTR for AI regressions: time to detect and revert a bad prompt or model "
         "change - target decreasing."),
        ("Oversight rate: the percentage of high-risk actions that received human "
         "approval - target 100 percent.",
         "Oversight rate: the percentage of high-impact changes that received human "
         "review - target 100 percent."),
        ("and they become the dashboard the capstone produces",
         "and they become the dashboard your stakeholders read"),
        ("is this agent governed?", "is this feature evaluated?"),
        ("a falling guardrail-efficacy number after a model update",
         "a falling catch-rate number after a model update"),
        ("as the lab dramatizes", "as every incident postmortem shows"),
        ("Transition: let us put the whole framework to work on one real agent. (~3 min)",
         "Transition: from program KPIs to the tooling that produces the numbers. "
         "(~3 min)"),
    ],
    ("58_4", 54): [
        ("TRANSITION: Evaluation habits, tuned for government work.",
         "TRANSITION: LLM-as-judge done properly — design and calibration, next."),
    ],
    ("58_7", 33): [
        ("Natural extension: score Lab 7.2's grounded vs ungrounded answers.",
         "Natural extension: score your own pipeline's grounded vs ungrounded answers."),
        ("TRANSITION: Where agencies are already pointing this.",
         "TRANSITION: The levels of evaluation — from single prompts to full agent "
         "trajectories."),
    ],
}

# --- rewritten text for kept carrier slides -----------------------------------
OBJECTIVES = [
    "Explain why generative systems need graded evaluation—and how evals differ from classic ML metrics and software tests",
    "Build and version golden datasets: real-usage sampling, synthetic generation, and sizing for statistical confidence",
    "Apply layered evaluation methods: programmatic assertions, rubric scoring, calibrated LLM-as-judge, and human review",
    "Structure evaluation across levels: prompt units, RAG components, integrations, and end-to-end agent trajectories",
    "Practice Eval-Driven Development: write the evals first and gate every prompt and model change in CI",
    "Operate evals with 2026 tooling: OpenAI eval tooling, Ragas, promptfoo, pytest suites, tracing, and scorecards",
]
AGENDA = [
    "Why Generative AI Needs Evaluation",
    "Building Evaluation Datasets",
    "Evaluation Methods: Assertions, Rubrics, and LLM-as-Judge",
    "Levels of Evaluation: From Prompts to Agent Trajectories",
    "Eval-Driven Development and Evaluation in CI",
    "Tooling, Observability, and Operating Evals",
    "Lab 6.1: Write an Eval Suite",
]
SUMMARY = [
    "GenAI has no compiler for correctness—evaluation is the harness, and 'it looked good' does not scale",
    "Golden datasets plus layered scorers—assertions, rubrics, calibrated judges, humans—make quality measurable",
    "Evaluate at every level: prompt units, RAG components, integration seams, and end-to-end agent trajectories",
    "Eval-Driven Development mirrors Chapter 5's TDD: write the evals first; gate every prompt and model change in CI",
    "Judges are models too—calibrate them against human review, watch their biases, and version them like prompts",
    "Next: Chapter 7 applies evals to agents; Chapter 8 adds the security view—evals for jailbreaks and injection",
]
OBJECTIVES_NOTES = (
    "Welcome to Chapter 6. Chapter 4 made prompting an engineering discipline and Chapter 5 "
    "put GenAI into the SDLC with test-driven development; this chapter supplies the missing "
    "half of that story—how you actually measure whether a generative system is any good. "
    "Walk the objectives: they trace the arc from why evals, to datasets, to methods, to "
    "levels, to the Eval-Driven Development loop that mirrors TDD, and finally to operating "
    "evals in production.")
SUMMARY_NOTES = (
    "Close on the twin message of the course: Chapter 5 taught TDD for code, and this chapter "
    "taught its quality twin for generative systems—Eval-Driven Development. The big three to "
    "retain: golden datasets make vibes measurable, layered scorers balance cost against "
    "fidelity, and evals belong in CI where they catch silent regressions. Then point forward: "
    "Chapter 7 evaluates agents and their trajectories, and Chapter 8 extends the same harness "
    "to adversarial cases.")

# --- authored slides ----------------------------------------------------------
# (title, bullets, notes) — 'Content with Header. Full Page' unless noted.
WHY = [
("Why Generative AI Needs Evaluation", [
    "A model is a function without a specification—no compiler checks the correctness of prose",
    "Same prompt, different outputs: non-determinism is a feature (temperature, sampling) and a testing problem",
    "Fluent is not correct: models produce confident, well-formed, wrong answers",
    "Small changes move many outputs at once—one prompt edit or model upgrade can regress a whole feature",
    "Evaluation is the correctness harness engineers bring to this gap",
],
 "Open with the contrast engineers feel in their bones: code has a compiler and tests, a prompt "
 "has neither. Non-determinism means one manual try proves nothing—you need repeated, graded "
 "measurement. Frame evaluation as the discipline that makes GenAI engineering-grade."),
("'It Felt Better' Is Not an Evaluation", [
    "The default industry workflow: tweak the prompt, eyeball three outputs, ship—'vibe checking'",
    "Vibes miss regressions: the fix for case A silently breaks cases B through Z",
    "Silent regressions are the norm—a provider-side model update changes behavior and nobody notices",
    "Field stories: support-bot tone drift; extraction JSON breaking after a 'harmless' prompt edit",
    "You cannot improve what you do not measure—and you cannot keep what you do not re-measure",
],
 "Every engineer in the room has shipped a prompt change on vibes—name the habit without shaming "
 "it. The two failure stories matter: tone drift (a quality regression nobody tested for) and a "
 "broken JSON contract (a structural regression a one-line assertion would have caught). The "
 "takeaway: regressions are silent until you instrument them."),
("What to Evaluate", [
    "Correctness—is the answer right for the input?",
    "Groundedness—is every claim supported by the provided context (no hallucination)?",
    "Instruction-following—format, length, constraints, schema compliance",
    "Tone and style—does it match the product voice and audience?",
    "Safety and policy—refusals, PII handling, disallowed content (Chapter 8 goes deep)",
    "Cost and latency—quality that blows the budget or the SLA is not quality",
],
 "Quality is multi-dimensional, and each dimension gets its own scoring method later in the "
 "chapter. Note that cost and latency are evaluation dimensions too—a perfect answer that takes "
 "40 seconds and a dollar per request fails the product. Defer the safety deep-dive to Chapter 8."),
("Evals vs Classic ML Metrics", [
    "Chapter 3 covered classic metrics: accuracy, precision/recall, F1, RMSE—one number per model, computed on labeled data",
    "They assume one correct label per input; GenAI has many acceptable outputs per prompt",
    "Classic metrics score a model; GenAI evals score a system—prompt + model + retrieval + tools",
    "GenAI evals are graded judgments (rubrics, judges), not just counts over a confusion matrix",
    "Keep both: classifiers and fine-tunes still use Chapter 3 metrics; prompts and pipelines use this chapter's evals",
],
 "This is the one contrast slide—do not reteach Chapter 3. The key differences: many acceptable "
 "outputs instead of one label, a whole system under test instead of a model, and graded judgment "
 "instead of counts. Engineers who own classic ML pipelines keep those metrics; this chapter adds "
 "the generative layer on top."),
# two-column contrast slide inserted here (see TWO_COL below)
("The Evaluation Loop", [
    "Assemble a golden dataset of representative inputs (plus references where they exist)",
    "Run the system on the dataset and capture outputs—with traces",
    "Score outputs with layered methods: assertions, rubrics, LLM-as-judge, human review",
    "Compare against thresholds and the previous baseline—then decide: ship, iterate, or revert",
    "Repeat on every change; this loop is the backbone of the chapter",
],
 "Preview the chapter structure through the loop: the next sections build each stage—datasets, "
 "then methods, then the levels of the system you run it on. Eval-Driven Development, the core "
 "section, is this loop wired into CI."),
]

TWO_COL = (
"Evals vs Software Tests",
["Classic software tests",
 "Deterministic: same input, same result",
 "Binary: pass or fail",
 "Assert exact values and states",
 "Written against a specification"],
["GenAI evals",
 "Stochastic: scores carry confidence intervals",
 "Graded: rubric scores and win rates",
 "Score semantic properties of text",
 "Written against examples plus judgment"],
 "Both columns live in the same toolbox and the same CI pipeline—evals complement tests, they do "
 "not replace them. The mindset shift: a failing eval is a measured drop in quality, not a broken "
 "assertion, so thresholds and baselines replace pass/fail.")

DATASETS = [
("Golden Datasets: Inputs and References", [
    "A golden dataset is the eval's fixed test bed: curated inputs plus expected outputs or references where available",
    "'Expected' need not mean exact text—a reference answer, key facts, or a constraint list all work",
    "Seed it from real usage: support tickets, real queries, sanitized production logs",
    "Every bug found in review or production becomes a new golden case—the dataset only grows",
    "Keep it in Git next to the prompts it tests; it is code",
],
 "The golden dataset is the single highest-leverage artifact of the chapter—everything else "
 "measures against it. Stress that 'reference' is flexible: for extraction you may have exact "
 "answers, for open-ended generation a list of must-include facts is enough for a judge to score "
 "against."),
("Sampling from Real Usage", [
    "Stratify across the traffic your system actually sees: intents, user segments, input lengths, difficulty",
    "Include the tail—rare intents and adversarial inputs are where systems fail",
    "Sample roughly in proportion, deliberately over-sampling high-risk and high-value slices",
    "Refresh on a cadence: usage drifts, and a stale sample measures last quarter's product",
    "Mind privacy: sanitize PII before production data becomes test data",
],
 "A golden set that only covers the happy path certifies the happy path. Stratification is what "
 "makes a suite's average score meaningful—otherwise the easy 80 percent of traffic drowns the "
 "slice that generates incidents. PII sanitization is non-negotiable before production data moves "
 "into a test fixture."),
("Synthetic Test Data with LLMs—and Its Risks", [
    "LLMs generate candidate cases at scale: paraphrases, edge cases, variations of real inputs",
    "Great for cold start and for expanding coverage of a thin slice",
    "Risk: the generator's blind spots become the dataset's blind spots—synthetic data is homogeneous",
    "Risk: circularity—testing a model on data written by the same model family flatters it",
    "Always anchor with real cases; treat synthetic data as augmentation, never the whole set",
],
 "Synthetic generation is how small teams get to a few hundred cases quickly, and Chapter 4's "
 "prompting skills apply directly. But name the two risks plainly: synthetic cases cluster around "
 "what the generator already handles, and a model judged on its own family's words looks better "
 "than it is. Real cases stay the anchor."),
("Dataset Versioning and Hygiene", [
    "Version the dataset like code: every eval run records the dataset version it used",
    "Never edit cases silently—a changed case invalidates every historical comparison",
    "Split dev vs holdout: iterate prompts against dev cases, report numbers on the holdout",
    "Watch for leakage: golden cases pasted into prompts or fine-tunes stop measuring anything",
    "Deduplicate, review, and re-balance periodically—datasets rot like any codebase",
],
 "Dataset versioning is what makes 'the score went up' meaningful: same data, same scorers, new "
 "system. The dev/holdout split prevents the classic failure of iterating a prompt until it "
 "overfits the visible cases—exactly the overfitting story from Chapter 3, one level up."),
("How Many Cases Do You Need?", [
    "Ten cases is a demo; a hundred starts to be a measurement",
    "Intuition: a score's margin of error shrinks like 1/√n—100 cases ≈ ±10%, 400 cases ≈ ±5%",
    "Size to the decision: detecting a 2-point regression needs more cases than detecting a 20-point one",
    "Small slices need their own counts—a 95% average can hide a failing minority slice",
    "When in doubt: more diverse real cases beat more synthetic ones",
],
 "Keep the statistics intuitive—the point is that a CI gate must see a real regression through "
 "sampling noise, and noisy judges add variance on top of model variance. If the team remembers "
 "one number: a hundred cases per important slice is where measurement begins."),
]

METHODS = [
("Exact Match and Programmatic Assertions", [
    "Unit-test style: assert the output contains 'refund policy', equals the expected ID, starts with 'Summary:'",
    "Perfect for extraction, classification-style prompts, and closed-form answers",
    "Cheap, deterministic, instant—run thousands of cases per CI build",
    "Brittle by nature: paraphrases fail exact match even when correct—scope it to truly closed outputs",
    "Write assertions like pytest tests: named, independent, one property each",
],
 "This is the layer engineers already know how to write, and it carries more of an eval suite "
 "than newcomers expect. The discipline is scoping: exact match only where the output is truly "
 "closed-form—otherwise you manufacture flaky tests and teach the team to ignore red builds."),
("Regex, Schema, and JSON-Validity Checks", [
    "Structure is verifiable: JSON parses, matches the schema, required fields present, types correct",
    "Regex covers contracts: ticket IDs, date formats, 'no markdown', maximum bullet count",
    "Pairs with Structured Outputs and function calling (Chapter 4)—schema checks are the safety net",
    "These checks grade form, not substance—a schema-valid answer can still be wrong",
    "Always the first layer: cheap gates that catch catastrophic breakage before judges spend tokens",
],
 "Structural checks are the highest value-per-token layer in the whole chapter: they catch the "
 "regressions that break downstream code, for free, deterministically. Remind them of Chapter 4's "
 "Structured Outputs—when the API guarantees the shape, these checks verify the guarantee held."),
# CLONE 58_4 s54 (rubric scoring) goes here
("LLM-as-Judge: Designing the Judge", [
    "The judge is a prompt too—write it with the same rigor: criteria, examples, an output schema",
    "Give it the reference and context: judging groundedness requires seeing the source, not just the answer",
    "Absolute scoring (1–5 against the rubric) for dashboards; pairwise (A vs B) for prompt/model selection",
    "Pairwise comparison is more reliable than absolute rating—models, like humans, compare better than they score",
    "Use a strong model as judge, pin its version, and set temperature to 0 for reproducibility",
],
 "LLM-as-judge is the workhorse for open-ended quality, and most of its failures are prompt "
 "failures: vague criteria, no reference, no schema. The absolute-vs-pairwise distinction is "
 "practical—use pairwise when choosing between two candidates, absolute when tracking a score "
 "over time."),
("LLM-as-Judge: Biases and Calibration", [
    "Position bias: judges favor the first (or last) answer in a pair—randomize order, run both ways",
    "Self-preference: a judge favors outputs written in its own model family's style",
    "Verbosity bias: longer answers score higher regardless of quality",
    "Calibrate: validate the judge against human labels on a sample; report agreement as you would for a human rater",
    "Re-validate on every judge model or judge-prompt change—the judge is versioned like the system under test",
],
 "An uncalibrated judge manufactures confidence. The calibration loop is simple: label 50–100 "
 "cases by hand, measure agreement, fix the judge prompt where it disagrees, repeat. Chapter 8 "
 "adds a sharper point—the judge is also an attack surface, another reason it can never be the "
 "only gate."),
("Human Review and Spot-Checking", [
    "Humans remain the gold standard for consequential judgments—and the calibration source for judges",
    "Review a small random sample of production outputs on a cadence (for example, 20 per week)",
    "Targeted review: every judge-disagreement, every low-confidence score, every guardrail trigger",
    "Make review cheap: a queue, a rubric, two clicks—not an essay per case",
    "Feed every human verdict back into the golden dataset and into judge calibration",
],
 "Human review is not a failure of automation—it is the source of truth that keeps the automated "
 "layers honest. The design goal is a sustainable trickle of high-value judgments, not a heroic "
 "weekly marathon; every verdict should pay twice by also improving the dataset and the judge."),
]
# TABLE slide (choosing the method) goes here
# CLONE 58_7 s33 (Ragas) goes here

LEVELS = [
("Prompt-Level Unit Evals", [
    "The smallest evaluable unit: one prompt template, one model config, a golden set of inputs",
    "Assert structure programmatically; judge content against the rubric",
    "Run on every prompt edit—this is the GenAI unit test, and Chapter 5's TDD loop applies here directly",
    "Fast and cheap: milliseconds of compute, seconds of judge time",
    "Lab 6.1 builds exactly this layer",
],
 "Start the levels section at the bottom of the pyramid where feedback is fastest. A prompt-level "
 "eval is a unit test whose assertion is a scorer instead of an equality check—engineers who "
 "internalize that mapping stop treating evals as exotic. Lab 6.1 is precisely this layer."),
("Component Evals: Retrieval Quality in RAG", [
    "A RAG answer is only as good as the retrieved context—evaluate retrieval separately from generation",
    "Retrieval: did the right chunks rank for the question? (Recall@k and MRR, Chapter 3's ranking cousins)",
    "Generation: is the answer faithful to the retrieved context, and relevant to the question?",
    "Ragas (previous section) automates exactly this split",
    "Fix the failing component, not the whole pipeline: retrieval misses need index or embedding work, not prompt edits",
],
 "Component evals are how you stop prompt-tuning a retrieval problem. The Ragas slide just showed "
 "the metrics; this slide is the engineering habit: when an end-to-end RAG eval fails, decompose "
 "it into retrieval versus generation and fix the stage that actually broke."),
("Integration Evals: Testing the Seams", [
    "Components pass alone; the pipeline fails together—test the seams",
    "Prompt + retrieval + post-processing as one unit: does the assembled context produce the right answer?",
    "Contract checks between stages: retriever output schema, token budget after stuffing, tool-result formats",
    "Use recorded fixtures for external tools and APIs so integration evals stay deterministic where possible",
    "Failures here are interface bugs—treat them with the same urgency as API contract breaks",
],
 "Integration evals catch the bugs that unit evals structurally cannot: the retriever returns "
 "chunks the prompt truncates, the tool changed its response shape, the context window overflows. "
 "Recorded fixtures keep the suite deterministic and fast; live calls belong in a smaller nightly "
 "tier."),
("End-to-End and Agent Trajectory Evals", [
    "End-to-end: task completion on realistic scenarios, scored on the final outcome",
    "Agents add a second dimension—the trajectory: right tools, right order, valid arguments, bounded steps",
    "Check intermediate state, not just the last message: did it write the file, call the API, stay in budget?",
    "Chapter 7 builds the agents; its evaluation slides build directly on this one",
    "Expensive and slow—keep a small, high-value suite and run it nightly or pre-release",
],
 "For agents the answer is only half the story: how it got there determines cost, safety, and "
 "reliability. A correct answer reached by a hallucinated tool call is a failure the trajectory "
 "eval must catch. Flag Chapter 7 as the deep dive and keep this slide to the vocabulary."),
("The Eval Pyramid", [
    "Same shape as the test pyramid: many cheap prompt-level evals, fewer integration, few end-to-end",
    "Prompt-level: thousands of cases on every commit—catches most regressions",
    "Component and integration: hundreds of cases on every pull request",
    "End-to-end: dozens of scenarios, nightly or pre-release—the final safety net",
    "If a failure can be caught lower in the pyramid, move it there—cost and flakiness both drop",
],
 "The test-pyramid mapping is deliberate: engineers already know why a suite of only end-to-end "
 "tests is slow, flaky, and unlocalizable—the same economics apply to evals. When a nightly "
 "end-to-end failure turns out to be a prompt-level issue, the fix includes a new case at the "
 "lower level."),
]

EDD = [
# CLONE 16 s23 (TEVV) then CLONE 16 s24 (EDD) go here
("Write the Evals First—TDD for GenAI", [
    "Chapter 5's loop, applied to prompts: red—write the eval that captures the desired behavior; green—make the system pass; refactor with the suite holding you honest",
    "Evals-first forces the specification conversation early: what does 'good' mean for this feature, in measurable terms?",
    "New feature → new golden cases in the same pull request as the prompt that serves them",
    "Bug report → failing eval case first, fix second—and the case stays in the suite forever",
    "The suite is living documentation of everything the system must never break again",
],
 "This is the chapter's core slide and the mirror of Chapter 5's TDD story: where TDD writes the "
 "test before the code, EDD writes the eval before the prompt. The PR habit is what makes it "
 "real—prompt change and eval change land together, or the suite drifts behind the behavior it "
 "is supposed to guard."),
("Eval Suites in CI", [
    "Run the suite on every PR that touches a prompt, model config, or retrieval setting",
    "Gate on thresholds: the suite score must not drop versus main; per-slice floors for high-risk slices",
    "Handle flaky evals like flaky tests: pin judge version and temperature 0, quarantine flaky cases, never silently delete",
    "Keep CI fast: dev split on PRs, full suite nightly; cache judge calls by content hash",
    "The build log is the audit trail: dataset version, model version, scores, decision",
],
 "An eval suite that runs ad hoc is a spreadsheet; wired into CI it becomes a regression gate. "
 "Flaky-eval handling deserves the room's attention—a judge at nonzero temperature will "
 "occasionally disagree with itself, and the answer is determinism plus quarantine, not "
 "tolerating red builds or lowering the bar."),
("Regression-Testing Prompts When Models Change", [
    "Provider model upgrades are dependency upgrades—treat them like a major version bump",
    "Re-run the full suite against the new model before switching; diff scores per slice, not just the average",
    "Pin the model version in configuration—'latest' is a regression delivery mechanism",
    "Averages hide regressions: +2 overall can be +5 on easy slices and −8 on your money slice",
    "Keep the old model's results as the baseline until the new one proves itself",
],
 "Silent provider-side upgrades are the most common cause of 'the demo worked last month'. The "
 "per-slice diff is the key habit: averages move slowly while a critical slice falls off a cliff. "
 "The suite from the previous slides is exactly what makes a model swap an evidence-based "
 "decision instead of a leap of faith."),
("Canarying New Prompts and Models", [
    "Ship changes like software: canary the new prompt or model on a small traffic slice first",
    "Shadow mode: run the candidate alongside production, score both, serve the old one",
    "Promote on evidence: canary scores at or above baseline for a defined window, then roll out",
    "Roll back on signal: the suite's online metrics trip the revert—not a customer complaint",
    "Same discipline as blue/green deploys; the eval suite supplies the health check",
],
 "Offline suites answer 'is it better on our cases?'; canaries answer 'is it better on real "
 "traffic right now?'. Both questions matter, and the rollback criterion should be decided before "
 "the rollout starts. Engineers who already run blue/green deploys will recognize every mechanic "
 "here—the eval suite just becomes the health signal."),
# CLONE 16 s25 (ARIA) then CLONE 16 s26 (KPIs) go here
]

TOOLING = [
("The Eval Tooling Landscape (2026)", [
    "OpenAI eval tooling: model-graded evals and scorecards in the platform you already call (Chapter 4)",
    "Ragas: reference-free RAG metrics—retrieval, faithfulness, answer quality (seen earlier this chapter)",
    "promptfoo: declarative YAML eval suites, prompt/model comparison matrices, CI-friendly",
    "pytest-based suites: evals as tests in your repo—assertions plus judge calls, reported like test results",
    "Pick by fit, not hype: they compose—pytest orchestrates, Ragas scores RAG, promptfoo compares versions",
],
 "The 2026 landscape has converged on a few shapes: platform-native tooling from OpenAI, "
 "specialist libraries like Ragas, declarative runners like promptfoo, and plain pytest for teams "
 "who want evals next to their tests. These tools compose rather than compete—Lab 6.1 uses the "
 "pytest style so the suite lives in the repo."),
("Traces and Observability: You Can't Eval What You Can't See", [
    "Capture per request: inputs, rendered prompt, retrieved context, tool calls, output, tokens, cost",
    "Traces turn a failed eval case into a debuggable story—which stage broke, with what context",
    "OpenTelemetry-style tracing is the 2026 standard; LLM observability platforms (LangSmith, Langfuse, Arize, OpenAI dashboard) store and search traces",
    "Sample production traces into the golden dataset—real failures are the best future test cases",
    "Log scores alongside traces so quality is queryable next to latency and cost",
],
 "Without traces a failed eval is a shrug; with them it is a stack trace for a probabilistic "
 "system. The pipeline runs both ways: traces explain eval failures, and production traces "
 "replenish the golden dataset. Name-check the platforms but stress the standard—OTel-shaped "
 "spans, not a vendor."),
("Storing and Comparing Results Over Time", [
    "Record every run: dataset version, prompt version, model version, scores—an immutable history",
    "Trend the numbers: quality over time is a timeseries, like test coverage or build duration",
    "Compare runs as diffs: which cases changed score and which slices moved—not just the average",
    "Keep the outputs, not just the scores—you will want to eyeball what '−3 points' actually meant",
    "This history is your defense when someone asks 'did the upgrade actually help?'",
],
 "Eval results without history are trivia; with history they are an engineering instrument. The "
 "diff view is the daily driver—average movement tells you something moved, the per-case diff "
 "tells you what to go fix. Keeping raw outputs matters because scores compress away exactly the "
 "detail you need when debugging."),
("Reporting Scorecards to Stakeholders", [
    "One page: headline scores, trend arrows, per-slice floors, cost and latency—no judge-prompt trivia",
    "Translate to product language: 'answer accuracy on billing questions', not 'rubric dimension 2'",
    "Show the gap honestly: confidence intervals and slice sizes belong on the scorecard",
    "The program KPIs from the previous section roll up here—coverage, catch rate, MTTR",
    "The scorecard's job is a decision—ship, hold, or invest; design it so the decision is obvious",
],
 "Stakeholders do not want your eval methodology; they want to know whether the feature is "
 "getting better and whether it is safe to ship. The honest-scorecard habits—intervals, slice "
 "sizes—are what survive the first time a scorecard predicts success and production disagrees."),
("The Cost of Evals and Sampling Strategies", [
    "Judge calls cost tokens: a 500-case suite against a frontier judge is a real line item—budget it like test infrastructure",
    "Stratified subsampling: run every case through cheap checks, judge a sample, deep-review a handful",
    "Cache judge verdicts by content hash—unchanged outputs should never be re-judged",
    "Tier the schedule: cheap layers every commit, judged layers on PRs, the full suite nightly",
    "Track eval spend like coverage: cutting eval coverage to save money is a false economy",
],
 "Eval economics mirror test economics: the cheapest layer that can catch a class of failure "
 "should own it. Content-hash caching is the single biggest win—most PRs change a handful of "
 "outputs, so most judge verdicts are reusable. The tiered schedule from the eval pyramid is also "
 "the cost control."),
("Online Evaluation: A/B Tests and Shadow Deployments", [
    "Offline suites miss what only users reveal—pair them with online measurement",
    "A/B test prompt or model variants on live traffic, with product metrics and sampled judge scores",
    "Shadow deployments score the candidate on real traffic without serving it to anyone",
    "Guardrail metrics decide rollback automatically: error rate, refusal rate, latency, cost per request",
    "Close the loop: online failures become golden cases—the suite grows from production evidence",
],
 "Online evaluation is where the chapter's loop meets the deployment pipeline: the canary slide "
 "gave the mechanics, this slide gives the measurement. The closing loop is the point—production "
 "failures flow back into the golden dataset, so the offline suite gets stronger every week the "
 "system runs."),
]

LAB_TITLE = "Lab 6.1: Write an Eval Suite"
LAB_BULLETS = [
    "Follow the detailed instructions in the Lab 6.1 notebook on your VM",
    "Build a small golden dataset: realistic inputs, references, and a few adversarial cases",
    "Implement two scorers—a 1–5 rubric and an LLM-as-judge prompt—and compare their verdicts",
    "Run a prompt change through the suite and catch the regression it introduces",
]
LAB_NOTES = (
    "The notebook provides a pytest-style starter suite over a small summarization prompt, with "
    "the OpenAI key pre-injected on the VM. The pedagogical moment is the last step: a prompt edit "
    "that 'reads better' gets measurably worse on one slice, and the suite catches it—that single "
    "experience is the whole chapter. Circulate during the judge-comparison step; students are "
    "often surprised how often their rubric and judge disagree.")


# --- formatting-preserving text rewrite helpers (same as build_ch07) -----------
def _set_para_text(p_el, text):
    runs = p_el.findall(qn("a:r"))
    if not runs:
        r = p_el.makeelement(qn("a:r"), {})
        r.append(p_el.makeelement(qn("a:t"), {}))
        p_el.append(r)
        runs = [r]
    first = runs[0]
    for extra in runs[1:]:
        p_el.remove(extra)
    for br in p_el.findall(qn("a:br")):
        p_el.remove(br)
    t = first.find(qn("a:t"))
    if t is None:
        t = first.makeelement(qn("a:t"), {})
        first.append(t)
    t.text = text


def rewrite_bullets(shape, lines):
    txBody = shape.text_frame._txBody
    paras = txBody.findall(qn("a:p"))
    for i, line in enumerate(lines):
        if i < len(paras):
            _set_para_text(paras[i], line)
        else:
            newp = copy.deepcopy(paras[-1])
            txBody.append(newp)
            paras.append(newp)
            _set_para_text(newp, line)
    for p in paras[len(lines):]:
        txBody.remove(p)


def _iter_paras(shape):
    """Yield every a:p element in a shape (recurses into groups and tables)."""
    if shape.shape_type == 6:  # GROUP
        for sub in shape.shapes:
            yield from _iter_paras(sub)
    elif getattr(shape, "has_table", False) and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                yield from cell.text_frame._txBody.findall(qn("a:p"))
    elif shape.has_text_frame:
        yield from shape.text_frame._txBody.findall(qn("a:p"))


def patch_slide(slide, pairs, tag):
    """Substring-replace inside paragraphs (incl. table cells), preserving the
    first run's formatting. Every pair must hit at least once."""
    paras = [p for sh in slide.shapes for p in _iter_paras(sh)]
    missed = []
    for old, new in pairs:
        hit = False
        for p in paras:
            joined = "".join(t.text or "" for t in p.iter(qn("a:t")))
            if old in joined:
                _set_para_text(p, joined.replace(old, new))
                hit = True
                break
        if not hit:
            missed.append(old[:60])
    assert not missed, f"patch MISS on {tag}: {missed}"


def patch_notes(slide, pairs, tag):
    """Substring-replace in speaker notes. Warn-only (long prose, low risk)."""
    if not slide.has_notes_slide:
        print(f"  !! notes patch MISS on {tag}: no notes slide")
        return
    tf = slide.notes_slide.notes_text_frame
    text = tf.text
    missed = [old[:60] for old, new in pairs if old not in text]
    for old, new in pairs:
        text = text.replace(old, new)
    tf.text = text
    if missed:
        print(f"  !! notes patch MISS on {tag}: {missed}")


def add_two_col_slide(prs, title, left, right, notes):
    """Contrast slide on 'Two Column Full Page': header line + short bullets
    in each column."""
    lay = sk.pick_layout(prs, "Two Column Full Page")
    assert lay is not None, "Two Column Full Page layout not found in carrier"
    slide = prs.slides.add_slide(lay)
    sk.set_title(slide, title)
    cols = {}
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx in (1, 13) and ph.has_text_frame:
            cols[idx] = ph
    assert len(cols) == 2, "two-column placeholders not found"
    for idx, lines in ((1, left), (13, right)):
        ph = cols[idx]
        ph.height = Inches(4.9)      # layout's column box is stubby; give it room
        tf = ph.text_frame
        tf.word_wrap = True
        for i, b in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
    sk.set_notes(slide, notes)
    return slide


def add_method_table(prs, title, header, rows, notes, widths_in):
    """Cost-vs-fidelity table on the content layout, sized for the 10in deck."""
    slide = sk.add_table_slide(prs, title, header, rows, notes=notes,
                               layout=sk.pick_layout(prs, "Content with Header. Full Page"))
    gf = next(sh for sh in slide.shapes if sh.has_table)
    gf.height = Inches(0.62 * (len(rows) + 1))
    for i, w in enumerate(widths_in):
        gf.table.columns[i].width = Inches(w)
    return slide


def add_exercise_slide(prs, title, bullets, notes):
    """Lab slide on the template's exercise layout (typing-hands graphic on the
    left, text column on the right)."""
    lay = sk.pick_layout(prs, "Exercise Reference Slide with Typing Hands")
    assert lay is not None, "exercise layout not found in carrier template"
    slide = prs.slides.add_slide(lay)
    sk.set_title(slide, title)
    body = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1 and ph.has_text_frame:
            body = ph
            break
    if body is None:
        body = pt._body_placeholder(slide)
    body.height = Inches(4.5)        # the layout's content placeholder is stubby
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
    sk.set_notes(slide, notes)
    return slide


def main():
    # 0. fresh carrier -> idempotent
    shutil.copy(CARRIER, OUT)
    prs = sk.open_prs(OUT)

    # 1. clone from the three sources (appended at end), patch + retitle
    cidx = {}
    decks = {tag: pt.open_deck(path) for tag, path in SRC.items()}
    for tag, keeps in CLONES.items():
        for k in keeps:
            dst = pt.clone_slide(prs, decks[tag].slides[k - 1])
            cidx[(tag, k)] = len(prs.slides) - 1
            if (tag, k) in PATCHES:
                patch_slide(dst, PATCHES[(tag, k)], f"{tag} s{k}")
            if (tag, k) in NOTES_PATCHES:
                patch_notes(dst, NOTES_PATCHES[(tag, k)], f"{tag} s{k}")
            if (tag, k) in RETITLES:
                sk.set_title(dst, RETITLES[(tag, k)])

    # 2. author the new slides (appended at end), recording indices
    def add(title, bullets, notes):
        pt.append_content(prs, title, bullets, notes)
        return len(prs.slides) - 1

    i_why = [add(t, b, n) for t, b, n in WHY[:4]]
    add_two_col_slide(prs, TWO_COL[0], TWO_COL[1], TWO_COL[2], TWO_COL[3])
    i_why.append(len(prs.slides) - 1)
    i_why.append(add(*WHY[4]))
    i_data = [add(t, b, n) for t, b, n in DATASETS]
    i_meth = [add(t, b, n) for t, b, n in METHODS[:2]]
    i_meth.append(cidx[("58_4", 54)])                     # rubric clone
    i_meth += [add(t, b, n) for t, b, n in METHODS[2:]]
    add_method_table(
        prs, "Choosing the Right Evaluation Method",
        ["Method", "Cost per case", "Fidelity", "Best for"],
        [["Exact match / assertions", "~free", "Low (strict)", "Closed-form answers, tool-call args"],
         ["Regex / schema / JSON checks", "~free", "Low–medium", "Structure and output contracts"],
         ["Rubric scoring (1–5)", "$", "Medium–high", "Comparing prompt or model versions"],
         ["LLM-as-judge", "$$", "High if calibrated", "Open-ended quality, tone, instruction-following"],
         ["Human review", "$$$", "Highest", "Consequential calls; judge validation"]],
        "The rule of thumb: start at the top of the table and escalate only where judgment is "
        "genuinely needed—every row below the first two spends either tokens or people. The "
        "fidelity column is a ceiling, not a guarantee: an uncalibrated judge scores below its "
        "row, which is why the last row is also the calibration source for the fourth.",
        [2.5, 1.4, 1.7, 3.2])
    i_meth.append(len(prs.slides) - 1)
    i_meth.append(cidx[("58_7", 33)])                     # Ragas clone
    i_levels = [add(t, b, n) for t, b, n in LEVELS]
    i_edd = [cidx[("16", 23)], cidx[("16", 24)]]          # TEVV, EDD clones
    i_edd += [add(t, b, n) for t, b, n in EDD]
    i_edd += [cidx[("16", 25)], cidx[("16", 26)]]         # ARIA, KPI clones
    i_tool = [add(t, b, n) for t, b, n in TOOLING]
    add_exercise_slide(prs, LAB_TITLE, LAB_BULLETS, LAB_NOTES)
    i_lab = len(prs.slides) - 1

    # 3. rewrite kept carrier slides in place (0-based: 0,1,2,31,32)
    s1 = prs.slides[0]
    _set_para_text(s1.shapes.title.text_frame._txBody.findall(qn("a:p"))[0],
                   "Evaluating Generative AI Systems")
    for sh in s1.shapes:                       # subtitle: Chapter 5 -> Chapter 6
        if sh.has_text_frame and "Chapter 5" in sh.text_frame.text:
            _set_para_text(sh.text_frame._txBody.findall(qn("a:p"))[0], "Chapter 6")
    rewrite_bullets(pt._body_placeholder(prs.slides[1]), OBJECTIVES)   # Objectives
    sk.set_notes(prs.slides[1], OBJECTIVES_NOTES)                      # (old notes were stale)
    rewrite_bullets(pt._body_placeholder(prs.slides[2]), AGENDA)       # Agenda
    rewrite_bullets(pt._body_placeholder(prs.slides[31]), SUMMARY)     # Summary
    sk.set_notes(prs.slides[31], SUMMARY_NOTES)                        # (old notes were stale)
    rewrite_bullets(pt._body_placeholder(prs.slides[32]), OBJECTIVES)  # objectives recap

    # 4. arrange into the final teaching order
    order = (
        [0, 1, 2]            # openers
        + i_why              # why evals (incl. two-column contrast + eval loop)
        + i_data             # eval datasets
        + i_meth             # evaluation methods (incl. rubric clone, table, Ragas clone)
        + i_levels           # levels of evaluation
        + i_edd              # eval-driven development (TEVV/EDD/ARIA/KPI clones + authored)
        + i_tool             # tooling & ops
        + [i_lab]            # Lab 6.1
        + [31, 32]           # summary + objectives recap
    )
    pt.arrange(prs, order)
    pt.save(prs, OUT)

    # 5. verify: reopen, titles, blanks, duplicate partnames, zip integrity
    prs2 = Presentation(OUT)
    n = len(prs2.slides)
    print(f"\n{OUT.name}: {n} slides")
    print(f"clones: 2016-Ch06={len(CLONES['16'])} 1258-Ch04={len(CLONES['58_4'])} "
          f"1258-Ch07={len(CLONES['58_7'])}  authored=33  carrier-kept=5")
    empty = []
    for i, s in enumerate(prs2.slides, 1):
        t = sk.slide_title(s)
        if not t.strip():
            empty.append(i)
        print(f"{i:>3} | {t}")
    assert not empty, f"empty titles at {empty}"
    assert 42 <= n <= 48, f"slide count {n} outside 42-48"
    # every authored slide carries speaker notes (kept openers excepted)
    no_notes = [i for i, s in enumerate(prs2.slides, 1)
                if not (s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip())]
    print(f"slides without notes: {no_notes} (expected: openers/closers only)")
    with zipfile.ZipFile(OUT) as z:
        dupes = [p for p, c in Counter(z.namelist()).items() if c > 1]
        assert not dupes, f"duplicate zip partnames: {dupes}"
        assert z.testzip() is None, "zip integrity check failed"
    print("\nOK: no empty titles, no duplicate partnames, zip integrity verified")


if __name__ == "__main__":
    main()
