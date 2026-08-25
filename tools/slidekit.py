"""
slidekit.py — shared, tested PowerPoint editing primitives for the Course 2016
enrichment pass. python-pptx has no native "delete slide", "move slide", or
"clone slide", so this module provides reliable versions plus convenience
builders for adding template-faithful content slides.

Use the venv interpreter:  ~/.venvs/pptx/bin/python

Typical usage (per chapter, on a COPY in enriched/):

    import slidekit as sk
    prs = sk.open_prs("Ch03-Pass5.pptx")
    sk.dump(prs)                         # inspect (idx | title | body-preview)

    sk.delete_slide(prs, 12)            # remove a duplicate "Contents" slide
    s = sk.add_bullets_slide(prs, "Direct vs Indirect Injection",
            ["Direct: attacker controls the user prompt",
             "Indirect: payload hides in a retrieved document",
             "Both exploit the missing trust boundary"],
            notes="Speaker note ...", at_index=17)   # insert at position 17
    sk.add_source(s, "OWASP LLM01", "https://genai.owasp.org")

    sk.save(prs, "Ch03-Pass5.pptx")
    sk.validate("Ch03-Pass5.pptx")      # reopens; prints slide count, asserts OK

DESIGN NOTES
- add_*_slide() clones the slide LAYOUT used by an existing reference slide, so
  new slides inherit the LT master (colors, fonts, footer). They are static
  (no custom animations / SmartArt) and may need light manual polish.
- Index arguments are 0-based over prs.slides.
- Always work on copies in enriched/. Never touch the originals one dir up.
"""
from __future__ import annotations
import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# --- Safe slide partname allocation -----------------------------------------
# python-pptx names a new slide part "slide{len(sldIdLst)+1}.xml". After you
# DELETE slides, that number collides with an existing high-numbered slide, so
# the next add writes a DUPLICATE "slideN.xml" and silently clobbers a real
# slide on save. We override it to always pick max(existing)+1, which never
# collides regardless of add/delete ordering.
import re as _re
from pptx.parts.presentation import PresentationPart as _PresPart
from pptx.opc.packuri import PackURI as _PackURI


def _safe_next_slide_partname(self):
    used = []
    for part in self.package.iter_parts():
        m = _re.match(r"/ppt/slides/slide(\d+)\.xml$", str(part.partname))
        if m:
            used.append(int(m.group(1)))
    n = (max(used) + 1) if used else 1
    return _PackURI("/ppt/slides/slide%d.xml" % n)


_PresPart._next_slide_partname = property(_safe_next_slide_partname)
# ----------------------------------------------------------------------------

# ---------- open / save / inspect ----------

def open_prs(path: str) -> Presentation:
    return Presentation(path)


def save(prs: Presentation, path: str) -> None:
    prs.save(path)


def slide_title(slide) -> str:
    t = slide.shapes.title
    if t is not None and t.has_text_frame:
        return t.text_frame.text.strip()
    # fall back to first text placeholder
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return sh.text_frame.text.strip().splitlines()[0]
    return ""


def slide_body_preview(slide, n: int = 90) -> str:
    parts = []
    title = slide_title(slide)
    for sh in slide.shapes:
        if sh.has_text_frame:
            txt = sh.text_frame.text.strip()
            if txt and txt != title:
                parts.append(txt.replace("\n", " / "))
    s = " | ".join(parts)
    return (s[:n] + "…") if len(s) > n else s


def dump(prs: Presentation) -> None:
    """Print idx | title | body preview for every slide."""
    for i, s in enumerate(prs.slides):
        print(f"{i:>3} | {slide_title(s)[:48]:<48} | {slide_body_preview(s)}")


def is_blank(slide) -> bool:
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return False
    return True


def find_contents(prs: Presentation):
    """Return indices of slides that look like a Contents/TOC slide."""
    import re
    out = []
    for i, s in enumerate(prs.slides):
        t = (slide_title(s) + " " + slide_body_preview(s, 200)).lower()
        if re.search(r"\bcontents\b|chapter outline|table of contents", t):
            out.append(i)
    return out


# ---------- low-level slide list surgery ----------

def _sldIdLst(prs):
    return prs.slides._sldIdLst


def delete_slide(prs: Presentation, index: int) -> None:
    """Remove the slide at `index` (0-based): drop its sldId and the rel."""
    sldIdLst = _sldIdLst(prs)
    ids = list(sldIdLst)
    el = ids[index]
    rId = el.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(el)


def move_slide(prs: Presentation, from_idx: int, to_idx: int) -> None:
    """Reorder: move the slide currently at from_idx so it lands at to_idx."""
    sldIdLst = _sldIdLst(prs)
    ids = list(sldIdLst)
    el = ids[from_idx]
    sldIdLst.remove(el)
    ids = list(sldIdLst)
    if to_idx >= len(ids):
        sldIdLst.append(el)
    else:
        sldIdLst.insert(list(sldIdLst).index(ids[to_idx]), el)


# ---------- layout helpers ----------

def layout_of(prs, ref_index: int):
    """The slide layout used by an existing reference slide (for brand match)."""
    return prs.slides[ref_index].slide_layout


def pick_layout(prs, *name_substrings):
    """Find a layout by case-insensitive name match across all masters.
    Pass several candidates, first hit wins. Returns None if not found."""
    subs = [s.lower() for s in name_substrings]
    for master in prs.slide_masters:
        for lay in master.slide_layouts:
            nm = (lay.name or "").lower()
            if any(s in nm for s in subs):
                return lay
    return None


def list_layouts(prs):
    out = []
    for mi, master in enumerate(prs.slide_masters):
        for li, lay in enumerate(master.slide_layouts):
            out.append((mi, li, lay.name))
    return out


# ---------- placeholder access ----------

def _title_ph(slide):
    if slide.shapes.title is not None:
        return slide.shapes.title
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            return ph
    return None


def _body_ph(slide):
    """First body/content/text placeholder (idx != 0)."""
    for ph in slide.placeholders:
        pf = ph.placeholder_format
        if pf.idx != 0 and ph.has_text_frame:
            return ph
    return None


def set_title(slide, text: str) -> None:
    ph = _title_ph(slide)
    if ph is not None:
        ph.text_frame.text = text


def set_bullets(slide, bullets, levels=None) -> None:
    """Fill the body placeholder with bullets. `bullets` = list[str];
    optional `levels` = list[int] (0-based indent) parallel to bullets."""
    ph = _body_ph(slide)
    if ph is None:
        return
    tf = ph.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        if levels:
            p.level = levels[i]


def set_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


# ---------- high-level builders ----------

def add_bullets_slide(prs, title, bullets, levels=None, notes=None,
                      ref_index=7, layout=None, at_index=None):
    """Add a title+bullets slide using the layout of an existing slide
    (ref_index) so it inherits LT branding. Optionally insert at_index and
    attach speaker notes. Returns the new slide."""
    lay = layout or layout_of(prs, ref_index)
    slide = prs.slides.add_slide(lay)
    set_title(slide, title)
    # ensure there is a body placeholder; many content layouts have one
    if _body_ph(slide) is None:
        _add_textbox_body(slide)
    set_bullets(slide, bullets, levels)
    if notes:
        set_notes(slide, notes)
    if at_index is not None:
        move_slide(prs, len(prs.slides._sldIdLst) - 1, at_index)
    return slide


def _add_textbox_body(slide):
    """Fallback body when a layout lacks a content placeholder."""
    left, top, width, height = Inches(0.7), Inches(1.6), Inches(8.6), Inches(5.0)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    return tb


def add_table_slide(prs, title, header, rows, notes=None, ref_index=7,
                    layout=None, at_index=None):
    """Add a slide with a title and a simple table (header + rows).
    `header` = list[str]; `rows` = list[list[str]]. Returns the new slide."""
    lay = layout or layout_of(prs, ref_index)
    slide = prs.slides.add_slide(lay)
    set_title(slide, title)
    ncol = len(header)
    nrow = len(rows) + 1
    left, top, width, height = Inches(0.6), Inches(1.7), Inches(8.8), Inches(0.4 * nrow)
    gframe = slide.shapes.add_table(nrow, ncol, left, top, width, height)
    tbl = gframe.table
    for c, h in enumerate(header):
        tbl.cell(0, c).text = str(h)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            tbl.cell(r, c).text = str(val)
    if notes:
        set_notes(slide, notes)
    if at_index is not None:
        move_slide(prs, len(prs.slides._sldIdLst) - 1, at_index)
    return slide


def add_image_slide(prs, title, img_path, notes=None, credit=None,
                    ref_index=7, layout=None, at_index=None):
    """Add a title + single-image slide. The picture is scaled to fit under the
    header while preserving aspect ratio and centered horizontally. Optional
    `credit` renders a small italic attribution line at the bottom. Inherits LT
    branding from `layout` (or the layout of ref_index). Returns the new slide."""
    from PIL import Image as _PILImage
    lay = layout or layout_of(prs, ref_index)
    slide = prs.slides.add_slide(lay)
    set_title(slide, title)
    sw, sh = prs.slide_width, prs.slide_height
    margin = Inches(0.5)
    top = Inches(1.55)
    bottom = Inches(1.0) if credit else Inches(0.45)
    avail_w = sw - 2 * margin
    avail_h = sh - top - bottom
    iw, ih = _PILImage.open(img_path).size
    ar = iw / ih
    if avail_w / avail_h > ar:          # height-bound
        h = int(avail_h); w = int(h * ar)
    else:                                # width-bound
        w = int(avail_w); h = int(w / ar)
    left = int((sw - w) / 2)
    slide.shapes.add_picture(img_path, left, top, width=w, height=h)
    if credit:
        tb = slide.shapes.add_textbox(margin, sh - Inches(0.9), avail_w, Inches(0.3))
        tf = tb.text_frame; tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        run.text = credit
        run.font.size = Pt(9); run.font.italic = True
        try:
            run.font.color.rgb = RGBColor(0x8A, 0x8A, 0x8A)
        except Exception:
            pass
    if notes:
        set_notes(slide, notes)
    if at_index is not None:
        move_slide(prs, len(prs.slides._sldIdLst) - 1, at_index)
    return slide


def add_source(slide, label: str, url: str) -> None:
    """Append a small 'Source: <label> — <url>' line at the bottom of a slide."""
    left, top, width, height = Inches(0.6), Inches(6.9), Inches(8.8), Inches(0.4)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Source: {label} — {url}"
    run.font.size = Pt(9)
    run.font.italic = True
    try:
        run.font.color.rgb = RGBColor(0x8A, 0x8A, 0x8A)
    except Exception:
        pass


def duplicate_slide(prs, index, at_index=None):
    """Best-effort clone of an existing slide (keeps its shapes/styling) so you
    can tweak the copy's text. Copies images/relationships. Returns new slide.
    Prefer add_bullets_slide for plain text content; use this to riff on a
    richly-styled existing slide."""
    source = prs.slides[index]
    new_slide = prs.slides.add_slide(source.slide_layout)
    # strip placeholders that add_slide created from the layout
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)
    # relate the new slide part to the same targets, recording old->new rId
    rid_map = {}
    for rId, rel in source.part.rels.items():
        if "notesSlide" in rel.reltype:
            continue
        if rel.is_external:
            new_rId = new_slide.part.relate_to(
                rel.target_ref, rel.reltype, is_external=True)
        else:
            new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
        rid_map[rId] = new_rId
    # deep-copy every shape element, remapping its r:embed / r:link / r:id refs
    ref_attrs = (qn("r:embed"), qn("r:link"), qn("r:id"))
    for shp in source.shapes:
        el = copy.deepcopy(shp._element)
        for node in el.iter():
            for attr in ref_attrs:
                val = node.get(attr)
                if val and val in rid_map:
                    node.set(attr, rid_map[val])
        new_slide.shapes._spTree.append(el)
    if at_index is not None:
        move_slide(prs, len(prs.slides._sldIdLst) - 1, at_index)
    return new_slide


# ---------- validation ----------

def validate(path: str, expect_min=None, expect_max=None) -> int:
    """Reopen a saved deck (proves it's not corrupt), print a summary, return
    slide count. Optionally assert count within [expect_min, expect_max]."""
    prs = Presentation(path)
    n = len(prs.slides)
    blanks = sum(1 for s in prs.slides if is_blank(s))
    toc = len(find_contents(prs))
    print(f"{path}: {n} slides | {blanks} blank | {toc} contents/TOC")
    if expect_min is not None:
        assert n >= expect_min, f"{n} < expected min {expect_min}"
    if expect_max is not None:
        assert n <= expect_max, f"{n} > expected max {expect_max}"
    return n


if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        validate(sys.argv[1])
