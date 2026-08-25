#!/usr/bin/env python3
"""build_engagement_ch01.py — Day-1 engagement pass on the A.2 Ch01 deck.

Carrier = _bak1/1851-Ch01.pptx (restored fresh every run -> idempotent).
Adds two unnumbered Do Now warm-ups ('Spot the AI' after the Agenda; 'Your AI
Toolbox' before 'AI Readiness for a Software Team'), adds the Lab 1.1 exercise
slide just before the Summary, and updates the Agenda's last bullet to mention
Lab 1.1. 38 -> 41 slides.

Run:  /Users/iahmad/Creator/Courses_and_conferences/LT/.venv-courseware/bin/python tools/build_engagement_ch01.py
"""
from __future__ import annotations
import copy
import shutil
import sys
import zipfile
from collections import Counter
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import slidekit as sk  # noqa: F401  (import activates the partname-collision fix)
import pptx_tools as pt
from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
BAK = ROOT / "decks" / "_bak1"
OUT = ROOT / "decks" / "1851-Ch01.pptx"

BODY_H = Inches(4.5)   # raise stubby layout-default body areas (house pattern)


def layout_named(prs, name):
    for lay in prs.slide_layouts:
        if (lay.name or "") == name:
            return lay
    raise KeyError(f"layout {name!r} not found")


def find_slide(prs, title_sub):
    """First slide index whose title contains title_sub (case-insensitive)."""
    sub = title_sub.lower()
    for i, s in enumerate(prs.slides):
        if sub in sk.slide_title(s).lower():
            return i
    raise KeyError(f"slide titled like {title_sub!r} not found")


def add_engagement_slide(prs, layout_name, title, bullets, notes, at_index):
    """Add a slide on a named Do Now / Discussion / Exercise layout, fill the
    main body placeholder (largest text area), raise its stubby height, attach
    notes, and move it to at_index."""
    slide = prs.slides.add_slide(layout_named(prs, layout_name))
    sk.set_title(slide, title)
    body = pt._body_placeholder(slide)               # largest text area (idx 1)
    assert body is not None, f"no body placeholder on layout {layout_name!r}"
    body.width = body.width    # pin inherited width: python-pptx writes cx=0 otherwise
    body.height = BODY_H
    tf = body.text_frame
    tf.word_wrap = True
    for k, b in enumerate(bullets):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.text = b
    sk.set_notes(slide, notes)
    sk.move_slide(prs, len(prs.slides._sldIdLst) - 1, at_index)
    return slide


def _set_para_text(p_el, text):
    """Replace a paragraph's text, keeping the first run's formatting."""
    runs = p_el.findall(qn("a:r"))
    if not runs:
        r = p_el.makeelement(qn("a:r"), {})
        r.append(p_el.makeelement(qn("a:t"), {}))
        p_el.append(r)
        runs = [r]
    first = runs[0]
    for extra in runs[1:]:
        p_el.remove(extra)
    for br in p_el.findall(qn("a:br")):
        p_el.remove(br)
    t = first.find(qn("a:t"))
    if t is None:
        t = first.makeelement(qn("a:t"), {})
        first.append(t)
    t.text = text


def replace_bullet(shape, old_sub, new_text):
    """Rewrite one agenda bullet in place, preserving its formatting."""
    for p in shape.text_frame._txBody.findall(qn("a:p")):
        if old_sub.lower() in "".join(
                t.text or "" for t in p.iter(qn("a:t"))).lower():
            _set_para_text(p, new_text)
            return
    raise KeyError(f"bullet containing {old_sub!r} not found")


def main():
    # 0. fresh carrier -> idempotent
    shutil.copy(BAK / "1851-Ch01.pptx", OUT)
    prs = sk.open_prs(OUT)
    assert len(prs.slides) == 38, "carrier deck should have 38 slides"

    # 1. Do Now: Spot the AI — right after the Agenda (position 4)
    add_engagement_slide(prs, "Discussion", "Do Now: Spot the AI",
        bullets=[
            "A teammate shows you five artifacts from last week’s sprint:",
            "1) A regex that validates postal codes   2) A unit test for the date parser",
            "3) A one-page design doc   4) A one-line bug fix   5) A commit message",
            "In pairs: which of these were AI-generated—and what evidence are you using?",
            "Time box: 5 minutes, then a show of hands per artifact",
        ],
        notes=("Run it as a rapid vote: read each artifact, take hands for human vs. AI. "
               "Answer key: it is deliberately unanswerable—any or all five could be AI-generated, "
               "and there is no reliable tell in the artifact itself. Debrief: detection is a losing "
               "game; provenance, review, and tests are what keep quality up. This sets up the "
               "trust-but-verify theme that Lab 1.1 makes concrete."),
        at_index=3)

    # 2. Do Now: Your AI Toolbox — just before 'AI Readiness for a Software Team'
    add_engagement_slide(prs, "Do Now Writing Offline", "Do Now: Your AI Toolbox",
        bullets=[
            "On your own, offline: list every AI tool you already use at work",
            "Assistants, chatbots, CI bots, search—anything that touches your workflow",
            "Next to each, note one task where it genuinely saves you time",
            "Then write one task you would NEVER hand to an AI—and why not",
            "Time box: 5 minutes; we will hear a few “never” items",
        ],
        notes=("Give them quiet writing time—this one works best offline and individual. "
               "Debrief by collecting the “never” items on the whiteboard; expect production "
               "credentials, incident comms, and anything with customer data. Map those answers "
               "to data boundaries and accountability, which is exactly what the AI-readiness "
               "slide that follows covers."),
        at_index=find_slide(prs, "AI Readiness for a Software Team"))

    # 3. Lab 1.1 — just before the Summary slide
    add_engagement_slide(prs, "Exercise Reference Slide with Typing Hands",
        "Lab 1.1: Assistant Recon — Where AI Helps, Where It Fails",
        bullets=[
            "Follow the detailed instructions in the Lab 1.1 notebook on your VM",
            "Probe a coding assistant on six task types—from boilerplate and tests to a novel algorithm and a security-sensitive change",
            "Log successes and failures on the scorecard, flagging output that was plausible but wrong",
            "Calibrate your trust: finish with your own “use freely / verify carefully / never” list",
            "Time: ~30 minutes",
        ],
        notes=("This is the day-one hands-on baseline: learners systematically probe where an "
               "assistant is strong and where it fails silently, instead of trusting vibes. "
               "Circulate during the lab and push anyone who finishes early to retry a failed task "
               "with a better prompt and note whether it changes the verdict. Debrief with the "
               "scorecards: most classes find the “plausible but wrong” column is the eye-opener."),
        at_index=find_slide(prs, "Summary"))

    # 4. Agenda: last bullet now carries the Lab 1.1 mention (still 7 bullets)
    replace_bullet(pt._body_placeholder(prs.slides[2]),
                   "Case study, course map, and Activity 1.2",
                   "Case study, course map, Activity 1.2, and Lab 1.1")

    pt.save(prs, OUT)

    # 5. verify: reopen, titles, blanks, duplicate partnames, zip integrity
    prs2 = Presentation(OUT)
    n = len(prs2.slides)
    print(f"\n{OUT.name}: {n} slides")
    empty = []
    for i, s in enumerate(prs2.slides, 1):
        t = sk.slide_title(s)
        if not t.strip():
            empty.append(i)
        print(f"{i:>3} | {t}")
    assert not empty, f"empty titles at {empty}"
    assert n == 41, f"slide count {n} != 41"
    with zipfile.ZipFile(OUT) as z:
        dupes = [p for p, c in Counter(z.namelist()).items() if c > 1]
        assert not dupes, f"duplicate zip partnames: {dupes}"
        assert z.testzip() is None, "zip integrity check failed"
    print("\nOK: no empty titles, no duplicate partnames, zip integrity verified")


if __name__ == "__main__":
    main()
