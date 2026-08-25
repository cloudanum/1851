"""
pptx_tools.py — deck surgery helpers for the 1258 a2->a3 revision.

python-pptx has no native "delete slide" or "copy slide", so this module
implements them via the underlying XML. All operations preserve the source
deck's masters, layouts, theme, and media because new decks START as a byte
copy of a source .pptx (shutil.copy) and we only delete/append from there.

Functions:
  inventory(path)                       -> list of (idx, title, n_words)
  delete_by_index(path, drop, out)      -> new deck with those 1-based slides removed
  keep_by_index(path, keep, out)        -> new deck with only those 1-based slides
  add_section(path, title, subtitle)    -> append a section-header slide
  add_content(path, title, bullets,notes,after=None) -> append a content slide
  clone_slide(dst_prs, src_slide)       -> deep-copy a slide (incl. images) into dst
  merge_decks(base, add_path, add_keep, out) -> base + selected slides of another deck
"""
from __future__ import annotations
import copy
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt


def _slide_title(slide) -> str:
    if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
        t = slide.shapes.title.text_frame.text.strip().replace("\n", " ")
        if t:
            return t
    # fall back to first non-empty text
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return sh.text_frame.text.strip().split("\n")[0][:80]
    return "(no title)"


def _slide_words(slide) -> int:
    n = 0
    for sh in slide.shapes:
        if sh.has_text_frame:
            n += len(sh.text_frame.text.split())
    return n


def inventory(path):
    prs = Presentation(path)
    return [(i + 1, _slide_title(s), _slide_words(s)) for i, s in enumerate(prs.slides)]


def _drop_slide(prs, index0):
    """Remove slide at 0-based index from the presentation's sldIdLst."""
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    sldId = slides[index0]
    rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def delete_by_index(path, drop_1based, out):
    prs = Presentation(path)
    for idx0 in sorted({i - 1 for i in drop_1based}, reverse=True):
        _drop_slide(prs, idx0)
    prs.save(out)
    return out


def keep_by_index(path, keep_1based, out):
    n = len(Presentation(path).slides._sldIdLst)
    drop = [i for i in range(1, n + 1) if i not in set(keep_1based)]
    return delete_by_index(path, drop, out)


from pptx.enum.shapes import PP_PLACEHOLDER

# Preferred content layouts by name, in order (LT template names).
_CONTENT_LAYOUT_NAMES = [
    "Content with Header. Full Page",
    "Content with Long Headline. Full Pg",
    "Full Page Blank",
]


def _layout_by_name(prs, names):
    by_name = {(lay.name or ""): lay for lay in prs.slide_layouts}
    for n in names:
        if n in by_name:
            return by_name[n]
    return None


def _content_layout(prs):
    lay = _layout_by_name(prs, _CONTENT_LAYOUT_NAMES)
    if lay is not None:
        return lay
    # fallback: a layout that has a real TITLE (idx 0) + a BODY placeholder
    for lay in prs.slide_layouts:
        idxs = {ph.placeholder_format.idx for ph in lay.placeholders}
        types = {ph.placeholder_format.type for ph in lay.placeholders}
        if 0 in idxs and PP_PLACEHOLDER.BODY in types:
            return lay
    return prs.slide_layouts[0]


def _body_placeholder(slide):
    """Return the MAIN content placeholder (largest text-capable area, not the
    title). In this template the main bullet area is the OBJECT placeholder
    (idx 1), while the BODY placeholder (idx 12) is a tiny footer caption — so we
    choose by area rather than by placeholder type."""
    candidates = []
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0 or not ph.has_text_frame:
            continue
        w = ph.width or 0
        h = ph.height or 0
        candidates.append((int(w) * int(h), ph))
    if not candidates:
        return None
    candidates.sort(key=lambda t: t[0], reverse=True)
    return candidates[0][1]


def _section_layout(prs):
    # New section-opener slides just use the content layout (a title + framing
    # bullets) to avoid the template's title-less section placeholders.
    return _content_layout(prs)


def _set_notes(slide, notes):
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def add_content(path, title, bullets, notes="", out=None):
    prs = Presentation(path)
    slide = prs.slides.add_slide(_content_layout(prs))
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
    # find a body placeholder
    body = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0 and ph.has_text_frame:
            body = ph
            break
    if body is None:
        body = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(11), Inches(5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
    _set_notes(slide, notes)
    prs.save(out or path)
    return out or path


def add_section(path, title, subtitle="", notes="", out=None):
    prs = Presentation(path)
    slide = prs.slides.add_slide(_section_layout(prs))
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
    if subtitle:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != 0 and ph.has_text_frame:
                ph.text = subtitle
                break
    _set_notes(slide, notes)
    prs.save(out or path)
    return out or path


def clone_slide(dst_prs, src_slide):
    """Deep-copy a slide from another presentation into dst_prs, correctly
    importing embedded images as fresh parts (rewriting r:embed references) so
    there are no media partname collisions across decks."""
    import io
    from pptx.oxml.ns import qn
    blank = _layout_by_name(dst_prs, ["Full Page Blank", "Blank"]) or dst_prs.slide_layouts[-1]
    dst = dst_prs.slides.add_slide(blank)
    for ph in list(dst.shapes):               # start from a clean slide
        ph._element.getparent().remove(ph._element)
    for shp in src_slide.shapes:
        el = copy.deepcopy(shp._element)
        for blip in el.findall(".//" + qn("a:blip")):
            embed = blip.get(qn("r:embed"))
            if not embed:
                continue
            try:
                src_img = src_slide.part.related_part(embed)
                image_part, new_rid = dst.part.get_or_add_image_part(io.BytesIO(src_img.blob))
                blip.set(qn("r:embed"), new_rid)
            except Exception:
                # if an image can't be imported, drop the ref rather than corrupt
                blip.getparent().getparent().getparent().remove(
                    blip.getparent().getparent())
        dst.shapes._spTree.append(el)
    if src_slide.has_notes_slide:
        dst.notes_slide.notes_text_frame.text = src_slide.notes_slide.notes_text_frame.text
    return dst


# --------------------------------------------------------------------------- #
# In-memory assembly: keep selected slides, append new ones, reorder to plan.
# --------------------------------------------------------------------------- #
def open_deck(path):
    return Presentation(path)


def keep_only(prs, keep_1based):
    n = len(list(prs.slides._sldIdLst))
    for idx0 in sorted({i for i in range(n) if (i + 1) not in set(keep_1based)}, reverse=True):
        _drop_slide(prs, idx0)
    return prs


def append_content(prs, title, bullets, notes=""):
    slide = prs.slides.add_slide(_content_layout(prs))
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
    body = _body_placeholder(slide)
    if body is None:
        body = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(11.5), Inches(5.2))
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
    _set_notes(slide, notes)
    return slide


def append_section(prs, title, subtitle="", notes=""):
    slide = prs.slides.add_slide(_section_layout(prs))
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
    if subtitle:
        body = _body_placeholder(slide)
        if body is not None:
            body.text_frame.text = subtitle
    _set_notes(slide, notes)
    return slide


def retitle(prs, index0, text):
    slide = list(prs.slides)[index0]
    if slide.shapes.title is not None:
        slide.shapes.title.text = text


def reorder(prs, order0):
    """Reorder slides so the sequence becomes [current index order0[0], order0[1], ...]."""
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    assert sorted(order0) == list(range(len(ids))), "order must be a permutation of all slides"
    for el in ids:
        sldIdLst.remove(el)
    for i in order0:
        sldIdLst.append(ids[i])


def arrange(prs, order0):
    """Keep only the slides at current 0-based indices in `order0`, in that exact
    order. Slides not referenced are dropped. Safe to call after appending new
    slides (append-before-delete pattern avoids partname collisions)."""
    from pptx.oxml.ns import qn
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    chosen = [ids[i] for i in order0]
    chosen_set = set(id(e) for e in chosen)
    for e in ids:
        if id(e) not in chosen_set:
            rId = e.get(qn("r:id"))
            prs.part.drop_rel(rId)
            sldIdLst.remove(e)
        else:
            sldIdLst.remove(e)
    for e in chosen:
        sldIdLst.append(e)


def save(prs, path):
    prs.save(path)
    return path


if __name__ == "__main__":
    import sys
    if sys.argv[1] == "inv":
        for idx, title, w in inventory(sys.argv[2]):
            print(f"{idx:>3}  {w:>4}w  {title}")
