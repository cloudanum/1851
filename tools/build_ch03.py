#!/usr/bin/env python3
"""build_ch03.py — assemble the A.2 'Machine Learning Essentials' chapter.

Carrier = _bak0/1851-Ch03.pptx (restored fresh every run -> idempotent).
Keeps selected carrier slides, clones selected slides from _bak0 Ch04/05/07
(cross-deck copy re-embedding images via pptx_tools.clone_slide), authors 4
new GenAI-bridge slides, then arranges everything into the target order.

Run:  /Users/iahmad/Creator/Courses_and_conferences/LT/.venv-courseware/bin/python tools/build_ch03.py
"""
from __future__ import annotations
import copy
import shutil
import sys
import zipfile
from collections import Counter
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import slidekit as sk  # noqa: F401  (import activates the partname-collision fix)
import pptx_tools as pt
from pptx import Presentation
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
BAK = ROOT / "decks" / "_bak0"
OUT = ROOT / "decks" / "1851-Ch03.pptx"

# --- keep lists (1-based, against _bak0 decks) -------------------------------
KEEP3 = [1, 2, 3, 4, 5, 6, 7, 8, 9, 10, 11, 12, 13, 15, 17, 20, 21, 22,
         25, 26, 28, 37, 39, 50, 51, 53, 54]                       # 27 carrier
KEEP4 = [4, 5, 21, 23, 25, 26, 27, 28, 29, 30, 32, 37]            # 12 data prep
KEEP5 = [5, 7, 12, 16, 17, 19, 22, 23]                            #  8 metrics
KEEP7 = [6, 10, 13, 14, 18, 19, 23]                               #  7 neural nets

# --- rewritten text for kept carrier slides ----------------------------------
OBJECTIVES = [
    "Explain how ML differs from traditional programming and distinguish supervised from unsupervised learning",
    "Apply core techniques—linear regression, decision trees, and K-Means clustering",
    "Prepare data for ML: quality checks, cleaning, encoding, scaling, and leakage-free splits",
    "Select and interpret evaluation metrics—confusion matrix, precision, recall, ROC-AUC, MAE/RMSE/R²",
    "Describe how neural networks learn and how to control overfitting",
    "Connect classic ML foundations to transformers, LLMs, and GenAI applications",
]
AGENDA = [
    "ML Core Concepts and Supervised Learning",
    "Regression and Decision Trees",
    "Unsupervised Learning and K-Means Clustering",
    "From Tokens to Embeddings: NLP Foundations",
    "Data Preparation and Quality",
    "Model Evaluation Metrics",
    "Neural Networks and Deep Learning",
    "From Classic ML to Generative AI",
]
SUMMARY = [
    "ML learns patterns from data rather than following explicit rules—supervised and unsupervised learning cover most practical problems",
    "Regression, decision trees, and K-Means remain workhorse techniques, especially on tabular data",
    "Model quality is data quality: cleaning, encoding, scaling, and leakage-free splits determine success",
    "Match the metric to the risk—precision, recall, and ROC-AUC for classifiers; MAE, RMSE, and R² for regression",
    "Neural networks learn via forward and backpropagation; validation and regularization control overfitting",
    "These foundations carry directly into transformers and LLMs—and into RAG data quality and evals later in this course",
]

# --- 4 new bridge slides (title, bullets, speaker notes) ---------------------
NEW_SLIDES = [
    ("From Neural Networks to Transformers",
     ["The RNN-era bottleneck: tokens processed one at a time—limited context, hard to parallelize",
      "Attention intuition: let every token look at every other token and weigh what matters",
      "Self-attention is computed in parallel—training scales with data and compute",
      "Transformer = stacked self-attention + feed-forward layers; one architecture for many tasks",
      "Embeddings become contextual: a token's vector now depends on the whole sentence"],
     "Bridge from the embeddings slides to modern NLP: embeddings are still the input, but transformers "
     "make them contextual. The key intuition is weighted relevance—each word decides which other words "
     "to listen to, all at once. Emphasize why this scales: no sequential bottleneck, so training on "
     "internet-scale text became practical."),
    ("From Transformers to LLMs",
     ["Step 1 — Pretraining: predict the next token on massive text corpora",
      "Step 2 — Instruction tuning: fine-tune on (instruction, response) pairs",
      "Step 3 — RLHF: human feedback shapes helpful, harmless behavior",
      "Same next-token engine at every stage—only the data and objectives change",
      "Next chapter: working with these models in practice—prompting, APIs, and RAG"],
     "Give the one-glance pipeline: pretrain, instruct, RLHF. Stress that all three stages use the same "
     "transformer architecture just covered—scale and data create the emergent capability. Demystify the "
     "LLM as a very large, very well-tuned next-token predictor, and flag that Chapter 4 goes hands-on."),
    ("Where Classic ML Still Wins",
     ["Tabular and small datasets: gradient-boosted trees and regression often beat LLMs",
      "Latency and cost: milliseconds and millicents vs. per-token pricing",
      "Determinism and auditability: reproducible scores, easier validation and sign-off",
      "Data privacy: train and run inside your own boundary—no external API",
      "Engineering guidance: start with the simplest model that meets the requirement"],
     "Be honest with engineers: GenAI is not always the right tool. Fraud scoring, churn prediction, and "
     "risk ranking on structured data are still classic-ML territory. Rule of thumb: reach for an LLM "
     "when the input or output is language; otherwise benchmark classic ML first."),
    ("ML Essentials in a GenAI World",
     ["Data preparation → data quality for RAG: chunking, freshness, and source trust",
      "Train/validation/test discipline → evals: golden sets and regression checks (Chapter 6)",
      "Confusion matrix, precision/recall → LLM quality and safety metrics",
      "Neural network training loop → how LLMs are trained and fine-tuned",
      "Interpretability and bias-variance → grounding, guardrails, and hallucination control"],
     "Map old to new so nothing in this chapter is wasted. Split discipline becomes the eval harnesses of "
     "Chapter 6; data-prep instincts become RAG data-quality practice; classic metrics thinking becomes "
     "LLM-eval thinking. The class now has the vocabulary the rest of the course builds on."),
]


# --- formatting-preserving text rewrite helpers ------------------------------
def _set_para_text(p_el, text):
    """Replace a paragraph's text, keeping the first run's formatting."""
    runs = p_el.findall(qn("a:r"))
    if not runs:
        r = p_el.makeelement(qn("a:r"), {})
        r.append(p_el.makeelement(qn("a:t"), {}))
        p_el.append(r)
        runs = [r]
    first = runs[0]
    for extra in runs[1:]:
        p_el.remove(extra)
    for br in p_el.findall(qn("a:br")):
        p_el.remove(br)
    t = first.find(qn("a:t"))
    if t is None:
        t = first.makeelement(qn("a:t"), {})
        first.append(t)
    t.text = text


def rewrite_bullets(shape, lines):
    """Rewrite a body placeholder's bullets, reusing paragraph formatting."""
    txBody = shape.text_frame._txBody
    paras = txBody.findall(qn("a:p"))
    for i, line in enumerate(lines):
        if i < len(paras):
            _set_para_text(paras[i], line)
        else:
            newp = copy.deepcopy(paras[-1])
            txBody.append(newp)
            paras.append(newp)
            _set_para_text(newp, line)
    for p in paras[len(lines):]:
        txBody.remove(p)


def main():
    # 0. fresh carrier -> idempotent
    shutil.copy(BAK / "1851-Ch03.pptx", OUT)
    prs = sk.open_prs(OUT)

    # 1. clone keeps from Ch04/05/07 (appended at end), record new indices
    clone_idx = {}
    for ch, keeps in [(4, KEEP4), (5, KEEP5), (7, KEEP7)]:
        src = pt.open_deck(BAK / f"1851-Ch{ch:02d}.pptx")
        for k in keeps:
            pt.clone_slide(prs, src.slides[k - 1])
            clone_idx[(ch, k)] = len(prs.slides) - 1

    # 2. author the 4 new bridge slides (appended at end)
    new_idx = []
    for title, bullets, notes in NEW_SLIDES:
        pt.append_content(prs, title, bullets, notes)
        new_idx.append(len(prs.slides) - 1)

    # 3. edit kept carrier slides in place (0-based indices still original)
    s1 = prs.slides[0]
    _set_para_text(s1.shapes.title.text_frame._txBody.findall(qn("a:p"))[0],
                   "Machine Learning Essentials")
    rewrite_bullets(pt._body_placeholder(prs.slides[1]), OBJECTIVES)   # Objectives
    rewrite_bullets(pt._body_placeholder(prs.slides[2]), AGENDA)       # Agenda
    rewrite_bullets(pt._body_placeholder(prs.slides[52]), SUMMARY)     # Summary
    rewrite_bullets(pt._body_placeholder(prs.slides[53]), OBJECTIVES)  # recap

    # 4. arrange: openers -> ML core -> data prep -> metrics -> NN -> bridge -> close
    order = (
        [0, 1, 2]                                                       # openers
        + [3, 4, 5, 6, 8, 9, 10, 7, 11, 12, 14, 16, 19]                 # supervised core
        + [20, 49, 50]                                                  # unsupervised
        + [21, 24, 25, 27]                                              # NLP -> embeddings
        + [clone_idx[(4, k)] for k in KEEP4]                            # data prep
        + [clone_idx[(5, k)] for k in KEEP5]                            # metrics
        + [36, 38]                                                      # deep-learning intro
        + [clone_idx[(7, k)] for k in KEEP7]                            # neural nets
        + new_idx                                                       # 4 bridge slides
        + [52, 53]                                                      # summary + recap
    )
    pt.arrange(prs, order)
    pt.save(prs, OUT)

    # 5. verify: reopen, titles, blanks, duplicate partnames, zip integrity
    prs2 = Presentation(OUT)
    n = len(prs2.slides)
    print(f"\n{OUT.name}: {n} slides")
    print(f"kept: Ch03={len(KEEP3)} Ch04={len(KEEP4)} Ch05={len(KEEP5)} "
          f"Ch07={len(KEEP7)} new={len(NEW_SLIDES)}")
    empty = []
    for i, s in enumerate(prs2.slides, 1):
        t = sk.slide_title(s)
        if not t.strip():
            empty.append(i)
        print(f"{i:>3} | {t}")
    assert not empty, f"empty titles at {empty}"
    assert 55 <= n <= 64, f"slide count {n} outside 55-64"
    with zipfile.ZipFile(OUT) as z:
        dupes = [p for p, c in Counter(z.namelist()).items() if c > 1]
        assert not dupes, f"duplicate zip partnames: {dupes}"
        assert z.testzip() is None, "zip integrity check failed"
    print("\nOK: no empty titles, no duplicate partnames, zip integrity verified")


if __name__ == "__main__":
    main()
