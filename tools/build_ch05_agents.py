#!/usr/bin/env python3
"""Enrich 1851-Ch05.pptx with content + figures from the Software Development Agents chapter."""
import copy, shutil
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = "/Users/iahmad/Creator/Courses_and_conferences/LT/courses/1851/1851-FTP-upload"
SRC = f"{BASE}/1851-Ch05.pptx"
BAK = f"{BASE}/1851-Ch05.backup.pptx"
A = f"{BASE}/ch05_assets"

# ---------- downscale oversized figures (keep file size sane) ----------
def prep(name, max_w=1800):
    path = f"{A}/{name}"
    im = Image.open(path)
    if im.width > max_w:
        im = im.resize((max_w, round(im.height * max_w / im.width)), Image.LANCZOS)
        out = path.replace(".png", "_s.png")
        im.save(out)
        return out
    return path

FIG91 = prep("p03_img0.png")   # ecosystem layers      2062x1328
FIG92 = prep("p05_img0.png")   # adoption curve        2096x1198
FIG93 = prep("p06_img0.png")   # TDD three phases      4139x1003
FIG94 = prep("p09_img0.png")   # LangGraph loop        5238x8313 (tall)
FIG95 = prep("p33_img0.png")   # self-improvement loop 1825x3502 (tall)

shutil.copy2(SRC, BAK)
p = Presentation(SRC)
layouts = {l.name: l for l in p.slide_layouts}
CONTENT = layouts["Content with Header. Full Page"]
DISCUSS = layouts["Discussion"]

GRAY = RGBColor(0x59, 0x59, 0x59)

def set_title(slide, text):
    slide.shapes.title.text_frame.text = text

def get_ph(slide, idx):
    for sh in slide.placeholders:
        if sh.placeholder_format.idx == idx:
            return sh
    return None

def fill_bullets(slide, bullets, idx=1):
    ph = get_ph(slide, idx)
    tf = ph.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = b
    return ph

def add_picture_centered(slide, path, top_in, width_in):
    im = Image.open(path)
    h_in = width_in * im.height / im.width
    left_in = (10.0 - width_in) / 2
    slide.shapes.add_picture(path, Inches(left_in), Inches(top_in), width=Inches(width_in))
    return top_in + h_in  # bottom edge

def add_picture_at(slide, path, left_in, top_in, width_in):
    im = Image.open(path)
    h_in = width_in * im.height / im.width
    slide.shapes.add_picture(path, Inches(left_in), Inches(top_in), width=Inches(width_in))
    return top_in + h_in

def add_caption(slide, text, left_in, top_in, width_in):
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    r = para.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(11)
    r.font.italic = True
    r.font.color.rgb = GRAY

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

new_slides = []

# ---- N1: three classes -------------------------------------------------
s = p.slides.add_slide(CONTENT)
set_title(s, "Three Classes of Software Development Agents")
fill_bullets(s, [
    "Code-Generation agents — turn natural-language specs into verified code: generate \u2192 test \u2192 refine",
    "Compliance-Driven agents — enforce security and policy as code is written: scan \u2192 evaluate \u2192 remediate",
    "Self-Improving agents — learn from outcomes and feedback: execute \u2192 observe \u2192 learn \u2192 adapt",
    "Same foundation for all three: specialized roles, structured feedback loops, human checkpoints",
    "The models are Chapter 4's models — what changes is the loop wrapped around them",
])
add_notes(s, "This is the map for the next ten slides. Each agent class is defined by its feedback loop, and each loop reuses a discipline the students already know: TDD for generation, policy-as-code for compliance, MLOps for self-improvement. Stress the last bullet: no new model magic — the architecture around the model is what creates reliability.")
new_slides.append(s)

# ---- N2: where agents operate ------------------------------------------
s = p.slides.add_slide(CONTENT)
set_title(s, "Where Development Agents Operate")
fill_bullets(s, [
    "IDE copilots — the inner loop: completions and function drafting in the editor (Copilot, Cursor)",
    "Pipeline agents — the outer loop: synthesize tests, repair broken builds, propose patches in CI",
    "Repository-aware assistants — index the whole codebase; multi-file edits that respect project patterns",
    "Grounding in real repository symbols is what stops invented APIs and phantom dependencies",
    "In production: Devin runs multi-step engineering tasks, SWE-agent resolves real GitHub issues, Snyk and Semgrep add semantic security analysis",
])
add_notes(s, "Connect back to the Levels of AI Assistance slide: copilots are level 1-2, pipeline and repo-aware agents are level 3. The hallucination traps slide named invented APIs as a top failure — repository grounding is the industry's answer to exactly that problem.")
new_slides.append(s)

# ---- N3: tooling stack + Fig 9.1 ---------------------------------------
s = p.slides.add_slide(CONTENT)
set_title(s, "The Agent Tooling Stack")
fill_bullets(s, [
    "Orchestration frameworks (LangGraph, LangChain) define the stateful workflow; the reasoning core is an LLM with retrieval and tools",
    "Quality gates and observability are what turn autonomy into something you can measure, audit, and trust",
])
bottom = add_picture_centered(s, FIG91, 2.55, 6.2)
add_caption(s, "Four functional layers of the software-agent ecosystem", 1.5, bottom + 0.05, 7.0)
add_notes(s, "Walk the four layers top to bottom. Students already know layers 1-2 from Chapter 4's API work; the new idea is that layers 3-4 — gates and observability — are not optional extras but the difference between a demo and a production system.")
new_slides.append(s)

# ---- N4: adoption curve + Fig 9.2 --------------------------------------
s = p.slides.add_slide(CONTENT)
set_title(s, "The Adoption Maturity Curve")
fill_bullets(s, [
    "Teams start with low-risk drafting and refactoring, then add test synthesis and quality gates in CI",
    "Mature stage: conditional autonomy — agents open PRs and automate merges; humans keep the architectural checkpoints",
    "Rule of thumb: the agent proposes; the developer disposes",
])
bottom = add_picture_centered(s, FIG92, 2.9, 6.9)
add_caption(s, "Adoption maturity curve for AI coding agents", 1.5, bottom + 0.05, 7.0)
add_notes(s, "This curve mirrors how teams adopted CI itself: assist first, automate later, gate everything. Conditional autonomy is the sweet spot for most organizations today — full merge authority stays with humans, everything before the merge can be agent-run.")
new_slides.append(s)

# ---- N5: TDG + Fig 9.3 --------------------------------------------------
s = p.slides.add_slide(CONTENT)
set_title(s, "Test-Driven Generation: TDD, Run by Agents")
fill_bullets(s, [
    "Red — a tester agent turns the requirement into a failing test suite: the executable spec",
    "Green — a developer agent writes the minimal code to pass; test output returns as structured feedback",
    "Refactor — clean up with tests green; any failure routes straight back to implementation",
    "The AI-pair TDD workflow from earlier in this chapter — with the whole loop automated",
])
bottom = add_picture_centered(s, FIG93, 3.9, 9.3)
add_caption(s, "Three phases of the agentic TDD loop", 1.5, bottom + 0.05, 7.0)
add_notes(s, "Point at the earlier slides 'TDD with an AI Pair': there the human wrote tests and the AI implemented. Here both roles are agents — distinct system prompts and tool sets, often the same underlying model. The governing signal is unchanged: the test suite decides when the work is done.")
new_slides.append(s)

# ---- N6: generate-test-refine + Fig 9.4 (side by side) ------------------
s = p.slides.add_slide(CONTENT)
set_title(s, "Inside the Generate\u2013Test\u2013Refine Loop")
ph = fill_bullets(s, [
    "An orchestrator decomposes the request and routes tasks to specialized agents",
    "Code and tests execute in a Docker sandbox; stack traces feed the next iteration's prompt",
    "State persists across iterations — task, code, tests, failure history (LangGraph checkpointing)",
    "Iteration caps stop infinite loops; the failure trace doubles as audit log and training data",
])
ph.width = Inches(5.9)
bottom = add_picture_at(s, FIG94, 6.55, 1.15, 3.3)
add_caption(s, "LangGraph workflow: the iterative generate\u2013test\u2013refine loop", 6.35, bottom + 0.03, 3.6)
add_notes(s, "Trace one iteration on the diagram: orchestrator assigns the task, developer agent writes code, tester agent writes tests, the sandbox runs them, and a failure routes the full stack trace back into the developer agent's next prompt. Emphasize that progress is impossible until all tests pass — the test runner, not the model's confidence, is the judge.")
new_slides.append(s)

# ---- N7: multi-agent feature team ---------------------------------------
s = p.slides.add_slide(CONTENT)
set_title(s, "Scaling Up: A Multi-Agent Feature Team")
fill_bullets(s, [
    "A project-manager agent decomposes the feature, tracks dependencies, and fans out independent tasks",
    "Backend agent (Python/Flask + pytest) and frontend agent (TypeScript/React + Jest) run the same loop in parallel",
    "A tester agent writes framework-appropriate tests per layer; an integration stage validates cross-layer contracts",
    "One model, many roles — specialization comes from prompts and tool sets, not different models",
    "Measured run: about 1.4 iterations per task, 100% coverage by construction — code cannot advance untested",
])
add_notes(s, "The scaling story is decomposition, not bigger prompts. Each specialized agent inherits the same generate-test-refine loop from the previous slide; the PM agent adds dependency ordering and parallel fan-out. The 100% coverage claim is structural: the workflow has no path that skips the test gate.")
new_slides.append(s)

# ---- N8: compliance agents ----------------------------------------------
s = p.slides.add_slide(CONTENT)
set_title(s, "Compliance-Driven Agents: Policy as Code")
fill_bullets(s, [
    "Tests ask 'does it work?'; compliance agents ask 'is it permissible?' — GDPR, PCI DSS, HIPAA, internal policy",
    "The loop runs in CI on every pull request: scan \u2192 evaluate \u2192 remediate",
    "Policy engines (OPA, Rego) are the test suite for rules; the LLM translates violations into developer-friendly fixes",
    "Semantic analysis catches what pattern matching misses — an 'anonymize' function that keeps emails and IPs",
    "Every flag, fix, and override is logged: the audit trail is built as you work",
])
add_notes(s, "Compliance is an orthogonal concern to correctness: a function can pass every unit test and still violate a data-retention policy. The key design move is treating policy as executable code — versioned, tested, and enforced in the same pipeline as functional tests. The anonymize example shows why pure pattern matching is not enough.")
new_slides.append(s)

# ---- N9: PCI DSS case study ---------------------------------------------
s = p.slides.add_slide(CONTENT)
set_title(s, "Case Study: PCI DSS Enforcement in CI")
fill_bullets(s, [
    "A fintech with dozens of teams shipping daily — manual security review became the bottleneck",
    "The agent fires on every PR: OPA policies plus static analysis over the diff, results in seconds",
    "Hard fail: clear violations block the merge, with specific remediation guidance posted on the PR",
    "Soft fail: ambiguous cases get a non-blocking comment asking for confirmation — a human checkpoint",
    "Six months in: violations down 85%; the annual PCI audit became a log review, not a scramble",
])
add_notes(s, "Note the two-tier enforcement: hard fails for certain, high-severity violations; soft fails that keep velocity while still creating an auditable human decision. The 85% drop is a learning effect — immediate, specific feedback trained the developers, not just the pipeline.")
new_slides.append(s)

# ---- N10: self-improving + Fig 9.5 (side by side) -----------------------
s = p.slides.add_slide(CONTENT)
set_title(s, "Self-Improving Agents: The Learning Loop")
ph = fill_bullets(s, [
    "A sensing layer gathers explicit (ratings, corrections), implicit (iterations, acceptance rates), and synthetic (benchmark) feedback",
    "A critic agent scores outcomes against KPIs; a planner turns patterns into improvement hypotheses",
    "A human checkpoint gates risky changes; approved ones update prompts, thresholds, retrieval, or weights",
    "Deploy & test validates against baselines before production — MLOps discipline for agent behavior",
])
ph.width = Inches(5.9)
bottom = add_picture_at(s, FIG95, 6.85, 1.1, 2.85)
add_caption(s, "The self-improvement closed loop", 6.6, bottom + 0.03, 3.3)
add_notes(s, "This is the same control-loop thinking as the generate-test-refine cycle, one level up: the thing being tested and refined is the agent's own behavior. The human-in-the-loop checkpoint is the load-bearing piece — low-risk prompt tweaks can auto-deploy, but anything touching models or policies waits for approval.")
new_slides.append(s)

# ---- N11: governing self-improvement ------------------------------------
s = p.slides.add_slide(CONTENT)
set_title(s, "Governing Self-Improvement")
fill_bullets(s, [
    "Risk-tiered updates: prompt tweaks auto-deploy; threshold changes get human review; fine-tuning needs cross-functional sign-off",
    "Immutable version history with instant rollback — every change linked to the feedback that motivated it",
    "Bias monitoring watches for behavior drift across user groups and flags it for review",
    "Beware over-personalization: diversity constraints and held-out evals keep the agent general",
    "The chapter's invariant still holds: autonomy scales only as fast as verification",
])
add_notes(s, "Governance closes the chapter's argument. Every safeguard here is a familiar engineering control — versioning, rollback, monitoring, tiered approval — reapplied to a system that changes itself. Land on the invariant: it is the same rule the quality-gates slide states for human-supervised GenAI.")
new_slides.append(s)

# ---- N12: discussion -----------------------------------------------------
s = p.slides.add_slide(DISCUSS)
set_title(s, "Do Now: How Much Autonomy?")
fill_bullets(s, [
    "Pick one real task from your current project or coursework.",
    "Which maturity stage fits it today: copilot, assistant, or conditional autonomy?",
    "What test, policy, or review gate would you demand before moving one stage up?",
    "Be ready to defend your gate: what failure does it catch?",
    "Time box: 6 minutes — 4 in pairs, 2 to hear two pairs defend theirs.",
])
add_notes(s, "This mirrors the earlier 'Where Would You Trust AI?' discussion, one autonomy level higher. Push students to name a concrete gate — a test suite, a policy rule, a required reviewer — rather than a vague comfort level. The gates they name are the maturity curve in action.")
new_slides.append(s)

# ---------- edits to existing slides ----------
def replace_bullets(slide, texts, idx):
    ph = get_ph(slide, idx)
    tf = ph.text_frame
    tf.clear()
    for i, t in enumerate(texts):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = t

# Agenda (slide 3): add agents line after CI/CD line
s3 = p.slides[2]
ph = get_ph(s3, 1)
agenda = ["".join(r.text for r in para.runs) for para in ph.text_frame.paragraphs]
pos = agenda.index("Review, documentation, CI/CD, and operations")
agenda.insert(pos + 1, "Software development agents: generate\u2013test\u2013refine, compliance, and self-improvement")
replace_bullets(s3, agenda, 1)

# Objectives (slides 2 and 52): append agent objective
OBJ = "Explain how code-generation, compliance, and self-improving agents close the loop — and where humans stay in charge"
for i in (1, 51):
    s_obj = p.slides[i]
    ph = get_ph(s_obj, 15)
    objs = ["".join(r.text for r in para.runs) for para in ph.text_frame.paragraphs]
    if OBJ not in objs:
        objs.append(OBJ)
    replace_bullets(s_obj, objs, 15)

# Slide 46: retarget the forward pointer to the new section
s46 = p.slides[45]
ph = get_ph(s46, 1)
for para in ph.text_frame.paragraphs:
    txt = "".join(r.text for r in para.runs)
    if txt.startswith("Chapter 7 builds these hands-on"):
        para.runs[0].text = "The next slides map the three agent classes; Chapter 7 builds them hands-on"
        for r in para.runs[1:]:
            r.text = ""

# Summary (slide 51): add agent takeaway
s51 = p.slides[50]
ph = get_ph(s51, 1)
summ = ["".join(r.text for r in para.runs) for para in ph.text_frame.paragraphs]
NEW = "Agents close the loop: generate\u2013test\u2013refine for code, scan\u2013evaluate\u2013remediate for compliance, observe\u2013learn\u2013adapt for improvement — human checkpoints throughout"
if NEW not in summ:
    pos = summ.index("Evals (Chapter 6) are TDD for AI behavior; agents (Chapter 7) automate the workflow toil")
    summ.insert(pos + 1, NEW)
replace_bullets(s51, summ, 1)

# ---------- move new slides to just after slide 46 ----------
sldIdLst = p.slides._sldIdLst
ids = list(sldIdLst)
new_ids = ids[-len(new_slides):]
for el in new_ids:
    sldIdLst.remove(el)
for offset, el in enumerate(new_ids):
    sldIdLst.insert(46 + offset, el)

p.save(SRC)
print(f"Saved. Total slides: {len(Presentation(SRC).slides.__iter__.__self__._sldIdLst)}")
print("Backup at:", BAK)
