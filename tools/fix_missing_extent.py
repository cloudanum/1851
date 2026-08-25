#!/usr/bin/env python3
"""Repair text placeholders with missing/zero extent (invisible body text).

A build idiom in the deck pipeline left some body placeholders with no
<a:ext> element at all (python-pptx: width is None) or cx=0. They render
as title-only slides. This pass restores body geometry; text is never edited.

Geometry rule per slide:
  - right-hand picture placeholder present -> left column
  - otherwise -> full-width body
Both match the decks' healthy reference slides.

Idempotent: shapes already at valid geometry are untouched.
Acceptance check is built in: after the pass, any text-bearing shape still
lacking extent is reported and the exit code is 1.

Usage: fix_missing_extent.py <deck.pptx> [...]
"""
import sys

from pptx import Presentation
from pptx.util import Inches, Emu

FULL = (0.29, 1.04, 9.51, 5.83)     # L, T, W, H in inches
LEFT_COL = (0.29, 1.04, 5.05, 5.83)
FOOTER_Y = 6.5                       # shapes starting below this are footers


def is_broken(sh):
    return sh.has_text_frame and sh.text_frame.text.strip() and (sh.width or 0) < 10000


def main(paths):
    rc = 0
    for path in paths:
        prs = Presentation(path)
        fixed = 0
        for slide in prs.slides:
            right_pic = any(
                'Picture' in sh.name and (sh.left or 0) > Inches(5) and (sh.width or 0) > Inches(2)
                for sh in slide.shapes)
            for sh in slide.shapes:
                if not is_broken(sh):
                    continue
                if sh.top is not None and sh.top > Inches(FOOTER_Y):
                    continue  # page-number footer placeholders: leave alone
                l, t, w, h = LEFT_COL if right_pic else FULL
                sh.left, sh.top, sh.width, sh.height = (Inches(l), Inches(t),
                                                        Inches(w), Inches(h))
                fixed += 1
        prs.save(path)
        # acceptance: nothing broken may remain
        prs2 = Presentation(path)
        remaining = sum(1 for s in prs2.slides for sh in s.shapes
                        if is_broken(sh) and not (sh.top is not None and sh.top > Inches(FOOTER_Y)))
        status = 'OK' if remaining == 0 else f'FAIL ({remaining} remain)'
        if remaining:
            rc = 1
        print(f'{path.split("/")[-1]}: {fixed} repaired, acceptance {status}')
    sys.exit(rc)


if __name__ == '__main__':
    main(sys.argv[1:])
