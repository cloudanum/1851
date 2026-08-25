#!/usr/bin/env python3
"""build_engagement_ch05.py — ENGAGEMENT pass on 1851-Ch05 (A.2, 50 -> 52 slides).

Restores the A.2 deck fresh from decks/_bak1 every run (idempotent), then:
  1. inserts 'Do Now: Where Would You Trust AI?' in the framing block, right
     after 'What Changed by 2026' (layout 'Discussion'),
  2. inserts 'Do Now: Red, Green, Refactor' at the start of the Testing & TDD
     section, right before 'Why Tests Matter More with AI-Generated Code'
     (layout 'Do Now Writing Offline').

Lab 5.1 already has its exercise slide — no lab-slide work here.

Run:  /Users/iahmad/Creator/Courses_and_conferences/LT/.venv-courseware/bin/python tools/build_engagement_ch05.py
"""
from __future__ import annotations
import shutil
import sys
import zipfile
from collections import Counter
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import slidekit as sk  # noqa: F401  (import activates the partname-collision fix)
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent
BAK1 = ROOT / "decks" / "_bak1"
OUT = ROOT / "decks" / "1851-Ch05.pptx"

# --- new slide content -------------------------------------------------------
DONOW_TRUST = ("Do Now: Where Would You Trust AI?", "Discussion",
    ["Look back at the GenAI-augmented SDLC phase map.",
     "Pick ONE phase where you would trust AI assistance today.",
     "Pick ONE phase where you would not — yet.",
     "Be ready to defend both picks: what makes them different?",
     "Time box: 6 minutes — 4 in pairs, 2 to hear two pairs defend theirs."],
    "Run it against the phase map a few slides back (requirements, design, implementation, testing, "
    "release, operations). There is no wrong answer; push past 'it depends' by asking what evidence "
    "would change their mind. Debrief: testing and documentation usually get trusted first because the "
    "output is checkable, while requirements and architecture lag because they need business context "
    "the model does not have — which is exactly this chapter's theme.")

DONOW_TDD = ("Do Now: Red, Green, Refactor", "Do Now Writing Offline",
    ["Number these six TDD steps 1–6 in the order they happen:",
     "__ Commit, then pick the next behavior and repeat",
     "__ Write the simplest code that passes the test",
     "__ Refactor — clean up while keeping every test green",
     "__ Run the new test and watch it fail",
     "__ Write one failing test for the next small behavior",
     "__ Run the whole suite and confirm it is green",
     "Time box: 5 minutes, pen and paper — laptops closed."],
    "Correct order: 1) write one failing test for the next small behavior, 2) run it and watch it fail "
    "(red), 3) write the simplest code that passes (green), 4) run the whole suite and confirm green, "
    "5) refactor while keeping every test green, 6) commit and repeat — so the slide rows top to bottom "
    "are 6, 3, 5, 2, 1, 4. Debrief: ask why we watch the test fail instead of assuming it fails — a test "
    "that never fails proves nothing, and that habit matters even more when an AI pair writes the "
    "implementation. This sets up the TDD-with-AI patterns in the next slides.")


# --- helpers ------------------------------------------------------------------
def pick_layout_exact(prs, name):
    """Exact layout-name match (pick_layout's substring match would return the
    'Numbered …' variant first)."""
    for master in prs.slide_masters:
        for lay in master.slide_layouts:
            if (lay.name or "") == name:
                return lay
    raise AssertionError(f"layout not found: {name!r}")


def add_donow_slide(prs, title, layout_name, bullets, notes):
    """Add a Do Now / Discussion slide on the named template layout (graphic on
    the left, text column on the right). Returns the new slide."""
    lay = pick_layout_exact(prs, layout_name)
    slide = prs.slides.add_slide(lay)
    sk.set_title(slide, title)
    body = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1 and ph.has_text_frame:
            body = ph
            break
    assert body is not None, f"content placeholder not found on {layout_name!r}"
    # the layout's content placeholder is stubby; give it room for 5+ lines
    body.height = Inches(4.5)
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
    sk.set_notes(slide, notes)
    return slide


def find_slide(prs, title_sub):
    """0-based index of the unique slide whose title contains title_sub."""
    hits = [i for i, s in enumerate(prs.slides) if title_sub in sk.slide_title(s)]
    assert len(hits) == 1, f"anchor {title_sub!r}: expected 1 hit, got {hits}"
    return hits[0]


def main():
    # 0. fresh restore from _bak1 -> idempotent
    shutil.copy(BAK1 / OUT.name, OUT)
    prs = sk.open_prs(OUT)
    assert len(prs.slides) == 50, f"expected 50-slide A.2 deck, got {len(prs.slides)}"

    # 1+2. append each new slide and immediately move it into place — descending
    # target index, so the second move cannot disturb the first placement.
    add_donow_slide(prs, DONOW_TDD[0], DONOW_TDD[1], DONOW_TDD[2], DONOW_TDD[3])
    at_tdd = find_slide(prs, "Why Tests Matter More with AI-Generated Code")  # before it
    sk.move_slide(prs, len(prs.slides) - 1, at_tdd)

    add_donow_slide(prs, DONOW_TRUST[0], DONOW_TRUST[1], DONOW_TRUST[2], DONOW_TRUST[3])
    at_trust = find_slide(prs, "What Changed by 2026") + 1                    # after it
    sk.move_slide(prs, len(prs.slides) - 1, at_trust)

    print(f"inserted at 1-based: {at_trust+1} (Where Would You Trust AI?), "
          f"{at_tdd+1} (Red, Green, Refactor)")

    sk.save(prs, OUT)

    # 3. verify: reopen, count, titles, blanks, duplicate partnames, zip integrity
    prs2 = Presentation(OUT)
    n = len(prs2.slides)
    print(f"\n{OUT.name}: {n} slides")
    empty = []
    for i, s in enumerate(prs2.slides, 1):
        t = sk.slide_title(s)
        if not t.strip():
            empty.append(i)
        print(f"{i:>3} | {t}")
    assert not empty, f"empty titles at {empty}"
    assert n == 52, f"slide count {n} != 52"
    titles = [sk.slide_title(s) for s in prs2.slides]
    # positions are relative to neighbors, so verify the adjacency contract
    assert titles.index(DONOW_TRUST[0]) == titles.index("What Changed by 2026") + 1
    assert titles.index(DONOW_TDD[0]) == \
        titles.index("Why Tests Matter More with AI-Generated Code") - 1
    with zipfile.ZipFile(OUT) as z:
        dupes = [p for p, c in Counter(z.namelist()).items() if c > 1]
        assert not dupes, f"duplicate zip partnames: {dupes}"
        assert z.testzip() is None, "zip integrity check failed"
    print("\nOK: 52 slides, no empty titles, no duplicate partnames, zip integrity verified")


if __name__ == "__main__":
    main()
