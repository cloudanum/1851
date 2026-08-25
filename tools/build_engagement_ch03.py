#!/usr/bin/env python3
"""build_engagement_ch03.py — Day-1 engagement pass on the A.2 Ch03 deck.

Carrier = _bak1/1851-Ch03.pptx (restored fresh every run -> idempotent).
Adds two unnumbered Do Now warm-ups ('Rules, Classic ML, or GenAI?' right after
the 'ML vs. Traditional Programming' content; 'Draw the Split' inside the
data-splitting block, after the train/val/test slides), retitles the existing
'Activity 3.1: Decision Tree Classifier for Predicting Income' slide to
'Lab 3.1: Decision Trees vs. an LLM' and rewrites its body/notes in place
(same exercise layout), and updates the Agenda to mention Lab 3.1. 58 -> 60.

Run:  /Users/iahmad/Creator/Courses_and_conferences/LT/.venv-courseware/bin/python tools/build_engagement_ch03.py
"""
from __future__ import annotations
import copy
import shutil
import sys
import zipfile
from collections import Counter
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import slidekit as sk  # noqa: F401  (import activates the partname-collision fix)
import pptx_tools as pt
from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
BAK = ROOT / "decks" / "_bak1"
OUT = ROOT / "decks" / "1851-Ch03.pptx"

BODY_H = Inches(4.5)   # raise stubby layout-default body areas (house pattern)


def layout_named(prs, name):
    for lay in prs.slide_layouts:
        if (lay.name or "") == name:
            return lay
    raise KeyError(f"layout {name!r} not found")


def find_slide(prs, title_sub):
    """First slide index whose title contains title_sub (case-insensitive)."""
    sub = title_sub.lower()
    for i, s in enumerate(prs.slides):
        if sub in sk.slide_title(s).lower():
            return i
    raise KeyError(f"slide titled like {title_sub!r} not found")


def add_engagement_slide(prs, layout_name, title, bullets, notes, at_index):
    """Add a slide on a named Do Now / Discussion / Exercise layout, fill the
    main body placeholder (largest text area), raise its stubby height, attach
    notes, and move it to at_index."""
    slide = prs.slides.add_slide(layout_named(prs, layout_name))
    sk.set_title(slide, title)
    body = pt._body_placeholder(slide)               # largest text area (idx 1)
    assert body is not None, f"no body placeholder on layout {layout_name!r}"
    body.width = body.width    # pin inherited width: python-pptx writes cx=0 otherwise
    body.height = BODY_H
    tf = body.text_frame
    tf.word_wrap = True
    for k, b in enumerate(bullets):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.text = b
    sk.set_notes(slide, notes)
    sk.move_slide(prs, len(prs.slides._sldIdLst) - 1, at_index)
    return slide


def _set_para_text(p_el, text):
    """Replace a paragraph's text, keeping the first run's formatting."""
    runs = p_el.findall(qn("a:r"))
    if not runs:
        r = p_el.makeelement(qn("a:r"), {})
        r.append(p_el.makeelement(qn("a:t"), {}))
        p_el.append(r)
        runs = [r]
    first = runs[0]
    for extra in runs[1:]:
        p_el.remove(extra)
    for br in p_el.findall(qn("a:br")):
        p_el.remove(br)
    t = first.find(qn("a:t"))
    if t is None:
        t = first.makeelement(qn("a:t"), {})
        first.append(t)
    t.text = text


def rewrite_bullets(shape, lines):
    """Rewrite a body placeholder's bullets, reusing paragraph formatting."""
    txBody = shape.text_frame._txBody
    paras = txBody.findall(qn("a:p"))
    for i, line in enumerate(lines):
        if i < len(paras):
            _set_para_text(paras[i], line)
        else:
            newp = copy.deepcopy(paras[-1])
            txBody.append(newp)
            paras.append(newp)
            _set_para_text(newp, line)
    for p in paras[len(lines):]:
        txBody.remove(p)


def replace_bullet(shape, old_sub, new_text):
    """Rewrite one agenda bullet in place, preserving its formatting."""
    for p in shape.text_frame._txBody.findall(qn("a:p")):
        if old_sub.lower() in "".join(
                t.text or "" for t in p.iter(qn("a:t"))).lower():
            _set_para_text(p, new_text)
            return
    raise KeyError(f"bullet containing {old_sub!r} not found")


def main():
    # 0. fresh carrier -> idempotent
    shutil.copy(BAK / "1851-Ch03.pptx", OUT)
    prs = sk.open_prs(OUT)
    assert len(prs.slides) == 58, "carrier deck should have 58 slides"

    # 1. Do Now: Rules, Classic ML, or GenAI? — right after 'ML vs. Traditional Programming'
    add_engagement_slide(prs, "Do Now with Typing Hands", "Do Now: Rules, Classic ML, or GenAI?",
        bullets=[
            "For each mini-problem, write R (rules), M (classic ML), or G (GenAI):",
            "1) Flag spam emails        2) Total an invoice with tax",
            "3) Summarize a 200-line function        4) Predict which customers will churn",
            "5) Lint commit messages for style and clarity",
            "Time box: 7 minutes, then answers together",
        ],
        notes=("Answers: 1) M—classic supervised classification; an LLM works but is overkill and "
               "per-token priced. 2) R—deterministic arithmetic; never ML. 3) G—language in, "
               "language out. 4) M—tabular prediction. 5) Deliberately debatable: fixed style "
               "rules say R, nuanced clarity says G—and that debate is the point. Debrief: the "
               "cheapest reliable tool wins; picking the right wave is the engineering skill "
               "Chapter 1 set up."),
        at_index=find_slide(prs, "ML vs. Traditional Programming") + 1)

    # 2. Do Now: Draw the Split — after the train/val/test slides, before cross-validation
    add_engagement_slide(prs, "Do Now Writing Offline", "Do Now: Draw the Split",
        bullets=[
            "Offline, pen and paper: sketch the train/validation/test split for a churn dataset of 10,000 rows",
            "Write actual row counts on each box—pick your ratios and be ready to justify them",
            "One twist: the churn rate is 8%—does that change HOW you split, not just how much?",
            "Time box: 5 minutes, then compare with your neighbor",
        ],
        notes=("Debrief: any sane split earns credit—70/15/15 (7,000/1,500/1,500) or 80/10/10 are "
               "the usual answers. The twist is where the real lesson is: with only ~800 positive "
               "rows, a plain random split can leave the test set with too few churners to evaluate, "
               "so the split must be stratified to keep ~8% positives in every part. That ties "
               "directly back to the stratified-splitting slide just covered."),
        at_index=find_slide(prs, "Cross-Validation Overview"))

    # 3. Retitle Activity 3.1 -> Lab 3.1 in place (same exercise layout)
    idx = find_slide(prs, "Activity 3.1: Decision Tree Classifier")
    slide = prs.slides[idx]
    sk.set_title(slide, "Lab 3.1: Decision Trees vs. an LLM")
    for ph in slide.placeholders:                    # corner tag: 'Activity' -> 'Lab'
        if ph.placeholder_format.idx == 13 and ph.has_text_frame:
            for p in ph.text_frame._txBody.findall(qn("a:p")):
                if "activity" in "".join(
                        t.text or "" for t in p.iter(qn("a:t"))).lower():
                    _set_para_text(p, "Lab")
    body = pt._body_placeholder(slide)
    body.width = body.width    # pin inherited width before resizing
    body.height = BODY_H
    rewrite_bullets(body, [
        "Follow the detailed instructions in the Lab 3.1 notebook on your VM",
        "Train a small decision tree on the income dataset and measure its accuracy",
        "Ask an LLM the same predictions on the same rows—zero-shot, no training data",
        "Compare accuracy, interpretability, and cost—and decide which one you would ship",
        "Time: ~30 minutes",
    ])
    sk.set_notes(slide,
        "The old income-classifier activity is now the day-one ML lab with a GenAI twist: same "
        "task, two very different tools. Expect the small decision tree to win on accuracy and "
        "cost and to be fully inspectable, while the LLM is surprisingly decent zero-shot but "
        "opaque and per-token priced. Debrief on the trade-off table learners build—this is the "
        "'Where Classic ML Still Wins' slide made hands-on, and it previews the eval mindset of "
        "Chapter 6.")

    # 4. Agenda: carry the Lab 3.1 mention (replaces the old activity reference)
    replace_bullet(pt._body_placeholder(prs.slides[2]),
                   "Regression and Decision Trees",
                   "Regression and Decision Trees (Lab 3.1)")

    pt.save(prs, OUT)

    # 5. verify: reopen, titles, blanks, duplicate partnames, zip integrity
    prs2 = Presentation(OUT)
    n = len(prs2.slides)
    print(f"\n{OUT.name}: {n} slides")
    empty = []
    for i, s in enumerate(prs2.slides, 1):
        t = sk.slide_title(s)
        if not t.strip():
            empty.append(i)
        print(f"{i:>3} | {t}")
    assert not empty, f"empty titles at {empty}"
    assert n == 60, f"slide count {n} != 60"
    with zipfile.ZipFile(OUT) as z:
        dupes = [p for p, c in Counter(z.namelist()).items() if c > 1]
        assert not dupes, f"duplicate zip partnames: {dupes}"
        assert z.testzip() is None, "zip integrity check failed"
    print("\nOK: no empty titles, no duplicate partnames, zip integrity verified")


if __name__ == "__main__":
    main()
