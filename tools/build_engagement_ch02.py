#!/usr/bin/env python3
"""build_engagement_ch02.py — Day-1 engagement pass on the A.2 Ch02 deck.

Carrier = _bak1/1851-Ch02.pptx (restored fresh every run -> idempotent).
Adds two unnumbered Do Now warm-ups ('Hallucination Hunt' after the Agenda;
'The Ethics Vote' before 'Embedding Governance Into AI Workflows'), adds the
Lab 2.1 exercise slide just before the Summary, and adds a Lab 2.1 mention to
the Agenda. 30 -> 33 slides.

Run:  /Users/iahmad/Creator/Courses_and_conferences/LT/.venv-courseware/bin/python tools/build_engagement_ch02.py
"""
from __future__ import annotations
import shutil
import sys
import zipfile
from collections import Counter
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import slidekit as sk  # noqa: F401  (import activates the partname-collision fix)
import pptx_tools as pt
from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
BAK = ROOT / "decks" / "_bak1"
OUT = ROOT / "decks" / "1851-Ch02.pptx"

BODY_H = Inches(4.5)   # raise stubby layout-default body areas (house pattern)


def layout_named(prs, name):
    for lay in prs.slide_layouts:
        if (lay.name or "") == name:
            return lay
    raise KeyError(f"layout {name!r} not found")


def find_slide(prs, title_sub):
    """First slide index whose title contains title_sub (case-insensitive)."""
    sub = title_sub.lower()
    for i, s in enumerate(prs.slides):
        if sub in sk.slide_title(s).lower():
            return i
    raise KeyError(f"slide titled like {title_sub!r} not found")


def add_engagement_slide(prs, layout_name, title, bullets, notes, at_index):
    """Add a slide on a named Do Now / Discussion / Exercise layout, fill the
    main body placeholder (largest text area), raise its stubby height, attach
    notes, and move it to at_index."""
    slide = prs.slides.add_slide(layout_named(prs, layout_name))
    sk.set_title(slide, title)
    body = pt._body_placeholder(slide)               # largest text area (idx 1)
    assert body is not None, f"no body placeholder on layout {layout_name!r}"
    body.width = body.width    # pin inherited width: python-pptx writes cx=0 otherwise
    body.height = BODY_H
    tf = body.text_frame
    tf.word_wrap = True
    for k, b in enumerate(bullets):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.text = b
    sk.set_notes(slide, notes)
    sk.move_slide(prs, len(prs.slides._sldIdLst) - 1, at_index)
    return slide


def _set_para_text(p_el, text):
    """Replace a paragraph's text, keeping the first run's formatting."""
    runs = p_el.findall(qn("a:r"))
    if not runs:
        r = p_el.makeelement(qn("a:r"), {})
        r.append(p_el.makeelement(qn("a:t"), {}))
        p_el.append(r)
        runs = [r]
    first = runs[0]
    for extra in runs[1:]:
        p_el.remove(extra)
    for br in p_el.findall(qn("a:br")):
        p_el.remove(br)
    t = first.find(qn("a:t"))
    if t is None:
        t = first.makeelement(qn("a:t"), {})
        first.append(t)
    t.text = text


def fill_empty_bullet(shape, new_text):
    """Put text into the agenda's empty trailing bullet, preserving formatting."""
    for p in shape.text_frame._txBody.findall(qn("a:p")):
        if not "".join(t.text or "" for t in p.iter(qn("a:t"))).strip():
            _set_para_text(p, new_text)
            return
    raise KeyError("no empty bullet found")


def main():
    # 0. fresh carrier -> idempotent
    shutil.copy(BAK / "1851-Ch02.pptx", OUT)
    prs = sk.open_prs(OUT)
    assert len(prs.slides) == 30, "carrier deck should have 30 slides"

    # 1. Do Now: Hallucination Hunt — right after the Agenda (position 4)
    add_engagement_slide(prs, "Do Now with Typing Hands", "Do Now: Hallucination Hunt",
        bullets=[
            "Three real assistant outputs from this week’s project—ONE contains a fabricated fact or API:",
            "A) “In Python, use requests.fetch(url), then call .parse() on the response object”",
            "B) “git rebase -i HEAD~3 opens an editor to squash the last three commits”",
            "C) “HTTP 404 means the resource was not found; 410 means it is permanently gone”",
            "Which one is the hallucination? Time box: 5 minutes, then vote",
        ],
        notes=("Answer: A. The requests library has requests.get(url); there is no requests.fetch "
               "and no .parse() on the response—that is a JavaScript-flavored API the model "
               "confidently invented. B and C are both accurate. Debrief: notice how fluent and "
               "plausible the fabrication is—that is exactly what makes hallucinations dangerous. "
               "The verification habit: run it, or check the docs. This previews Lab 2.1, where "
               "they will elicit these on purpose."),
        at_index=3)

    # 2. Do Now: The Ethics Vote — just before the governance block
    add_engagement_slide(prs, "Discussion", "Do Now: The Ethics Vote",
        bullets=[
            "Scenario: your support team wants to paste raw customer tickets into a free public chatbot to draft replies faster",
            "The tickets contain customer names, order numbers, and the occasional rant about a coworker",
            "Vote: YES / NO / YES-WITH-CONDITIONS—and be ready to defend your vote",
            "In pairs: argue the strongest case AGAINST your own vote, then switch sides",
            "Time box: 8 minutes, then a room vote",
        ],
        notes=("Take the room vote first, then pairs argue against their own position—that flip is "
               "where the learning happens. Expected landing: as posed, this is a NO (or a heavily "
               "conditioned yes)—pasting personal data into a public external tool breaks data "
               "minimization and likely GDPR/PIPEDA obligations. Conditions that flip it: an "
               "approved enterprise tool with a data-processing agreement, redaction, no secrets, "
               "and human review before sending. This runs straight into the governance slides that "
               "follow."),
        at_index=find_slide(prs, "Embedding Governance Into AI Workflows"))

    # 3. Lab 2.1 — just before the Summary slide
    add_engagement_slide(prs, "Exercise Reference Slide with Typing Hands",
        "Lab 2.1: Hallucination and Bias Audit",
        bullets=[
            "Follow the detailed instructions in the Lab 2.1 notebook on your VM",
            "Elicit fabricated APIs and facts from an assistant—make it confidently wrong on purpose",
            "Run bias probes: same prompt, different names and demographics; compare the outputs",
            "Rate each finding’s severity and propose a mitigation your team could actually enforce",
            "Time: ~30 minutes",
        ],
        notes=("Learners deliberately induce the failure modes this chapter described: hallucinated "
               "references and demographically skewed outputs. Encourage them to keep the prompts "
               "realistic—the point is that these failures show up in ordinary use, not just in "
               "adversarial stunts. Debrief on the severity ratings: a wrong date-parsing API is "
               "annoying; a biased screening suggestion is a governance problem. That contrast "
               "motivates the mitigation column."),
        at_index=find_slide(prs, "Summary"))

    # 4. Agenda: fill the empty trailing bullet with the Lab 2.1 mention
    fill_empty_bullet(pt._body_placeholder(prs.slides[2]),
                      "Lab 2.1: Hallucination and Bias Audit (hands-on)")

    pt.save(prs, OUT)

    # 5. verify: reopen, titles, blanks, duplicate partnames, zip integrity
    prs2 = Presentation(OUT)
    n = len(prs2.slides)
    print(f"\n{OUT.name}: {n} slides")
    empty = []
    for i, s in enumerate(prs2.slides, 1):
        t = sk.slide_title(s)
        if not t.strip():
            empty.append(i)
        print(f"{i:>3} | {t}")
    assert not empty, f"empty titles at {empty}"
    assert n == 33, f"slide count {n} != 33"
    with zipfile.ZipFile(OUT) as z:
        dupes = [p for p, c in Counter(z.namelist()).items() if c > 1]
        assert not dupes, f"duplicate zip partnames: {dupes}"
        assert z.testzip() is None, "zip integrity check failed"
    print("\nOK: no empty titles, no duplicate partnames, zip integrity verified")


if __name__ == "__main__":
    main()
