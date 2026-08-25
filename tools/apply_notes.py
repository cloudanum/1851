#!/usr/bin/env python3
"""Add speaker notes to all slides missing them + fix Ch07 slide 26 feedback edge."""
from pptx import Presentation
from pptx.util import Inches

BASE = "/Users/iahmad/Creator/Courses_and_conferences/LT/courses/1851/1851-FTP-upload"

NOTES = {
("1851-Ch00.pptx", 1): "Welcome the group and set expectations: this course treats AI as an engineering tool across the whole software lifecycle, not a magic box. Do a quick round of introductions — name, role, and one place AI already touches their work. Cover logistics: timing, breaks, and the VM every participant will use for labs.",
("1851-Ch00.pptx", 2): "Read these objectives as a contract for the course — each one maps to a chapter and at least one hands-on lab on the VM. Point out the arc: foundations first (Chapters 1–3), then prompting and the API (4), the SDLC deep-dive (5), evals (6), agents (7), security (8), and the capstone (9).",
("1851-Ch00.pptx", 4): "Walk the table of contents briefly — do not read every line. Emphasize that each chapter ends with a lab or activity run on the course VM, and that Chapter 9 is a capstone tying the whole lifecycle together. Confirm everyone can log into their VM before moving on.",

("1851-Ch01.pptx", 1): "Chapter 1 sets the vocabulary for the whole course. Goal for the session: everyone leaves able to say precisely what kind of AI they are looking at — predictive, generative, or agentic — and what that means for their SDLC.",
("1851-Ch01.pptx", 2): "Frame these objectives as the day's exit criteria. The last one — trust-but-verify — is the thread running through the entire course; flag it now and return to it in every later chapter.",
("1851-Ch01.pptx", 21): "Debrief the role-play here. The insight to land: fact questions are weak probes because machines ace facts — the differences show up in emotion, creativity, and inconsistency. Ask each pair which single response surprised them most.",
("1851-Ch01.pptx", 22): "Open-floor discussion. Push past 'did it fool you?' to the design question: if fluency is not intelligence, what would you measure instead? This seeds the evaluation mindset that Chapter 6 formalizes.",
("1851-Ch01.pptx", 23): "Run this exactly as written in the VM activity. Keep the pace brisk: the point is the felt experience of the imitation game, not winning. Timebox each round and rotate roles so everyone interrogates at least once. Debrief on the next two slides.",
("1851-Ch01.pptx", 24): "Six phases — and notice the first three are all understanding and data; the modeling phase is the small middle, not the main event. Students who have shipped ML will recognize this; for the rest, contrast it with a classic software project plan. This lifecycle returns in Chapter 6 when we build evals.",
("1851-Ch01.pptx", 40): "Hit the through-lines: AI is commodity tooling now, accountability did not move, and tokens, context, and cost are the physics of every GenAI design. One-minute whip-around: each person names one SDLC phase where they will try AI assistance this month.",
("1851-Ch01.pptx", 41): "Use this as a self-check, not a formality. Read each objective and ask for a show of hands — confident, shaky, or lost. Anything shaky gets revisited before the lab or in the next morning's recap.",

("1851-Ch02.pptx", 1): "Chapter 2 is the trust chapter: what makes AI output good, and what makes it responsible. Keep it concrete — every concept gets a software-engineering example, and the lab makes hallucination and bias tangible on the VM.",
("1851-Ch02.pptx", 33): "Close the loop on the opening objectives. For each one, ask a volunteer for a one-sentence answer or example. Bias detection and documentation are the two that matter most for the lab — make sure both landed.",

("1851-Ch03.pptx", 1): "Chapter 3 is the only pure-ML chapter, and it is deliberately hands-on. The point is not to turn developers into data scientists — it is to give them the mechanics (data, splits, metrics, overfitting) that LLM evals and agents build on later.",
("1851-Ch03.pptx", 13): "Build one tiny tree live on the whiteboard from four or five rows before showing any library code. The depth trade-off is the concept to nail: depth is where interpretability and overfitting fight each other.",
("1851-Ch03.pptx", 14): "Two slides on trees, then move on — trees are the teaching vehicle, not the destination. The limitation list (overfitting, instability, dominant features) is exactly why evaluation and data hygiene matter; that is the bridge to the metrics section.",
("1851-Ch03.pptx", 15): "Do the arithmetic by hand once with a 2x2 on the board — twenty rows, count TP/FP/FN/TN together. Precision versus recall only sticks when tied to a cost: ask which error is expensive in spam filtering, in cancer screening, in code review.",
("1851-Ch03.pptx", 20): "Unsupervised learning in one slide: no labels, just geometry. Stress the two practical traps — normalize your features first, and never trust k without plotting. The elbow-method demo is in the VM notebook if the group wants it.",
("1851-Ch03.pptx", 34): "Short but critical slide. The one-sentence version: on imbalanced data, a random split can put zero rare-class examples in the test set and your metrics lie. stratify=y is the fix — show it in the notebook.",
("1851-Ch03.pptx", 37): "Leakage is the most common way smart people ship broken models. Give the canonical example: scaling or imputing before the split lets test statistics bleed into training. Rule of thumb — fit preprocessors on the training fold only, always.",
("1851-Ch03.pptx", 43): "Connect back to the confusion matrix: every threshold gives a different (TPR, FPR) pair, and the ROC curve is all of them at once; AUC is the threshold-free summary. Warn against quoting AUC on heavily imbalanced data — precision-recall curves tell the truth there.",
("1851-Ch03.pptx", 60): "Run the self-check against the opening objectives. The last one matters most for the rest of the course: students should be able to say in one sentence how a decision tree's training loop relates to how an LLM is trained — and how both get evaluated.",

("1851-Ch04.pptx", 1): "This is the pivotal chapter — everything before was concepts; everything from here is engineering. By the end of the day students will make real API calls from the VM and engineer prompts with the same discipline as code.",
("1851-Ch04.pptx", 3): "Preview the shape of the day: landscape and API mechanics first, then techniques, then agent-grade prompting with the PTCF blueprint. Both the lab and the clinic activity are hands-on in the VM — flag that API keys are preconfigured there.",
("1851-Ch04.pptx", 55): "Seven takeaways; the ones to say out loud are the first and the fifth: prompting is engineering, and prompts get versioned and tested like code. That framing separates this course from a prompt-recipe list, and Chapter 6 turns it into evals.",
("1851-Ch04.pptx", 56): "Self-check before the lab. Everyone should be able to name the four PTCF elements without looking, and explain temperature in one sentence. If either is shaky, revisit the relevant section now.",

("1851-Ch06.pptx", 1): "Chapter 6 is where 'trust but verify' becomes a discipline. The message for the day: if you cannot measure an LLM feature, you do not ship it. Everything here is practiced in Lab 6.1.",
("1851-Ch06.pptx", 3): "The agenda runs narrow to wide: datasets first, then scoring methods, then levels of evaluation up to full agent trajectories, then operations. Eval-Driven Development is the cultural ask — evals first, just like tests.",
("1851-Ch06.pptx", 46): "Closing self-check. The two objectives that carry forward are golden datasets and CI gating — students will reuse both in the agent chapters. Ask each person to name one behavior in their own product that needs an eval.",

("1851-Ch07.pptx", 1): "Agents are the payoff chapter — Chapter 4's API skills and Chapter 6's evals combine into systems that act. Set the safety frame early: every capability we add gets a matching guardrail, and Lab 7.1 is guarded by construction.",
("1851-Ch07.pptx", 2): "Read the objectives as two columns: capability (loops, tools, orchestration) and control (allow-lists, caps, human approval, evaluation). The course thesis in miniature: autonomy scales only as fast as verification.",
("1851-Ch07.pptx", 3): "Walk the arc: what an agent is, how it acts through tools, how patterns scale to multi-agent, then failure modes and evaluation. The vendor-landscape section is a map, not a catalog — the goal is orientation, not memorization.",
("1851-Ch07.pptx", 53): "Self-check against the opening objectives. Before Lab 7.1, confirm everyone can state the anatomy (LLM + tools + loop) and name two guardrails. Those two answers are the lab in miniature.",

("1851-Ch08.pptx", 1): "Chapter 8 flips perspective: we have spent seven chapters building with AI; now we attack what we built. The red-team lab at the end reuses the Chapter 7 agent, so keep that architecture fresh in mind all day.",
("1851-Ch08.pptx", 2): "Frame the objectives around one question: how does an AI system fail — accidentally, adversarially, or politically? The four-dimension model on the agenda is the organizing spine for everything that follows.",
("1851-Ch08.pptx", 3): "Point out the structure: vulnerability dimensions first, then attacks, then the ATLAS taxonomy as shared vocabulary, then defenses and controls. Activities 8.1 and 8.2 plus the red-team lab are all on the VM.",
("1851-Ch08.pptx", 8): "This is the 'so what' slide for the four dimensions — each row is an action, not a concept. Ask the group which of the four dimensions their own organization under-invests in; that answer predicts their real exposure.",
("1851-Ch08.pptx", 11): "The NIST taxonomy is the reference — students do not need to memorize it, they need to be able to look things up in it. Walk one example per category: availability, integrity, privacy, abuse. The citation on the slide is the source to bookmark.",
("1851-Ch08.pptx", 16): "ATLAS is the adversary's playbook, structured like ATT&CK so security teams can reuse what they know. Key message: these are observed, real-world techniques, not theory. The Do Now on the next slide gets hands on the actual matrix.",
("1851-Ch08.pptx", 17): "Give this real time — ten minutes minimum. Have everyone pick one technique, read its case study, and be ready to explain it in two sentences. Debrief question: which tactic would hit your current project first?",
("1851-Ch08.pptx", 19): "Integrity is the dimension AI stresses hardest — a poisoned model is an integrity failure, not just bad output. Hashing, signatures, and version control are familiar; the new part is applying them to datasets and model weights.",
("1851-Ch08.pptx", 20): "Availability for AI systems has a twist: the dependency chain includes GPUs, model endpoints, and rate limits, not just servers. Keep it brief — the mechanisms (redundancy, backups, DoS protection) are classic; the attack surface is new.",
("1851-Ch08.pptx", 22): "A checklist slide — do not read it line by line. The pattern to extract: secure the data, the model, and the pipeline as three separate assets, each with its own monitoring and feedback loop.",
("1851-Ch08.pptx", 27): "Run from the VM instructions. The mindset shift to coach: minimization is not deletion for its own sake — it is keeping exactly what the task needs and being able to justify each field. The peer-comparison step is where the learning happens.",
("1851-Ch08.pptx", 29): "Three layers: data, model, infrastructure. Emphasize the model row — versioned models with rollback turn a bad deployment from an incident into a non-event. Ask who in the room can roll back their model today.",
("1851-Ch08.pptx", 30): "The classic control trio applied to AI assets: authentication, authorization, encryption. The AI-specific nuance is scope — these controls must cover prompts, datasets, and models, not just servers and users.",
("1851-Ch08.pptx", 31): "Operations is where controls live or die. Logging and threat detection close the loop with the Chapter 6 observability material — the same traces that evaluate quality also reveal attacks.",
("1851-Ch08.pptx", 49): "Straight from the VM instructions. Two coaching points: pattern matching catches the obvious PII, and the interesting failures are contextual — a name in free text. That gap is exactly why semantic detection matters.",
("1851-Ch08.pptx", 51): "Insider threats feel awkward to discuss — frame them as blast-radius control, not distrust. Each mitigation pairs a detection control with a friction control; the two-person rule for production changes is the one to champion back at work.",
("1851-Ch08.pptx", 54): "Closing self-check before the red-team lab. Everyone should be able to name the four vulnerability dimensions and one ATLAS tactic. The lab is the objective test: if they can break the agent, they understood the chapter.",

("1851-Ch09.pptx", 1): "The wrap-up chapter is deliberately short: a capstone lab that runs one ticket through the entire AI-augmented lifecycle, then next steps. Keep energy high — this is where the course becomes their Monday morning.",
("1851-Ch09.pptx", 4): "Run this as the course-long retrospective: read each objective and take a quick confidence poll. Any objective the group rates shaky gets a five-minute review now — the material is all in the chapters and VM notebooks.",
}

# --- fix Ch07 slide 26 feedback edge (shape id 231: reposition left, drop flipH)
p7 = Presentation(f"{BASE}/1851-Ch07.pptx")
s26 = p7.slides[25]
for sh in s26.shapes:
    if sh.shape_id == 231:
        sh.left = Inches(1.613079615048119)
        xfrm = sh._element.spPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm")
        if xfrm is not None and xfrm.get("flipH"):
            del xfrm.attrib["flipH"]
        print("Ch07 s26 feedback edge repositioned:", sh.left, sh.width)
p7.save(f"{BASE}/1851-Ch07.pptx")

# --- apply notes
from collections import defaultdict
by_file = defaultdict(dict)
for (f, i), note in NOTES.items():
    by_file[f][i] = note

for f, slides in by_file.items():
    p = Presentation(f"{BASE}/{f}")
    applied = 0
    for i, note in slides.items():
        s = p.slides[i - 1]
        s.notes_slide.notes_text_frame.text = note
        applied += 1
    p.save(f"{BASE}/{f}")
    print(f"{f}: {applied} notes written")
