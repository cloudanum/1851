#!/usr/bin/env python3
"""build_ch04.py — assemble the A.2 'Prompt Engineering with OpenAI' chapter.

Carrier = _bak0/1851-Ch06.pptx (A.1 'Prompt Engineering for Generative AI',
restored fresh every run -> idempotent). Keeps the carrier's prompting-core
slides, drops agentic (now Ch07), data-science (now Ch03), dated IDE/LangChain
tooling, and the superseded lifecycle/API slides. Clones the 2026 model
landscape + technique slides from 1258-Ch04 and the OpenAI API anatomy block
from 1258-Ch07 (cross-deck copy re-embedding images via pptx_tools.clone_slide).
Authors 4 new slides (Why OpenAI First?, Prompting as an Engineering
Discipline, Agents: Preview, Lab 4.1), then arranges everything to plan.

Run:  /Users/iahmad/Creator/Courses_and_conferences/LT/.venv-courseware/bin/python tools/build_ch04.py
"""
from __future__ import annotations
import copy
import shutil
import sys
import zipfile
from collections import Counter
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import slidekit as sk  # noqa: F401  (import activates the partname-collision fix)
import pptx_tools as pt
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent
BAK = ROOT / "decks" / "_bak0"
OUT = ROOT / "decks" / "1851-Ch04.pptx"
CARRIER = BAK / "1851-Ch06.pptx"
SRC1258 = Path("/Users/iahmad/Creator/Courses_and_conferences/LT/courses/"
               "1258-applied-ai-for-government-it-professionals/"
               "1258a4-author-input/decks")

# --- clone lists (1-based, against the 1258 decks) ---------------------------
CLONE4 = [7, 8, 12,                       # 2026 model landscape + context windows
          22, 23, 24,                     # zero/one/few-shot (with figures)
          27,                             # chain-of-thought (with figure)
          30, 34,                         # reasoning models + reasoning.effort
          41,                             # role & system prompts
          43, 44,                         # JSON / schema-valid structured output
          45,                             # grounding
          47,                             # reusable prompt templates
          55]                             # fact-checking / verification habits
CLONE7 = [8, 9, 10, 11, 12, 13, 14,       # OpenAI API anatomy (s15 FedRAMP dropped: gov-only)
          17]                             # LLM application layer-stack figure

# --- rewritten text for kept carrier slides ----------------------------------
OBJECTIVES = [
    "Navigate the 2026 model landscape—families, tiers, and context windows",
    "Make OpenAI API calls with the Responses API: messages, roles, tokens, and temperature",
    "Apply core techniques: zero/one/few-shot prompting, chain-of-thought, and reasoning effort",
    "Engineer prompts with roles, structure, constraints, grounding, and schema-valid output",
    "Version, test, and evaluate prompts as engineering artifacts",
    "Apply ethics, privacy, and security practices—and preview RAG and agents",
]
AGENDA = [
    "The 2026 model landscape and why we start with OpenAI",
    "The OpenAI API: messages, tokens, temperature, structured outputs",
    "Prompting techniques from zero-shot to reasoning models",
    "Grounding, templates, evaluation, and prompts in production",
    "RAG, ethics and privacy, Lab 4.1, and what's next",
]
SUMMARY = [
    "Prompt engineering turns intent into reliable model behavior—it is an engineering discipline",
    "2026 frontier: tiered model families, ~1M-token context windows, reasoning-effort controls",
    "The Responses API: role-tagged messages, per-token pricing, temperature, strict JSON Schema",
    "Core techniques: zero/one/few-shot, chain-of-thought, system prompts, grounding, templates",
    "Version, test, and review prompts like code—Chapter 6 adds rigorous evaluation",
    "RAG grounds prompts in your documents; ethics and privacy checks apply to every prompt",
    "Next: Chapter 5 applies GenAI across the SDLC; Chapter 7 builds agents on these APIs",
]
OBJECTIVES_NOTES = (
    "Welcome to Chapter 4. Earlier chapters established what modern AI is and how classic ML "
    "underpins it; this chapter is where engineers get hands-on with the models themselves—first "
    "through the OpenAI API, then through disciplined prompting. Walk the objectives: they trace "
    "the chapter's arc from landscape, to API, to techniques, to engineering practice.")

# --- 4 new slides ------------------------------------------------------------
NEW_WHY = ("Why OpenAI First?",
    ["The most mature hosted API: a decade of production use and stable SDKs (openai>=1.x)",
     "The Responses API is now the primary interface—one call shape for text, reasoning, and tools",
     "Structured Outputs and function calling are platform features with schema guarantees",
     "The ecosystem converged on it: Azure OpenAI and many open-weight servers speak an OpenAI-compatible API",
     "Not the only option—Claude, Gemini, and open-weight models; every concept here transfers"],
    "Why center a hands-on chapter on one vendor? Practicality, not endorsement: the OpenAI API is "
    "the most mature and the most imitated interface in the industry. Make the neutrality explicit—"
    "roles, tokens, temperature, and structured output map onto Claude, Gemini, and open-weight "
    "stacks; only the parameter names differ.")
NEW_ENG = ("Prompting as an Engineering Discipline",
    ["Prompts are behavior-changing configuration: store them in Git with the code they serve",
     "Review prompt changes in pull requests—a one-word change can change every output",
     "Keep a fixed test set per prompt; re-run it on every edit and every model upgrade",
     "Pin model and parameters (model name, temperature, reasoning effort) for reproducibility",
     "Chapter 6 builds the full evaluation practice: rubrics, LLM-as-judge, and regression gates"],
    "Engineers already know this discipline—they just have not always applied it to prompts. Version "
    "prompts, review prompt diffs, and test prompts against fixed cases, exactly like code. This "
    "slide is the on-ramp to Chapter 6, where ad-hoc testing becomes a real evaluation practice.")
NEW_AGENTS = ("Agents: Preview",
    ["So far: one prompt, one response—agents put the model in a loop with tools",
     "Function calling (this chapter's API section) is the mechanism agents are built on",
     "Grounding and structured output become agent guardrails: what it may read, do, and return",
     "Chapter 7 builds agents hands-on: tool use, the ReAct loop, guardrails, and failure modes"],
    "A forward pointer only—do not teach agents here. The message: everything in this chapter "
    "(messages, function calling, grounding, structured output) is exactly what Chapter 7 assembles "
    "into an agent loop. One slide, then move on.")
NEW_LAB = ("Lab 4.1: First OpenAI API Calls",
    ["Follow the detailed instructions in the Lab 4.1 notebook on your VM",
     "Make your first Responses API call—send messages with roles and read the reply",
     "Sweep temperature and reasoning effort to see the same prompt's output change",
     "Get schema-valid JSON back with Structured Outputs, and print token usage and cost"],
    "The classroom VM has OPENAI_API_KEY pre-injected, so students call live current models from "
    "the notebook without handling secrets. The lab rehearses the whole API section: roles and "
    "messages, temperature and reasoning effort, strict JSON, and token cost. Circulate and make "
    "sure everyone has a successful first call before the later steps.")

# --- text fixes on cloned slides (audience + cross-reference repairs) ---------
# (clone key, old substring, new substring); old must appear verbatim in one run
# or one paragraph of the cloned slide.
CLONE_FIXES = [
    ((4, 12), "you retrieve the relevant parts (RAG, Chapter 7)",
              "you retrieve the relevant parts (RAG, later in this chapter)"),
    ((4, 41), "'You are a FOIA officer. Cite the exemption for every redaction.'",
              "'You are a senior code reviewer. Cite the rule behind every requested change.'"),
    ((4, 41), "In code (Chapter 7) the system prompt",
              "In code (the API section ahead) the system prompt"),
    ((4, 43), "(Lab 7.1 does this against a memo)", "(Lab 4.1 does this)"),
    ((4, 44), "(Lab 7.1)", "(Lab 4.1)"),
    ((4, 45), "Chapter 7 automates it over many documents",
              "the RAG section ahead automates it over many documents"),
    ((4, 45), "For government work, grounded + cited",
              "For production systems, grounded + cited"),
    ((4, 34), "eligibility analysis = high; form letter = low or none",
              "code review on a payment path = high; commit-message draft = low or none"),
    ((4, 55), "may be federal records", "may be auditable deliverables"),
    ((7, 8),  "This is how agencies build AI into case systems, portals, and pipelines",
              "This is how teams build AI into products, services, and pipelines"),
    ((7, 8),  "Lab 7.1 makes your first working calls",
              "Lab 4.1 makes your first working calls"),
    ((7, 10), "Lab 7.1 prints token usage and cost",
              "Lab 4.1 prints token usage and cost"),
    ((7, 11), "Lab 7.1 sweeps temperature", "Lab 4.1 sweeps temperature"),
    ((7, 12), "The mechanism the agent lab (7.3) is built on",
              "The mechanism agents (Chapter 7) are built on"),
]


# --- formatting-preserving text rewrite helpers ------------------------------
def _set_para_text(p_el, text):
    """Replace a paragraph's text, keeping the first run's formatting."""
    runs = p_el.findall(qn("a:r"))
    if not runs:
        r = p_el.makeelement(qn("a:r"), {})
        r.append(p_el.makeelement(qn("a:t"), {}))
        p_el.append(r)
        runs = [r]
    first = runs[0]
    for extra in runs[1:]:
        p_el.remove(extra)
    for br in p_el.findall(qn("a:br")):
        p_el.remove(br)
    t = first.find(qn("a:t"))
    if t is None:
        t = first.makeelement(qn("a:t"), {})
        first.append(t)
    t.text = text


def rewrite_bullets(shape, lines):
    """Rewrite a body placeholder's bullets, reusing paragraph formatting."""
    txBody = shape.text_frame._txBody
    paras = txBody.findall(qn("a:p"))
    for i, line in enumerate(lines):
        if i < len(paras):
            _set_para_text(paras[i], line)
        else:
            newp = copy.deepcopy(paras[-1])
            txBody.append(newp)
            paras.append(newp)
            _set_para_text(newp, line)
    for p in paras[len(lines):]:
        txBody.remove(p)


def replace_in_slide(slide, old, new):
    """Replace `old` with `new` anywhere in the slide's text. Run-level first
    (preserves inline formatting); falls back to a paragraph-level rewrite.
    Returns the number of replacements (asserted > 0 by the caller)."""
    n = 0
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if old in r.text:
                    r.text = r.text.replace(old, new)
                    n += 1
    if n:
        return n
    # fallback: the string spans runs — rewrite the whole paragraph's text
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p_el in sh.text_frame._txBody.findall(qn("a:p")):
            full = "".join(t.text or "" for t in p_el.iter(qn("a:t")))
            if old in full:
                _set_para_text(p_el, full.replace(old, new))
                n += 1
    return n


def add_exercise_slide(prs, title, bullets, notes):
    """Add a slide on the template's exercise layout (typing-hands graphic on
    the left, text column on the right). Returns the new slide."""
    lay = sk.pick_layout(prs, "Exercise Reference Slide with Typing Hands")
    assert lay is not None, "exercise layout not found in carrier template"
    slide = prs.slides.add_slide(lay)
    sk.set_title(slide, title)
    body = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1 and ph.has_text_frame:
            body = ph
            break
    if body is None:
        body = pt._body_placeholder(slide)
    # the layout's content placeholder is stubby; give it room for 4+ lines
    body.height = Inches(4.5)
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
    sk.set_notes(slide, notes)
    return slide


def main():
    # 0. fresh carrier -> idempotent
    shutil.copy(CARRIER, OUT)
    prs = sk.open_prs(OUT)

    # 1. clone keeps from 1258 Ch04/Ch07 (appended at end), record new indices
    idx4, idx7 = {}, {}
    src4 = pt.open_deck(SRC1258 / "1258-Ch04-ModernGenAI-PromptEng.pptx")
    for k in CLONE4:
        pt.clone_slide(prs, src4.slides[k - 1])
        idx4[k] = len(prs.slides) - 1
    src7 = pt.open_deck(SRC1258 / "1258-Ch07-Building-with-LLMs.pptx")
    for k in CLONE7:
        pt.clone_slide(prs, src7.slides[k - 1])
        idx7[k] = len(prs.slides) - 1

    # 2. author the 4 new slides (appended at end)
    new_why = pt.append_content(prs, NEW_WHY[0], NEW_WHY[1], NEW_WHY[2])
    i_why = len(prs.slides) - 1
    pt.append_content(prs, NEW_ENG[0], NEW_ENG[1], NEW_ENG[2])
    i_eng = len(prs.slides) - 1
    pt.append_content(prs, NEW_AGENTS[0], NEW_AGENTS[1], NEW_AGENTS[2])
    i_agents = len(prs.slides) - 1
    add_exercise_slide(prs, NEW_LAB[0], NEW_LAB[1], NEW_LAB[2])
    i_lab = len(prs.slides) - 1

    # 3. edit kept carrier slides in place (0-based indices still original)
    s1 = prs.slides[0]
    _set_para_text(s1.shapes.title.text_frame._txBody.findall(qn("a:p"))[0],
                   "Prompt Engineering with OpenAI")
    for sh in s1.shapes:                       # subtitle: Chapter 6 -> Chapter 4
        if sh.has_text_frame and "Chapter 6" in sh.text_frame.text:
            _set_para_text(sh.text_frame._txBody.findall(qn("a:p"))[0], "Chapter 4")
    rewrite_bullets(pt._body_placeholder(prs.slides[1]), OBJECTIVES)   # Objectives
    sk.set_notes(prs.slides[1], OBJECTIVES_NOTES)                      # (old notes were stale)
    rewrite_bullets(pt._body_placeholder(prs.slides[2]), AGENDA)       # Agenda
    rewrite_bullets(pt._body_placeholder(prs.slides[35]), SUMMARY)     # Summary
    rewrite_bullets(pt._body_placeholder(prs.slides[36]), OBJECTIVES)  # objectives recap

    # 4. audience / cross-reference fixes on cloned slides
    slides = list(prs.slides)
    pos = {**{(4, k): v for k, v in idx4.items()}, **{(7, k): v for k, v in idx7.items()}}
    for key, old, new in CLONE_FIXES:
        n = replace_in_slide(slides[pos[key]], old, new)
        assert n > 0, f"fix not applied on clone {key}: {old!r}"
    sk.set_title(slides[idx4[55]], "Verification and Fact-Checking Habits")

    # 5. arrange: openers -> landscape/why -> API -> techniques -> structured
    #    output/grounding -> engineering discipline -> RAG -> agents preview ->
    #    ethics/privacy -> lab -> summary -> recap
    order = (
        [0, 1, 2]                                                   # openers
        + [idx4[7], idx4[8], idx4[12], i_why]                       # landscape + why OpenAI
        + [idx7[k] for k in [8, 9, 10, 11, 12, 13, 14]]             # API anatomy
        + [idx7[17]]                                                # LLM app stack figure
        + [10, 11, 12, 13, idx4[41]]                                # what/why/best/role/system
        + [idx4[22], idx4[23], idx4[24]]                            # zero/one/few-shot
        + [idx4[27], idx4[30], idx4[34]]                            # CoT -> reasoning models/effort
        + [17, 18, 19]                                              # structure/constraints/context
        + [20, 21, idx4[47], 22]                                    # pitfalls/iteration/templates/tuning
        + [idx4[43], idx4[44], idx4[45]]                            # structured output + grounding
        + [i_eng, 25, 26, 33, idx4[55]]                             # engineering discipline
        + [34, i_agents]                                            # RAG -> agents preview
        + [29, 30]                                                  # ethics / privacy
        + [i_lab]                                                   # Lab 4.1
        + [35, 36]                                                  # summary + recap
    )
    pt.arrange(prs, order)
    pt.save(prs, OUT)

    # 6. verify: reopen, titles, blanks, duplicate partnames, zip integrity
    prs2 = Presentation(OUT)
    n = len(prs2.slides)
    print(f"\n{OUT.name}: {n} slides")
    print(f"kept carrier: 21 | cloned: Ch04={len(CLONE4)} Ch07={len(CLONE7)} | authored: 4")
    empty = []
    for i, s in enumerate(prs2.slides, 1):
        t = sk.slide_title(s)
        if not t.strip():
            empty.append(i)
        print(f"{i:>3} | {t}")
    assert not empty, f"empty titles at {empty}"
    assert 46 <= n <= 54, f"slide count {n} outside 46-54"
    with zipfile.ZipFile(OUT) as z:
        dupes = [p for p, c in Counter(z.namelist()).items() if c > 1]
        assert not dupes, f"duplicate zip partnames: {dupes}"
        assert z.testzip() is None, "zip integrity check failed"
    print("\nOK: no empty titles, no duplicate partnames, zip integrity verified")


if __name__ == "__main__":
    main()
